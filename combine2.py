#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Rebuild a PPTX from the current layer JSON protocol.

- Text layers are rebuilt from `paragraphs[].runs`.
- Image layers are placed from final RGBA/image assets plus final ratio boxes.
- Shape layers are rebuilt from exported XML snippets.

Layers are inserted bottom -> top because the JSON order is top -> bottom.
"""

import argparse
import json
import os
import math
import shutil
import tempfile
import zipfile
from collections import defaultdict
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Pt

EMU_PER_PT = 12700
XML_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"

# Global knob for rel-only exports. Default 1.0 keeps current behavior.
_REL_SCALE = 1.0
_REL_BASE = "height"  # height|min|geom


def _rel_base_pt(slide_size_pt: Tuple[float, float]) -> float:
    """Return the baseline (in pt) for rel-only metrics."""
    slide_w_pt, slide_h_pt = slide_size_pt
    token = str(_REL_BASE).lower()
    if token == "min":
        return min(float(slide_w_pt), float(slide_h_pt))
    if token == "geom":
        # Geometric mean baseline: sqrt(w * h)
        try:
            return math.sqrt(float(slide_w_pt) * float(slide_h_pt))
        except Exception:
            return float(slide_h_pt)
    return float(slide_h_pt)


def _text_rel_base_pt(slide_size_pt: Tuple[float, float]) -> float:
    """Font/text metrics use geometric mean regardless of global geometry base."""
    try:
        slide_w_pt, slide_h_pt = slide_size_pt
        return math.sqrt(float(slide_w_pt) * float(slide_h_pt))
    except Exception:
        return _rel_base_pt(slide_size_pt)


# #region debug-point A:report
def _dbg_report_method_ppt_style(hypothesis_id: str, location: str, msg: str, data: Optional[Dict[str, Any]] = None) -> None:
    try:
        import json as _json
        import urllib.request as _urlreq
        import time as _time

        _u, _s, _r = "http://127.0.0.1:7777/event", "method-ppt-style", os.environ.get("DEBUG_RUN_ID", "pre-fix")
        try:
            with open(".dbg/method-ppt-style.env", "r", encoding="utf-8") as _f:
                for _line in _f.read().splitlines():
                    if _line.startswith("DEBUG_SERVER_URL="):
                        _u = _line.split("=", 1)[1] or _u
                    elif _line.startswith("DEBUG_SESSION_ID="):
                        _s = _line.split("=", 1)[1] or _s
                    elif _line.startswith("DEBUG_RUN_ID="):
                        _r = _line.split("=", 1)[1] or _r
        except Exception:
            pass

        _payload = {
            "sessionId": _s,
            "runId": _r,
            "hypothesisId": hypothesis_id,
            "location": location,
            "msg": msg,
            "data": data or {},
            "ts": int(_time.time() * 1000),
        }
        _urlreq.urlopen(
            _urlreq.Request(
                _u,
                data=_json.dumps(_payload).encode("utf-8"),
                headers={"Content-Type": "application/json"},
            ),
            timeout=1.0,
        ).read()
    except Exception:
        pass
# #endregion


def _ratio_box_to_pt(
    box: Optional[Dict[str, Any]], slide_size_pt: Optional[Tuple[float, float]] = None
) -> Optional[Tuple[Pt, Pt, Pt, Pt]]:
    if not box or slide_size_pt is None:
        return None
    slide_w_pt, slide_h_pt = slide_size_pt
    try:
        left_ratio = box.get("left")
        top_ratio = box.get("top")
        width_ratio = box.get("width")
        height_ratio = box.get("height")
        if None in (left_ratio, top_ratio, width_ratio, height_ratio):
            return None
        return (
            Pt(float(left_ratio) * slide_w_pt),
            Pt(float(top_ratio) * slide_h_pt),
            Pt(float(width_ratio) * slide_w_pt),
            Pt(float(height_ratio) * slide_h_pt),
        )
    except Exception:
        return None


def _slide_size_pt_from_slide(slide) -> Tuple[float, float]:
    try:
        pres = slide.part.package.presentation_part.presentation
        return float(pres.slide_width) / EMU_PER_PT, float(pres.slide_height) / EMU_PER_PT
    except Exception:
        return 720.0, 540.0


def _decode_scaled_protocol_value(value: Any, key: Optional[str] = None, parent_key: Optional[str] = None) -> Any:
    scale = 1000.0
    if isinstance(value, dict):
        current_parent = key if key is not None else parent_key
        decoded = {k: _decode_scaled_protocol_value(v, k, current_parent) for k, v in value.items()}
        keys = set(decoded.keys())
        if keys and keys.issubset({"left", "top", "width", "height"}):
            for k in list(decoded.keys()):
                if isinstance(decoded[k], int):
                    decoded[k] = float(decoded[k]) / scale
        if keys and keys.issubset({"x", "y"}):
            for k in list(decoded.keys()):
                if isinstance(decoded[k], int):
                    decoded[k] = float(decoded[k]) / scale
        return decoded
    if isinstance(value, list):
        decoded_list = [_decode_scaled_protocol_value(v, None, key) for v in value]
        if key and any(str(key).endswith(suffix) for suffix in ("_pt", "_px", "_rel", "_deg")):
            return [float(v) / scale if isinstance(v, int) else v for v in decoded_list]
        return decoded_list
    if isinstance(value, int):
        token = str(key or "")
        parent = str(parent_key or "")
        if token.endswith(("_pt", "_px", "_rel", "_deg")):
            return float(value) / scale
        if token == "rel" and parent in {"line_spacing", "space_before", "space_after"}:
            return float(value) / scale
        if token in {"w", "h"} and parent == "pathLst":
            return float(value) / scale
        if token in {"x", "y"} and parent in {"pt", "pts"}:
            return float(value) / scale
        if token in {"left", "top", "width", "height"} and parent == "box":
            return float(value) / scale
        if token in {"x", "y"} and parent == "pts":
            return float(value) / scale
    return value


def _decode_scaled_protocol_payload(payload: Any) -> Any:
    return _decode_scaled_protocol_value(payload)


def _normalize_layers_payload(payload: Any) -> List[Dict[str, Any]]:
    """
    Support two historical `layers.json` formats:
    1) Flat list of layer dicts.
    2) List-of-slides, where each slide is a list of layer dicts (no `slide` field).
    """
    if not isinstance(payload, list):
        raise SystemExit("layers.json must contain a JSON array.")

    # Format (2): [ [layer, ...], [layer, ...], ... ]
    if payload and all(isinstance(item, list) for item in payload):
        flat: List[Dict[str, Any]] = []
        global_index = 0
        for slide_idx, slide_layers in enumerate(payload, start=1):
            if not isinstance(slide_layers, list):
                continue
            for layer_idx, layer in enumerate(slide_layers):
                if not isinstance(layer, dict):
                    continue
                normalized = dict(layer)
                normalized.setdefault("slide", slide_idx)
                normalized.setdefault("layer_index_in_slide", layer_idx)
                normalized.setdefault("layer_index_global", global_index)
                global_index += 1
                flat.append(normalized)
        return flat

    # Format (1): [ {layer}, ... ]
    flat_layers: List[Dict[str, Any]] = []
    for layer in payload:
        if isinstance(layer, dict):
            flat_layers.append(layer)
    return flat_layers


def _run_font_size_pt(run_info: Dict[str, Any], slide_size_pt: Optional[Tuple[float, float]]) -> Optional[float]:
    """Return font size in points, supporting pt-first and rel-only exports."""
    try:
        v = run_info.get("font_size_pt")
        if v is not None:
            pt = float(v)
            if pt > 0:
                return pt
    except Exception:
        pass
    try:
        rel = run_info.get("font_size_rel")
        if rel is None or slide_size_pt is None:
            return None
        base = _text_rel_base_pt(slide_size_pt)
        pt = float(rel) * base * float(_REL_SCALE)
        return pt if pt > 0 else None
    except Exception:
        return None


def _infer_assets_dir(layers_json: Path, assets_dir_arg: Optional[str]) -> Path:
    if assets_dir_arg:
        return Path(assets_dir_arg).expanduser().resolve()
    return (layers_json.parent / "assets").resolve()


def _embed_extracted_fonts_into_pptx(layers_json: Path, dst_pptx: Path) -> int:
    """Embed fonts using only extracted `fonts.json` + `fonts/` assets."""
    fonts_json = layers_json.parent / "fonts.json"
    fonts_dir = layers_json.parent / "fonts"
    if not fonts_json.exists() or not fonts_dir.exists() or not dst_pptx.exists():
        return 0

    try:
        font_entries = json.loads(fonts_json.read_text(encoding="utf-8"))
    except Exception:
        return 0
    if not isinstance(font_entries, list) or not font_entries:
        return 0

    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
    NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
    pres_xml_path = "ppt/presentation.xml"
    pres_rels_path = "ppt/_rels/presentation.xml.rels"
    content_types_path = "[Content_Types].xml"

    def _next_rid(existing: set[str]) -> str:
        n = 1
        while True:
            rid = f"rId{n}"
            if rid not in existing:
                return rid
            n += 1

    def _find_embedded_font_lst(root: ET.Element) -> Optional[ET.Element]:
        for el in root.iter():
            if el.tag == f"{{{NS_P}}}embeddedFontLst":
                return el
        return None

    def _remove_child(root: ET.Element, target: ET.Element) -> None:
        for parent in root.iter():
            for child in list(parent):
                if child is target:
                    parent.remove(child)
                    return

    def _insert_embedded_font_lst(root: ET.Element, font_lst: ET.Element) -> None:
        preferred_after = {
            f"{{{NS_P}}}notesSz",
            f"{{{NS_P}}}sldSz",
            f"{{{NS_P}}}sldIdLst",
        }
        for idx in range(len(root) - 1, -1, -1):
            if root[idx].tag in preferred_after:
                root.insert(idx + 1, font_lst)
                return
        for idx, child in enumerate(list(root)):
            if child.tag == f"{{{NS_P}}}defaultTextStyle":
                root.insert(idx, font_lst)
                return
        root.append(font_lst)

    prepared: List[Dict[str, Any]] = []
    for item in font_entries:
        if not isinstance(item, dict):
            continue
        saved_path = Path(str(item.get("saved_path") or "")).expanduser()
        if not saved_path.is_absolute():
            saved_path = (fonts_dir / saved_path.name).resolve()
        if not saved_path.exists():
            alt = (fonts_dir / saved_path.name).resolve()
            if alt.exists():
                saved_path = alt
            else:
                continue
        rel_target = str(item.get("rel_target") or "").strip() or f"fonts/{saved_path.name}"
        typeface = str(item.get("typeface") or saved_path.stem.split("__", 1)[0]).strip() or "Unknown Font"
        style = str(item.get("style") or "regular").strip() or "regular"
        rel_type = str(item.get("rel_type") or "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font")
        prepared.append(
            {
                "saved_path": saved_path,
                "rel_target": rel_target,
                "typeface": typeface,
                # Keep the original protocol: style tag comes from extracted metadata
                # and typeface stays exactly as exported.
                "style": style,
                "rel_type": rel_type,
            }
        )
    if not prepared:
        return 0

    tmp_dir = dst_pptx.parent
    with tempfile.NamedTemporaryFile(prefix=dst_pptx.stem + ".", suffix=".tmp.pptx", dir=str(tmp_dir), delete=False) as tf:
        tmp_path = Path(tf.name)

    injected = 0
    try:
        with zipfile.ZipFile(dst_pptx, "r") as zin, zipfile.ZipFile(tmp_path, "w") as zout:
            try:
                dst_pres_root = ET.fromstring(zin.read(pres_xml_path))
                dst_rels_root = ET.fromstring(zin.read(pres_rels_path))
                dst_ct_root = ET.fromstring(zin.read(content_types_path))
            except Exception:
                return 0

            dst_existing = _find_embedded_font_lst(dst_pres_root)
            if dst_existing is not None:
                _remove_child(dst_pres_root, dst_existing)
            for rel in list(dst_rels_root):
                rtype = rel.get("Type") or ""
                if rtype.endswith("/font") or "relationships/font" in rtype:
                    dst_rels_root.remove(rel)

            dst_pres_root.set("embedTrueTypeFonts", "1")
            dst_pres_root.set("saveSubsetFonts", "1")

            existing_rids = {rel.get("Id") for rel in list(dst_rels_root) if rel.get("Id")}
            embedded_font_lst = ET.Element(f"{{{NS_P}}}embeddedFontLst")

            # Group by typeface, then attach available styles under one embeddedFont node.
            by_face: Dict[str, List[Dict[str, Any]]] = defaultdict(list)
            for item in prepared:
                by_face[item["typeface"]].append(item)

            extra_written_paths: set[str] = set()
            font_exts_needed: set[str] = set()
            for typeface, items in by_face.items():
                ef = ET.Element(f"{{{NS_P}}}embeddedFont")
                font_el = ET.SubElement(ef, f"{{{NS_P}}}font")
                font_el.set("typeface", typeface)
                used_styles = set()
                for item in items:
                    style = item["style"]
                    if style in used_styles:
                        continue
                    used_styles.add(style)
                    rid = _next_rid(existing_rids)
                    existing_rids.add(rid)
                    style_el = ET.SubElement(ef, f"{{{NS_P}}}{style}")
                    style_el.set(f"{{{NS_R}}}id", rid)

                    rel_el = ET.Element(f"{{{NS_REL}}}Relationship")
                    rel_el.set("Id", rid)
                    rel_el.set("Type", item["rel_type"])
                    rel_el.set("Target", item["rel_target"])
                    dst_rels_root.append(rel_el)

                    part_zip_path = f"ppt/{item['rel_target'].lstrip('/')}"
                    data = item["saved_path"].read_bytes()
                    zout.writestr(part_zip_path, data)
                    extra_written_paths.add(part_zip_path)
                    injected += 1

                    ext = Path(part_zip_path).suffix.lower().lstrip(".")
                    if ext:
                        font_exts_needed.add(ext)

                embedded_font_lst.append(ef)

            _insert_embedded_font_lst(dst_pres_root, embedded_font_lst)

            existing_defaults = set()
            for el in dst_ct_root.iter():
                if el.tag.endswith("Default") and el.get("Extension"):
                    existing_defaults.add((el.get("Extension") or "").lower())
            existing_overrides = set()
            for el in dst_ct_root.iter():
                if el.tag.endswith("Override") and el.get("PartName"):
                    existing_overrides.add(str(el.get("PartName")))
            default_content_types = {
                "fntdata": "application/x-fontdata",
                "odttf": "application/vnd.openxmlformats-officedocument.obfuscatedFont",
                "ttf": "application/x-font-ttf",
                "otf": "application/x-font-otf",
            }
            for ext in sorted(font_exts_needed):
                if ext in existing_defaults:
                    continue
                ctype = default_content_types.get(ext)
                if not ctype:
                    continue
                el = ET.Element(f"{{{NS_CT}}}Default")
                el.set("Extension", ext)
                el.set("ContentType", ctype)
                dst_ct_root.append(el)

            # Also add per-part overrides for font parts. PowerPoint is sensitive to
            # missing/incorrect content-type declarations on these binaries.
            for part_zip_path in sorted(extra_written_paths):
                part_name = "/" + part_zip_path
                if part_name in existing_overrides:
                    continue
                ext = Path(part_zip_path).suffix.lower().lstrip(".")
                ctype = default_content_types.get(ext)
                if not ctype:
                    continue
                ov = ET.Element(f"{{{NS_CT}}}Override")
                ov.set("PartName", part_name)
                ov.set("ContentType", ctype)
                dst_ct_root.append(ov)

            patched_paths = {pres_xml_path, pres_rels_path, content_types_path}
            for zi in zin.infolist():
                if zi.filename in patched_paths or zi.filename in extra_written_paths:
                    continue
                zout.writestr(zi, zin.read(zi.filename))

            zout.writestr(pres_xml_path, ET.tostring(dst_pres_root, encoding="utf-8", xml_declaration=True))
            zout.writestr(pres_rels_path, ET.tostring(dst_rels_root, encoding="utf-8", xml_declaration=True))
            zout.writestr(content_types_path, ET.tostring(dst_ct_root, encoding="utf-8", xml_declaration=True))

        shutil.move(str(tmp_path), str(dst_pptx))
    finally:
        try:
            if tmp_path.exists():
                tmp_path.unlink()
        except Exception:
            pass

    return injected


def _read_canvas_size_pt(layers: List[Dict[str, Any]]) -> Optional[Tuple[float, float]]:
    """
    Return the canonical slide size from a synthetic slide_canvas layer if present.
    The exported JSON adds one such layer per slide, all with the same page size.
    """
    for layer in layers:
        if layer.get("kind") != "slide_canvas":
            continue
        width_emu = layer.get("canvas_width_emu")
        height_emu = layer.get("canvas_height_emu")
        if width_emu is None or height_emu is None:
            continue
        try:
            return float(width_emu) / EMU_PER_PT, float(height_emu) / EMU_PER_PT
        except Exception:
            continue
    return None


def _extract_slide_background_fill(slide_layers: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    for layer in slide_layers:
        if layer.get("kind") != "slide_canvas":
            continue
        bg = layer.get("background_fill")
        if isinstance(bg, dict):
            return bg
    return None


def _extract_slide_background_spec(slide_layers: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    for layer in slide_layers:
        if layer.get("kind") != "slide_canvas":
            continue
        spec = layer.get("background_spec")
        if isinstance(spec, dict):
            return spec
    return None


def _compute_canvas_shift(
    layers: List[Dict[str, Any]], slide_size_pt: Optional[Tuple[float, float]] = None
) -> Tuple[float, float, float, float]:
    """
    Compute (shift_x, shift_y, slide_w, slide_h) in EMU.
    We normalize negative coords by shifting everything into positive space.
    """
    min_left = 0
    min_top = 0
    max_right = 0
    max_bottom = 0
    canvas_size = _read_canvas_size_pt(layers)
    for layer in layers:
        if layer.get("kind") == "slide_canvas":
            continue
        box = layer.get("box")
        b = _ratio_box_to_pt(box, slide_size_pt or canvas_size)
        if not b:
            continue
        l, t, w, h = b
        min_left = min(min_left, float(l.pt))
        min_top = min(min_top, float(t.pt))
        max_right = max(max_right, float(l.pt) + float(w.pt))
        max_bottom = max(max_bottom, float(t.pt) + float(h.pt))
    shift_x = -min_left if min_left < 0 else 0
    shift_y = -min_top if min_top < 0 else 0
    slide_w = max_right + shift_x
    slide_h = max_bottom + shift_y
    if canvas_size is not None:
        canvas_w, canvas_h = canvas_size
        slide_w = max(slide_w, canvas_w)
        slide_h = max(slide_h, canvas_h)
    # Guard against empty values
    slide_w = max(slide_w, 720.0)  # 10 inches default
    slide_h = max(slide_h, 540.0)  # 7.5 inches default
    return shift_x, shift_y, slide_w, slide_h


def _shift_box_pt(
    box: Optional[Dict[str, Any]],
    shift_x: float,
    shift_y: float,
    slide_size_pt: Optional[Tuple[float, float]] = None,
) -> Optional[Tuple[Pt, Pt, Pt, Pt]]:
    b = _ratio_box_to_pt(box, slide_size_pt)
    if not b:
        return None
    l, t, w, h = b
    return Pt(float(l.pt) + shift_x), Pt(float(t.pt) + shift_y), w, h


def _xml_spec_to_oxml(spec: Dict[str, Any]):
    tag = spec.get("tag")
    if not tag:
        return None
    ns_name = spec.get("ns") or "a"
    el = OxmlElement(f"{str(ns_name)}:{str(tag)}")
    for k, v in (spec.get("attrs") or {}).items():
        el.set(str(k), str(v))
    if spec.get("text") is not None:
        el.text = str(spec.get("text"))
    for child_spec in spec.get("children") or []:
        if not isinstance(child_spec, dict):
            continue
        child_el = _xml_spec_to_oxml(child_spec)
        if child_el is not None:
            el.append(child_el)
    return el


def _apply_slide_background_fill(slide, background_fill: Optional[Dict[str, Any]], background_spec: Optional[Dict[str, Any]] = None) -> None:
    if isinstance(background_spec, dict):
        try:
            bg_el = _xml_spec_to_oxml(background_spec)
            c_sld = slide._element.find(qn("p:cSld"))
            if bg_el is not None and c_sld is not None:
                old_bg = c_sld.find(qn("p:bg"))
                if old_bg is not None:
                    c_sld.remove(old_bg)
                # Put p:bg before spTree when possible.
                inserted = False
                for idx, child in enumerate(list(c_sld)):
                    if child.tag == qn("p:spTree"):
                        c_sld.insert(idx, bg_el)
                        inserted = True
                        break
                if not inserted:
                    c_sld.append(bg_el)
                return
        except Exception:
            pass

    if not background_fill:
        return
    if background_fill.get("type") != "solid":
        return
    try:
        fill = slide.background.fill
        fill.solid()
    except Exception:
        return

    color_type = background_fill.get("color_type")
    if color_type == "RGB" and background_fill.get("rgb"):
        try:
            fill.fore_color.rgb = RGBColor.from_string(str(background_fill["rgb"]))
        except Exception:
            pass
    elif color_type == "SCHEME" and background_fill.get("scheme"):
        try:
            sppr = slide._element.find(".//p:bg/p:bgPr", XML_NS)
            if sppr is not None:
                for child in list(sppr):
                    if child.tag == qn("a:solidFill"):
                        sppr.remove(child)
                solid = OxmlElement("a:solidFill")
                clr = OxmlElement("a:schemeClr")
                clr.set("val", str(background_fill["scheme"]))
                solid.append(clr)
                sppr.append(solid)
        except Exception:
            pass


def _parse_alignment(value: Optional[str]):
    if not value:
        return None
    token = str(value).split(" ", 1)[0].upper()
    mapping = {
        "LEFT": PP_ALIGN.LEFT,
        "CENTER": PP_ALIGN.CENTER,
        "RIGHT": PP_ALIGN.RIGHT,
        "JUSTIFY": PP_ALIGN.JUSTIFY,
        "JUSTIFY_LOW": PP_ALIGN.JUSTIFY_LOW,
        "DISTRIBUTE": PP_ALIGN.DISTRIBUTE,
        "THAI_DISTRIBUTE": PP_ALIGN.THAI_DISTRIBUTE,
    }
    return mapping.get(token)


def _apply_spacing_node(
    p_pr,
    tag_name: str,
    spacing: Optional[Dict[str, Any]],
    slide_size_pt: Optional[Tuple[float, float]] = None,
) -> None:
    if p_pr is None or not spacing:
        return

    tag = qn(f"a:{tag_name}")
    existing = p_pr.find(tag)
    if existing is not None:
        p_pr.remove(existing)

    mode = spacing.get("mode")
    raw = spacing.get("raw")
    rel = spacing.get("rel")
    if mode not in {"points", "percent"}:
        return

    # Support rel-only exports: for point-based spacing, `val` is in 1/100 pt.
    if raw is None and mode == "points" and rel is not None and slide_size_pt is not None:
        try:
            raw = int(round(float(rel) * _text_rel_base_pt(slide_size_pt) * 100.0))
        except Exception:
            raw = None

    if raw is None:
        return

    container = OxmlElement(f"a:{tag_name}")
    child = OxmlElement("a:spcPts" if mode == "points" else "a:spcPct")
    child.set("val", str(int(raw)))
    container.append(child)
    p_pr.append(container)


def _apply_paragraph_bullet(p_pr, para_info: Dict[str, Any]) -> None:
    if p_pr is None:
        return
    for tag in ("buNone", "buChar", "buAutoNum", "buFont"):
        node = p_pr.find(qn(f"a:{tag}"))
        if node is not None:
            p_pr.remove(node)

    bullet = para_info.get("bullet")
    if not isinstance(bullet, dict):
        return
    btype = str(bullet.get("type") or "").lower()
    if btype == "none":
        p_pr.append(OxmlElement("a:buNone"))
        return
    if btype == "char":
        font = bullet.get("font")
        if isinstance(font, dict) and font.get("typeface"):
            el_font = OxmlElement("a:buFont")
            el_font.set("typeface", str(font.get("typeface")))
            p_pr.append(el_font)
        el = OxmlElement("a:buChar")
        if bullet.get("char") is not None:
            el.set("char", str(bullet.get("char")))
        p_pr.append(el)
        return
    if btype == "autonum":
        el = OxmlElement("a:buAutoNum")
        if bullet.get("scheme"):
            el.set("type", str(bullet.get("scheme")))
        if bullet.get("startAt") is not None:
            try:
                el.set("startAt", str(int(bullet.get("startAt"))))
            except Exception:
                pass
        p_pr.append(el)


def _apply_paragraph_format(p, para_info: Dict[str, Any], slide_size_pt: Optional[Tuple[float, float]] = None) -> None:
    alignment = _parse_alignment(para_info.get("alignment"))
    if alignment is not None:
        try:
            p.alignment = alignment
        except Exception:
            pass

    if para_info.get("level") is not None:
        try:
            p.level = int(para_info["level"])
        except Exception:
            pass

    p_pr = p._p.get_or_add_pPr()
    ppr_attrs = para_info.get("ppr_attrs")
    if isinstance(ppr_attrs, dict):
        for key in ("marL", "marR", "indent", "lvl", "algn"):
            val = ppr_attrs.get(key)
            if val is None:
                continue
            try:
                p_pr.set(str(key), str(val))
            except Exception:
                pass
        if slide_size_pt is not None:
            base_emu = _text_rel_base_pt(slide_size_pt) * EMU_PER_PT
            for key in ("marL", "marR", "indent"):
                rel = ppr_attrs.get(f"{key}_rel")
                if rel is None:
                    continue
                try:
                    p_pr.set(key, str(int(round(float(rel) * base_emu))))
                except Exception:
                    pass
    _apply_paragraph_bullet(p_pr, para_info)
    _apply_spacing_node(p_pr, "lnSpc", para_info.get("line_spacing"), slide_size_pt)
    _apply_spacing_node(p_pr, "spcBef", para_info.get("space_before"), slide_size_pt)
    _apply_spacing_node(p_pr, "spcAft", para_info.get("space_after"), slide_size_pt)

    # Apply endParaRPr color as paragraph default when present (theme inheritance).
    end_para = para_info.get("end_para_rpr")
    if isinstance(end_para, dict) and isinstance(end_para.get("color"), dict):
        try:
            end_el = p._p.find(qn("a:endParaRPr"))
            if end_el is None:
                end_el = OxmlElement("a:endParaRPr")
                p._p.append(end_el)
            # remove existing solidFill
            for child in list(end_el):
                if child.tag == qn("a:solidFill"):
                    end_el.remove(child)
            c = end_para.get("color") or {}
            if c.get("type") == "RGB" and c.get("rgb"):
                solid = OxmlElement("a:solidFill")
                clr = OxmlElement("a:srgbClr")
                clr.set("val", str(c.get("rgb")))
                for mod in c.get("mods") or []:
                    if isinstance(mod, dict) and mod.get("op") and mod.get("val") is not None:
                        m = OxmlElement(f"a:{str(mod['op'])}")
                        m.set("val", str(mod.get("val")))
                        clr.append(m)
                solid.append(clr)
                end_el.append(solid)
            elif c.get("type") == "SCHEME" and c.get("scheme"):
                solid = OxmlElement("a:solidFill")
                clr = OxmlElement("a:schemeClr")
                clr.set("val", str(c.get("scheme")))
                for mod in c.get("mods") or []:
                    if isinstance(mod, dict) and mod.get("op") and mod.get("val") is not None:
                        m = OxmlElement(f"a:{str(mod['op'])}")
                        m.set("val", str(mod.get("val")))
                        clr.append(m)
                solid.append(clr)
                end_el.append(solid)
        except Exception:
            pass


def _set_run_font(run, run_info: Dict[str, Any], slide_size_pt: Optional[Tuple[float, float]] = None) -> None:
    # Use exactly what layers.json provides. Do not infer style from font_name,
    # because templates often encode style in the typeface itself (e.g. "Aileron Bold").
    bold_val = run_info.get("bold")
    italic_val = run_info.get("italic")
    def _spec_to_oxml(spec: Dict[str, Any]):
        tag = spec.get("tag")
        if not tag:
            return None
        el = OxmlElement(f"a:{str(tag)}")
        for k, v in (spec.get("attrs") or {}).items():
            el.set(str(k), str(v))
        if spec.get("text") is not None:
            el.text = str(spec.get("text"))
        for child_spec in spec.get("children") or []:
            if not isinstance(child_spec, dict):
                continue
            child_el = _spec_to_oxml(child_spec)
            if child_el is not None:
                el.append(child_el)
        return el

    raw_rpr = run_info.get("raw_rpr")
    if isinstance(raw_rpr, dict) and raw_rpr.get("from_fields"):
        try:
            rpr = OxmlElement("a:rPr")
            # Basic attrs that are not present in our run fields.
            if raw_rpr.get("lang"):
                rpr.set("lang", str(raw_rpr["lang"]))
            # Keep strike unset when it's the default "noStrike" to better match source XML.
            if raw_rpr.get("strike") and str(raw_rpr.get("strike")) != "noStrike":
                rpr.set("strike", str(raw_rpr["strike"]))

            # Bold/italic only when explicitly present in the protocol.
            if bold_val is not None:
                rpr.set("b", "true" if bool(bold_val) else "false")
            if italic_val is not None:
                rpr.set("i", "true" if bool(italic_val) else "false")

            # sz is in hundredth-points (e.g. 1700 == 17pt)
            font_size_pt = _run_font_size_pt(run_info, slide_size_pt)
            if font_size_pt is not None:
                try:
                    rpr.set("sz", str(int(round(float(font_size_pt) * 100))))
                except Exception:
                    pass

            # Underline
            u = run_info.get("underline")
            if isinstance(u, str):
                rpr.set("u", u)
            elif u is False:
                rpr.set("u", "none")
            elif u is True:
                rpr.set("u", "sng")

            # solidFill from run fields
            color = run_info.get("color") or {}
            if color:
                try:
                    solid = OxmlElement("a:solidFill")
                    if color.get("type") == "RGB" and color.get("rgb"):
                        c = OxmlElement("a:srgbClr")
                        c.set("val", str(color["rgb"]))
                        solid.append(c)
                        rpr.append(solid)
                    elif color.get("type") == "SCHEME" and color.get("scheme"):
                        c = OxmlElement("a:schemeClr")
                        c.set("val", str(color["scheme"]))
                        solid.append(c)
                        rpr.append(solid)
                except Exception:
                    pass

            # Fonts: match the source style nodes (latin/ea/cs/sym) using font_name
            font_name = run_info.get("font_name")
            if font_name:
                for tag in ("latin", "ea", "cs", "sym"):
                    el = OxmlElement(f"a:{tag}")
                    el.set("typeface", str(font_name))
                    rpr.append(el)

            # Preserve character spacing/kern. These are frequently present on templates
            # and must be set on the rPr we are constructing (this branch returns early).
            if run_info.get("char_spacing_raw") is not None:
                try:
                    rpr.set("spc", str(int(run_info["char_spacing_raw"])))
                except Exception:
                    pass
            elif run_info.get("char_spacing_rel") is not None and slide_size_pt is not None:
                # DrawingML run spacing is in 1/100 pt.
                try:
                    rel = float(run_info["char_spacing_rel"])
                    spc = int(round(rel * _text_rel_base_pt(slide_size_pt) * 100.0))
                    rpr.set("spc", str(spc))
                except Exception:
                    pass
            if run_info.get("kern_raw") is not None:
                try:
                    rpr.set("kern", str(int(run_info["kern_raw"])))
                except Exception:
                    pass

            old_rpr = run._r.find(qn("a:rPr"))
            if old_rpr is not None:
                run._r.remove(old_rpr)
            run._r.insert(0, rpr)
            return
        except Exception:
            pass

    if isinstance(raw_rpr, str) and raw_rpr.strip():
        try:
            new_rpr = parse_xml(str(raw_rpr).encode("utf-8"))
            old_rpr = run._r.find(qn("a:rPr"))
            if old_rpr is not None:
                run._r.remove(old_rpr)
            # rPr must stay before the text-bearing children if present.
            run._r.insert(0, new_rpr)
            return
        except Exception:
            pass

    raw_rpr_xml = run_info.get("raw_rpr_xml")
    if raw_rpr_xml:
        try:
            new_rpr = parse_xml(str(raw_rpr_xml).encode("utf-8"))
            old_rpr = run._r.find(qn("a:rPr"))
            if old_rpr is not None:
                run._r.remove(old_rpr)
            run._r.insert(0, new_rpr)
            return
        except Exception:
            pass

    f = run.font
    if run_info.get("font_name"):
        font_name = run_info["font_name"]
        f.name = font_name
        # PowerPoint distinguishes font for different scripts (ascii/hAnsi/eastAsia/cs).
        # Setting only run.font.name can still render "different" due to fallback.
        try:
            rPr = run._r.get_or_add_rPr()
            rFonts = rPr.find(qn("a:rFonts"))
            if rFonts is None:
                rFonts = OxmlElement("a:rFonts")
                rPr.insert(0, rFonts)
            rFonts.set("ascii", font_name)
            rFonts.set("hAnsi", font_name)
            rFonts.set("eastAsia", font_name)
            rFonts.set("cs", font_name)
        except Exception:
            # Best-effort; keep going even if python-pptx internals differ.
            pass
    font_size_pt = _run_font_size_pt(run_info, slide_size_pt)
    if font_size_pt is not None:
        try:
            f.size = Pt(float(font_size_pt))
        except Exception:
            pass
    for key in ("bold", "italic", "underline"):
        if run_info.get(key) is not None:
            setattr(f, key, bool(run_info[key]))

    color = run_info.get("color") or {}
    if color:
        try:
            r_pr = run._r.get_or_add_rPr()
            for child in list(r_pr):
                if child.tag == qn("a:solidFill"):
                    r_pr.remove(child)
            def _apply_mods(color_el, mods):
                for mod in mods or []:
                    if not isinstance(mod, dict) or not mod.get("op"):
                        continue
                    el = OxmlElement(f"a:{str(mod['op'])}")
                    el.set("val", str(mod.get("val")))
                    color_el.append(el)
            if color.get("type") == "RGB" and color.get("rgb"):
                solid = OxmlElement("a:solidFill")
                c = OxmlElement("a:srgbClr")
                c.set("val", str(color["rgb"]))
                _apply_mods(c, color.get("mods"))
                solid.append(c)
                r_pr.append(solid)
            elif color.get("type") == "SCHEME" and color.get("scheme"):
                solid = OxmlElement("a:solidFill")
                c = OxmlElement("a:schemeClr")
                c.set("val", str(color["scheme"]))
                _apply_mods(c, color.get("mods"))
                solid.append(c)
                r_pr.append(solid)
        except Exception:
            pass

    try:
        r_pr = run._r.get_or_add_rPr()
        if run_info.get("char_spacing_raw") is not None:
            r_pr.set("spc", str(int(run_info["char_spacing_raw"])))
        if run_info.get("kern_raw") is not None:
            r_pr.set("kern", str(int(run_info["kern_raw"])))
    except Exception:
        pass
    # #region debug-point C:text-token-rebuild
    if "token" in (run_info.get("text") or "").lower():
        try:
            _dbg_report_method_ppt_style(
                "C",
                "combine.py:_set_run_font",
                "[DEBUG] rebuild token run",
                {
                    "text": run_info.get("text"),
                    "font_name": run_info.get("font_name"),
                    "font_size_pt": run_info.get("font_size_pt"),
                    "font_size_rel": run_info.get("font_size_rel"),
                    "font_size_pt_derived": font_size_pt,
                    "bold": run_info.get("bold"),
                    "italic": run_info.get("italic"),
                    "underline": run_info.get("underline"),
                    "color": run_info.get("color"),
                    "rpr_xml": ET.tostring(run._r.get_or_add_rPr(), encoding="unicode"),
                },
            )
        except Exception:
            pass
    # #endregion


def _add_text_layer(slide, layer: Dict[str, Any], shift_x: float, shift_y: float) -> None:
    base_box = layer.get("box")
    box = _shift_box_pt(base_box, shift_x, shift_y, _slide_size_pt_from_slide(slide))
    if not box:
        return
    left, top, width, height = box
    tb = slide.shapes.add_textbox(left, top, width, height)
    # Keep original shape name for debugging and for stable downstream references.
    try:
        name = layer.get("shape_name")
        if name:
            tb._element.nvSpPr.cNvPr.set("name", str(name))
    except Exception:
        pass

    # Preserve raw xfrm rot when available (some templates use negative rot values and
    # PowerPoint can render them with tiny differences vs normalized 0..21600000 space).
    try:
        rot_raw = layer.get("xfrm_rot_raw")
        if rot_raw is not None:
            tb._element.spPr.xfrm.set("rot", str(rot_raw))
    except Exception:
        pass

    # Apply bodyPr (insets/wrap/anchor/autofit) from exported spec, if present.
    try:
        bp = layer.get("body_pr")
        if isinstance(bp, dict) and isinstance(bp.get("attrs"), dict):
            body_pr_el = tb._element.txBody.bodyPr
            # Clear existing attrs, then set exactly what we exported.
            body_pr_el.attrib.clear()
            for k, v in (bp.get("attrs") or {}).items():
                if v is None:
                    continue
                # Some export protocols carry derived "*_rel" helpers that are not valid
                # OOXML attributes. Keeping them can cause PowerPoint to ignore bodyPr.
                if str(k).endswith("_rel"):
                    continue
                body_pr_el.set(str(k), str(v))
            # Clear existing autofit children, then set exported mode.
            for child in list(body_pr_el):
                if child.tag in (qn("a:spAutoFit"), qn("a:normAutoFit"), qn("a:noAutoFit")):
                    body_pr_el.remove(child)
            mode = bp.get("autofit")
            if mode in {"spAutoFit", "normAutoFit", "noAutoFit"}:
                body_pr_el.append(OxmlElement(f"a:{mode}"))
    except Exception:
        pass

    # Apply textbox background fill/line when exported (keeps "boxed text" visuals).
    try:
        sppr = tb._element.spPr
        fill_spec = layer.get("shape_fill")
        line_style = layer.get("shape_line")

        def _apply_mods(color_el, mods) -> None:
            for mod in mods or []:
                if not isinstance(mod, dict) or not mod.get("op") or mod.get("val") is None:
                    continue
                m = OxmlElement(f"a:{str(mod['op'])}")
                m.set("val", str(mod.get("val")))
                color_el.append(m)

        # Fill: default to explicit noFill to avoid theme default fills.
        for child in list(sppr):
            if child.tag in (qn("a:solidFill"), qn("a:noFill")):
                sppr.remove(child)
        if isinstance(fill_spec, dict) and fill_spec.get("type"):
            solid = OxmlElement("a:solidFill")
            if fill_spec.get("type") == "srgb" and fill_spec.get("rgb"):
                clr = OxmlElement("a:srgbClr")
                clr.set("val", str(fill_spec.get("rgb")).lstrip("#"))
                _apply_mods(clr, fill_spec.get("mods"))
                solid.append(clr)
                sppr.append(solid)
            elif fill_spec.get("type") == "scheme" and fill_spec.get("scheme"):
                clr = OxmlElement("a:schemeClr")
                clr.set("val", str(fill_spec.get("scheme")))
                _apply_mods(clr, fill_spec.get("mods"))
                solid.append(clr)
                sppr.append(solid)
            else:
                sppr.append(OxmlElement("a:noFill"))
        else:
            sppr.append(OxmlElement("a:noFill"))

        # Line: default to explicit noFill to avoid theme default outlines.
        ln = sppr.get_or_add_ln()
        for child in list(ln):
            if child.tag in (qn("a:solidFill"), qn("a:noFill"), qn("a:prstDash"), qn("a:headEnd"), qn("a:tailEnd")):
                ln.remove(child)
        if isinstance(line_style, dict):
            if line_style.get("width_pt") is not None:
                try:
                    ln.set("w", str(int(round(float(line_style["width_pt"]) * EMU_PER_PT))))
                except Exception:
                    pass
            c = line_style.get("color")
            if isinstance(c, dict) and c.get("type"):
                solid = OxmlElement("a:solidFill")
                if c.get("type") == "srgb" and c.get("rgb"):
                    clr = OxmlElement("a:srgbClr")
                    clr.set("val", str(c.get("rgb")).lstrip("#"))
                    _apply_mods(clr, c.get("mods"))
                    solid.append(clr)
                    ln.append(solid)
                elif c.get("type") == "scheme" and c.get("scheme"):
                    clr = OxmlElement("a:schemeClr")
                    clr.set("val", str(c.get("scheme")))
                    _apply_mods(clr, c.get("mods"))
                    solid.append(clr)
                    ln.append(solid)
                else:
                    ln.append(OxmlElement("a:noFill"))
            else:
                ln.append(OxmlElement("a:noFill"))
            if line_style.get("dash"):
                d = OxmlElement("a:prstDash")
                d.set("val", str(line_style.get("dash")))
                ln.append(d)
            for tag, key in (("headEnd", "head_end"), ("tailEnd", "tail_end")):
                spec = line_style.get(key)
                if isinstance(spec, dict) and spec.get("type"):
                    e = OxmlElement(f"a:{tag}")
                    e.set("type", str(spec.get("type")))
                    if spec.get("w"):
                        e.set("w", str(spec.get("w")))
                    if spec.get("len"):
                        e.set("len", str(spec.get("len")))
                    ln.append(e)
        else:
            ln.append(OxmlElement("a:noFill"))
    except Exception:
        pass

    # If we have an explicit raw xfrm rot, keep it as-is (don't normalize via tb.rotation).
    if layer.get("xfrm_rot_raw") is None and layer.get("rotation_deg") is not None:
        try:
            tb.rotation = float(layer.get("rotation_deg") or 0.0)
        except Exception:
            pass
    tf = tb.text_frame
    _apply_shape_flip(tb, layer)
    _fill_text_frame(tf, layer, _slide_size_pt_from_slide(slide))


def _fill_text_frame(tf, text_info: Dict[str, Any], slide_size_pt: Optional[Tuple[float, float]] = None) -> None:
    tf.clear()
    body_pr = text_info.get("body_pr") if isinstance(text_info.get("body_pr"), dict) else None
    # If we have a full bodyPr spec (attrs/autofit), we apply it at the OOXML level
    # in _add_text_layer() to avoid python-pptx injecting defaults like wrap="square".
    # Keep the text-frame defaults minimal here.
    if not (isinstance(body_pr, dict) and isinstance(body_pr.get("attrs"), dict)):
        try:
            tf.word_wrap = True
        except Exception:
            pass
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
    paragraphs = text_info.get("paragraphs") or []
    if not paragraphs:
        tf.text = text_info.get("text", "") or ""
        return
    first = True
    for para_info in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        _apply_paragraph_format(p, para_info, slide_size_pt)
        para_runs = para_info.get("runs") or []
        if not para_runs and para_info.get("text"):
            p.text = para_info.get("text", "") or ""
            continue
        for rinfo in sorted(para_runs, key=lambda x: int(x.get("run_index", 0))):
            rr = p.add_run()
            rr.text = rinfo.get("text", "") or ""
            _set_run_font(rr, rinfo, slide_size_pt)


def _append_color_fill(parent, color_spec: Optional[Dict[str, Any]]) -> bool:
    if not isinstance(color_spec, dict) or not color_spec.get("type"):
        return False
    solid = OxmlElement("a:solidFill")
    kind = str(color_spec.get("type") or "").lower()
    if kind == "srgb" and color_spec.get("rgb"):
        clr = OxmlElement("a:srgbClr")
        clr.set("val", str(color_spec.get("rgb")).lstrip("#"))
    elif kind == "scheme" and color_spec.get("scheme"):
        clr = OxmlElement("a:schemeClr")
        clr.set("val", str(color_spec.get("scheme")))
    else:
        return False
    for mod in color_spec.get("mods") or []:
        if not isinstance(mod, dict) or not mod.get("op") or mod.get("val") is None:
            continue
        m = OxmlElement(f"a:{str(mod['op'])}")
        m.set("val", str(mod.get("val")))
        clr.append(m)
    solid.append(clr)
    parent.append(solid)
    return True


def _apply_table_cell_fill(cell, fill_spec: Optional[Dict[str, Any]]) -> None:
    try:
        # If no explicit fill spec was exported, keep table style/theme defaults.
        if not (isinstance(fill_spec, dict) and fill_spec.get("type")):
            return

        fill = cell.fill
        fill.solid()
        kind = str(fill_spec.get("type") or "").lower()
        if kind == "srgb" and fill_spec.get("rgb"):
            fill.fore_color.rgb = RGBColor.from_string(str(fill_spec.get("rgb")).lstrip("#"))
            return

        # Non-sRGB fills (scheme/mods) need OOXML injection.
        tcpr = cell._tc.get_or_add_tcPr()
        for child in list(tcpr):
            if child.tag in (qn("a:solidFill"), qn("a:noFill")):
                tcpr.remove(child)
        if not _append_color_fill(tcpr, fill_spec):
            tcpr.append(OxmlElement("a:noFill"))
    except Exception:
        pass


def _apply_table_cell_borders(cell, borders: Optional[Dict[str, Any]]) -> None:
    try:
        # If no explicit border spec was exported, keep table style/theme defaults.
        if not (isinstance(borders, dict) and borders):
            return

        tcpr = cell._tc.get_or_add_tcPr()
        for tag in ("lnL", "lnR", "lnT", "lnB"):
            for child in list(tcpr):
                if child.tag == qn(f"a:{tag}"):
                    tcpr.remove(child)
        for xml_tag, key in (("lnL", "l"), ("lnR", "r"), ("lnT", "t"), ("lnB", "b")):
            spec = borders.get(key) if isinstance(borders, dict) else None
            if not isinstance(spec, dict):
                continue
            ln = OxmlElement(f"a:{xml_tag}")
            # Match PPT-exported defaults to avoid PowerPoint rendering differences vs WPS.
            ln.set("cmpd", "sng")
            ln.set("algn", "ctr")
            ln.set("cap", "flat")
            if spec.get("width_pt") is not None:
                try:
                    ln.set("w", str(int(round(float(spec["width_pt"]) * EMU_PER_PT))))
                except Exception:
                    pass
            if not _append_color_fill(ln, spec.get("color")):
                ln.append(OxmlElement("a:noFill"))
            d = OxmlElement("a:prstDash")
            d.set("val", str(spec.get("dash") or "solid"))
            ln.append(d)
            ln.append(OxmlElement("a:round"))
            head = OxmlElement("a:headEnd")
            head.set("type", "none")
            head.set("w", "med")
            head.set("len", "med")
            ln.append(head)
            tail = OxmlElement("a:tailEnd")
            tail.set("type", "none")
            tail.set("w", "med")
            tail.set("len", "med")
            ln.append(tail)
            tcpr.append(ln)
    except Exception:
        pass


def _add_table_layer(
    slide,
    layer: Dict[str, Any],
    shift_x: float,
    shift_y: float,
) -> None:
    base_box = layer.get("box")
    box = _shift_box_pt(base_box, shift_x, shift_y, _slide_size_pt_from_slide(slide))
    if not box:
        return
    rows = int(layer.get("rows") or 0)
    cols = int(layer.get("cols") or 0)
    if rows <= 0 or cols <= 0:
        return
    left, top, width, height = box
    shp = slide.shapes.add_table(rows, cols, left, top, width, height)
    try:
        name = layer.get("shape_name")
        if name:
            shp._element.nvGraphicFramePr.cNvPr.set("name", str(name))
    except Exception:
        pass
    table = shp.table
    _apply_shape_flip(shp, layer)

    # Table style: keep explicit style id when exported; otherwise strip the default style id
    # python-pptx injects. Sticker Chart source uses no tableStyleId and PowerPoint can apply
    # a theme style that hides borders when a default style is present.
    try:
        tbl_prs = shp._element.xpath(".//a:tblPr")
        if tbl_prs:
            tbl_pr = tbl_prs[0]
            sid = tbl_pr.find(qn("a:tableStyleId"))
            style_id = layer.get("table_style_id")
            if style_id:
                if sid is None:
                    sid = OxmlElement("a:tableStyleId")
                    tbl_pr.append(sid)
                sid.text = str(style_id)
            else:
                if sid is not None:
                    tbl_pr.remove(sid)
    except Exception:
        pass
    try:
        if layer.get("rotation_deg") is not None:
            shp.rotation = float(layer.get("rotation_deg") or 0.0)
    except Exception:
        pass

    box_w_pt = max(0.0, float(width.pt))
    box_h_pt = max(0.0, float(height.pt))
    raw_col_widths = [float(v or 0.0) for v in (layer.get("col_widths_pt") or [])[:cols]]
    raw_row_heights = [float(v or 0.0) for v in (layer.get("row_heights_pt") or [])[:rows]]

    if len(raw_col_widths) < cols:
        raw_col_widths.extend([0.0] * (cols - len(raw_col_widths)))
    if len(raw_row_heights) < rows:
        raw_row_heights.extend([0.0] * (rows - len(raw_row_heights)))

    col_total = sum(v for v in raw_col_widths if v > 0)
    row_total = sum(v for v in raw_row_heights if v > 0)

    if box_w_pt > 0:
        if col_total > 0:
            col_widths = [(v if v > 0 else 0.0) * box_w_pt / col_total for v in raw_col_widths]
        else:
            col_widths = [box_w_pt / cols] * cols
    else:
        col_widths = raw_col_widths

    if box_h_pt > 0:
        if row_total > 0:
            row_heights = [(v if v > 0 else 0.0) * box_h_pt / row_total for v in raw_row_heights]
        else:
            row_heights = [box_h_pt / rows] * rows
    else:
        row_heights = raw_row_heights

    for i, w_pt in enumerate(col_widths):
        if i >= len(table.columns):
            break
        try:
            table.columns[i].width = Pt(float(w_pt))
        except Exception:
            pass
    for i, h_pt in enumerate(row_heights):
        if i >= len(table.rows):
            break
        try:
            table.rows[i].height = Pt(float(h_pt))
        except Exception:
            pass

    cells = layer.get("cells") or []
    for ri, row in enumerate(cells):
        if ri >= rows or not isinstance(row, list):
            continue
        for ci, cell_info in enumerate(row):
            if ci >= cols or not isinstance(cell_info, dict):
                continue
            cell = table.cell(ri, ci)
            _apply_table_cell_fill(cell, cell_info.get("fill"))
            _apply_table_cell_borders(cell, cell_info.get("borders"))
            tf = cell.text_frame
            try:
                tf.margin_left = 0
                tf.margin_right = 0
                tf.margin_top = 0
                tf.margin_bottom = 0
            except Exception:
                pass
            _fill_text_frame(
                tf,
                {"text": cell_info.get("text"), "paragraphs": cell_info.get("paragraphs") or []},
                _slide_size_pt_from_slide(slide),
            )


def _apply_shape_flip(shape, layer: Dict[str, Any]) -> None:
    flip_h = bool(layer.get("flip_h"))
    flip_v = bool(layer.get("flip_v"))
    if not flip_h and not flip_v:
        return
    try:
        xfrm = shape._element.spPr.xfrm
    except Exception:
        xfrm = None
    if xfrm is None:
        return
    if flip_h:
        xfrm.set("flipH", "true")
    elif "flipH" in xfrm.attrib:
        del xfrm.attrib["flipH"]
    if flip_v:
        xfrm.set("flipV", "true")
    elif "flipV" in xfrm.attrib:
        del xfrm.attrib["flipV"]


def _next_shape_id(slide) -> int:
    ids: List[int] = []
    for el in slide.shapes._spTree.iter():
        if el.tag == qn("p:cNvPr"):
            try:
                ids.append(int(el.get("id")))
            except Exception:
                pass
    return (max(ids) + 1) if ids else 1


def _clamp_box_to_slide(box: Tuple[Pt, Pt, Pt, Pt], slide) -> Optional[Tuple[Pt, Pt, Pt, Pt]]:
    """Clamp an axis-aligned box to the slide canvas.

    This is intentionally simple (no rotation-aware clipping): if a line/geo shape's
    box exceeds the slide bounds, we shrink it so its right/bottom stay within
    the canvas. If fully outside, return None.
    """
    try:
        slide_w, slide_h = _slide_size_pt_from_slide(slide)
    except Exception:
        return box

    left, top, width, height = box
    l = float(left.pt)
    t = float(top.pt)
    r = l + float(width.pt)
    b = t + float(height.pt)

    cl = max(0.0, l)
    ct = max(0.0, t)
    cr = min(slide_w, r)
    cb = min(slide_h, b)
    if cr <= cl or cb <= ct:
        return None
    return Pt(cl), Pt(ct), Pt(cr - cl), Pt(cb - ct)


def _add_graph_layer(slide, layer: Dict[str, Any], shift_x: float, shift_y: float) -> None:
    shape_xml = layer.get("shape_xml")
    if not shape_xml:
        return
    base_box = layer.get("box")
    box = _shift_box_pt(base_box, shift_x, shift_y, _slide_size_pt_from_slide(slide))

    is_raw_xml = isinstance(shape_xml, str)
    # New compact graph schema: { "spPr": { "xfrm": {...}, "custGeom": {...} } }
    if isinstance(shape_xml, dict) and isinstance(shape_xml.get("spPr"), dict):
        if not box:
            return
        sppr = shape_xml.get("spPr") or {}
        xfrm_spec = sppr.get("xfrm") or {}
        cust = sppr.get("custGeom") or {}
        path_lst = cust.get("pathLst") or []

        # Compute placement from `box` (no clamping; allow out-of-canvas).
        left, top, width, height = box
        off_x = int(round(float(left.pt) * EMU_PER_PT))
        off_y = int(round(float(top.pt) * EMU_PER_PT))
        ext_cx = int(round(float(width.pt) * EMU_PER_PT))
        ext_cy = int(round(float(height.pt) * EMU_PER_PT))

        flip_h_attr = ' flipH="true"' if bool(xfrm_spec.get("flipH")) else ""
        flip_v_attr = ' flipV="true"' if bool(xfrm_spec.get("flipV")) else ""
        rot_attr = ""
        try:
            rot = float(layer.get("rotation_deg") or 0.0)
            if rot:
                rot_attr = f' rot="{int(round(rot * 60000))}"'
        except Exception:
            rot_attr = ""

        # Canvas size for decoding relative w/h.
        slide_w_pt, slide_h_pt = _slide_size_pt_from_slide(slide)
        canvas_w_emu = int(round(float(slide_w_pt) * EMU_PER_PT))
        canvas_h_emu = int(round(float(slide_h_pt) * EMU_PER_PT))

        path_chunks: List[str] = []
        for path in path_lst:
            if not isinstance(path, dict):
                continue
            try:
                pw = float(path.get("w") or 0.0)
                ph = float(path.get("h") or 0.0)
            except Exception:
                continue
            if pw <= 0 or ph <= 0:
                continue
            path_w = max(1, int(round(pw * canvas_w_emu)))
            path_h = max(1, int(round(ph * canvas_h_emu)))

            cmd_chunks: List[str] = []
            for cmd in path.get("commands") or []:
                if not isinstance(cmd, dict):
                    continue
                ctype = cmd.get("type")
                if ctype in {"moveTo", "lnTo"} and isinstance(cmd.get("pt"), dict):
                    pt = cmd["pt"]
                    x = int(round(float(pt.get("x") or 0.0) * path_w))
                    y = int(round(float(pt.get("y") or 0.0) * path_h))
                    cmd_chunks.append(f'<a:{ctype}><a:pt x="{x}" y="{y}"/></a:{ctype}>')
                elif ctype == "cubicBezTo" and isinstance(cmd.get("pts"), list) and len(cmd.get("pts")) == 3:
                    pts_xml = ""
                    for pt in cmd["pts"]:
                        if not isinstance(pt, dict):
                            continue
                        x = int(round(float(pt.get("x") or 0.0) * path_w))
                        y = int(round(float(pt.get("y") or 0.0) * path_h))
                        pts_xml += f'<a:pt x="{x}" y="{y}"/>'
                    cmd_chunks.append(f"<a:cubicBezTo>{pts_xml}</a:cubicBezTo>")
                elif ctype == "close":
                    cmd_chunks.append("<a:close/>")
            if cmd_chunks:
                path_chunks.append(f'<a:path w="{path_w}" h="{path_h}">' + "".join(cmd_chunks) + "</a:path>")

        if not path_chunks:
            return

        def _mods_xml(spec: Any) -> str:
            mods = ""
            if not isinstance(spec, dict):
                return mods
            for mod in spec.get("mods") or []:
                if isinstance(mod, dict) and mod.get("op") and mod.get("val") is not None:
                    mods += f'<a:{str(mod["op"])} val="{str(mod["val"])}"/>'
            return mods

        def _solid_fill_xml_from_color_spec(spec: Any) -> str:
            if not isinstance(spec, dict) or not spec.get("type"):
                return ""
            mods = _mods_xml(spec)
            if spec.get("type") == "srgb" and spec.get("rgb"):
                return f'<a:solidFill><a:srgbClr val="{str(spec.get("rgb")).lstrip("#")}">{mods}</a:srgbClr></a:solidFill>'
            if spec.get("type") == "scheme" and spec.get("scheme"):
                return f'<a:solidFill><a:schemeClr val="{str(spec.get("scheme"))}">{mods}</a:schemeClr></a:solidFill>'
            return ""

        # Fill/line are stored at layer level (not inside shape_xml).
        fill = layer.get("fill")
        fill_xml = "<a:noFill/>"
        if isinstance(fill, dict):
            # Legacy: {"type":"solidFill","color":{"type":"srgbClr","val":"#RRGGBB"}}
            if fill.get("type") == "solidFill":
                color = fill.get("color") if isinstance(fill.get("color"), dict) else None
                if isinstance(color, dict) and color.get("type") == "srgbClr" and color.get("val"):
                    val = str(color.get("val")).lstrip("#")
                    fill_xml = f'<a:solidFill><a:srgbClr val="{val}"/></a:solidFill>'
            else:
                # Current: {"type":"srgb"/"scheme", ... , "mods":[...]}
                solid = _solid_fill_xml_from_color_spec(fill)
                if solid:
                    fill_xml = solid

        line_style = layer.get("line")
        ln_xml = "<a:ln><a:noFill/></a:ln>"
        if isinstance(line_style, dict):
            ln_attrs = ""
            if line_style.get("width_pt") is not None:
                try:
                    ln_attrs = f' w="{int(round(float(line_style["width_pt"]) * EMU_PER_PT))}"'
                except Exception:
                    ln_attrs = ""
            ln_fill = "<a:noFill/>"
            solid = _solid_fill_xml_from_color_spec(line_style.get("color"))
            if solid:
                ln_fill = solid
            dash_xml = ""
            if line_style.get("dash"):
                dash_xml = f'<a:prstDash val="{str(line_style["dash"])}"/>'
            # Arrowheads for non-line geo shapes are rare but keep it consistent.
            end_xml = ""
            head = line_style.get("head_end")
            if isinstance(head, dict) and head.get("type"):
                end_xml += f'<a:headEnd type="{str(head.get("type"))}"' + (
                    f' w="{str(head.get("w"))}"' if head.get("w") else ""
                ) + (f' len="{str(head.get("len"))}"' if head.get("len") else "") + "/>"
            tail = line_style.get("tail_end")
            if isinstance(tail, dict) and tail.get("type"):
                end_xml += f'<a:tailEnd type="{str(tail.get("type"))}"' + (
                    f' w="{str(tail.get("w"))}"' if tail.get("w") else ""
                ) + (f' len="{str(tail.get("len"))}"' if tail.get("len") else "") + "/>"
            ln_xml = f"<a:ln{ln_attrs}>{ln_fill}{dash_xml}{end_xml}</a:ln>"

        shape_name = str(layer.get("shape_name") or "shape")
        shape_name_xml = (
            shape_name.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&apos;")
        )
        sp_xml = (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f"<p:nvSpPr><p:cNvPr id=\"{_next_shape_id(slide)}\" name=\"{shape_name_xml}\"/>"
            "<p:cNvSpPr/><p:nvPr/></p:nvSpPr>"
            f"<p:spPr><a:xfrm{rot_attr}{flip_h_attr}{flip_v_attr}>"
            f'<a:off x="{off_x}" y="{off_y}"/><a:ext cx="{ext_cx}" cy="{ext_cy}"/>'
            "</a:xfrm>"
            "<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>"
            '<a:rect l="l" t="t" r="r" b="b"/>'
            f"<a:pathLst>{''.join(path_chunks)}</a:pathLst></a:custGeom>"
            f"{fill_xml}{ln_xml}</p:spPr></p:sp>"
        )
        try:
            shape_el = parse_xml(sp_xml.encode("utf-8"))
        except Exception:
            return
        slide.shapes._spTree.insert_element_before(shape_el, "p:extLst")
        return

    if isinstance(shape_xml, dict) and shape_xml.get("tag"):
        try:
            shape_el = _xml_spec_to_oxml(shape_xml)
            if shape_el is None:
                return
        except Exception:
            return
    elif isinstance(shape_xml, dict):
        fmt = shape_xml.get("format")
        stype = shape_xml.get("type")
        if fmt != "shape_spec_v1" or not isinstance(stype, str):
            return

        def _apply_color_spec(parent_el, color_spec: Any) -> None:
            if not isinstance(color_spec, dict) or not color_spec.get("type"):
                return
            for child in list(parent_el):
                if child.tag == qn("a:solidFill"):
                    parent_el.remove(child)
            solid = OxmlElement("a:solidFill")
            if color_spec.get("type") == "srgb" and color_spec.get("rgb"):
                c = OxmlElement("a:srgbClr")
                c.set("val", str(color_spec.get("rgb")))
            elif color_spec.get("type") == "scheme" and color_spec.get("scheme"):
                c = OxmlElement("a:schemeClr")
                c.set("val", str(color_spec.get("scheme")))
            else:
                return
            for mod in color_spec.get("mods") or []:
                if not isinstance(mod, dict) or not mod.get("op"):
                    continue
                op = str(mod.get("op"))
                val = mod.get("val")
                el = OxmlElement(f"a:{op}")
                el.set("val", str(val))
                c.append(el)
            solid.append(c)
            parent_el.append(solid)

        if stype == "line":
            p1 = shape_xml.get("p1") or {}
            p2 = shape_xml.get("p2") or {}
            try:
                x1r, y1r = float(p1.get("x")), float(p1.get("y"))
                x2r, y2r = float(p2.get("x")), float(p2.get("y"))
            except Exception:
                return
            # Lines may have a 0-width/0-height bbox; rely on endpoints only.
            slide_w, slide_h = _slide_size_pt_from_slide(slide)
            x1 = Pt(x1r * slide_w)
            y1 = Pt(y1r * slide_h)
            x2 = Pt(x2r * slide_w)
            y2 = Pt(y2r * slide_h)
            ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
            line_style = shape_xml.get("line") or {}
            if isinstance(line_style, dict):
                ln_el = ln._element.spPr.get_or_add_ln()
                if line_style.get("width_pt") is not None:
                    try:
                        ln_el.set("w", str(int(round(float(line_style["width_pt"]) * EMU_PER_PT))))
                    except Exception:
                        pass
                color = line_style.get("color")
                _apply_color_spec(ln_el, color)
                # Dash style if available.
                if line_style.get("dash"):
                    dash_map = {
                        "solid": MSO_LINE_DASH_STYLE.SOLID,
                        "dash": MSO_LINE_DASH_STYLE.DASH,
                        "dashDot": MSO_LINE_DASH_STYLE.DASH_DOT,
                        "dashDotDot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
                        "lgDash": MSO_LINE_DASH_STYLE.LONG_DASH,
                        "lgDashDot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
                        "roundDot": MSO_LINE_DASH_STYLE.ROUND_DOT,
                        "sqDot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
                    }
                    try:
                        ds = dash_map.get(str(line_style["dash"]))
                        if ds is not None:
                            ln.line.dash_style = ds
                    except Exception:
                        pass

                # Arrowheads are not exposed as enums in this python-pptx build; inject OOXML.
                def _apply_end(tag: str, spec: Any) -> None:
                    if not isinstance(spec, dict) or not spec.get("type"):
                        return
                    ln_el = ln._element.spPr.get_or_add_ln()
                    # remove existing
                    for child in list(ln_el):
                        if child.tag == qn(f"a:{tag}"):
                            ln_el.remove(child)
                    end_el = OxmlElement(f"a:{tag}")
                    end_el.set("type", str(spec.get("type")))
                    if spec.get("w"):
                        end_el.set("w", str(spec.get("w")))
                    if spec.get("len"):
                        end_el.set("len", str(spec.get("len")))
                    ln_el.append(end_el)

                _apply_end("headEnd", line_style.get("head_end"))
                _apply_end("tailEnd", line_style.get("tail_end"))
            # #region debug-point B:line-rebuild
            try:
                _dbg_report_method_ppt_style(
                    "B",
                    "combine.py:_add_graph_layer:line",
                    "[DEBUG] rebuild line",
                    {
                        "shape_name": layer.get("shape_name"),
                        "p1": p1,
                        "p2": p2,
                        "line_style": line_style,
                        "ln_xml": ET.tostring(ln._element.spPr.get_or_add_ln(), encoding="unicode"),
                    },
                )
            except Exception:
                pass
            # #endregion
            return

        if stype != "geo":
            return

        base_box = layer.get("box")
        box = _shift_box_pt(base_box, shift_x, shift_y, _slide_size_pt_from_slide(slide))
        if box is not None:
            box = _clamp_box_to_slide(box, slide)
            if box is None:
                return
        if box is None:
            return
        left, top, width, height = box
        off_x = int(round(float(left.pt) * EMU_PER_PT))
        off_y = int(round(float(top.pt) * EMU_PER_PT))
        ext_cx = int(round(float(width.pt) * EMU_PER_PT))
        ext_cy = int(round(float(height.pt) * EMU_PER_PT))

        # Build a minimal p:sp with custGeom + solidFill + optional line style.
        geom = shape_xml.get("geom") or {}
        if not isinstance(geom, dict) or "type" not in geom:
            return

        def _is_rect_cust_geom(g: Dict[str, Any]) -> bool:
            if g.get("type") != "custGeom":
                return False
            paths_local = g.get("paths")
            if not isinstance(paths_local, list) or len(paths_local) != 1:
                return False
            path0 = paths_local[0]
            if not isinstance(path0, dict):
                return False
            cmds = path0.get("commands")
            if not isinstance(cmds, list) or len(cmds) < 5:
                return False
            pts = []
            for cmd in cmds:
                if not isinstance(cmd, dict):
                    return False
                op = cmd.get("op")
                if op == "close":
                    continue
                if op not in {"moveTo", "lnTo"}:
                    return False
                p = (cmd.get("pts") or [])
                if len(p) != 1:
                    return False
                try:
                    x = round(float(p[0].get("x")), 6)
                    y = round(float(p[0].get("y")), 6)
                except Exception:
                    return False
                pts.append((x, y))
            uniq = set(pts)
            corners = {(0.0, 0.0), (1.0, 0.0), (1.0, 1.0), (0.0, 1.0)}
            return corners.issubset(uniq)

        def _apply_shape_fill_and_line(sppr, fill_spec: Any, line_style: Any) -> None:
            try:
                # Normalize fill to explicit noFill when absent so PowerPoint does not
                # apply the theme's default blue shape fill.
                for child in list(sppr):
                    if child.tag in (qn("a:solidFill"), qn("a:noFill")):
                        sppr.remove(child)
                if isinstance(fill_spec, dict) and fill_spec.get("type"):
                    _apply_color_spec(sppr, fill_spec)
                else:
                    sppr.append(OxmlElement("a:noFill"))
            except Exception:
                pass

            try:
                ln_el = sppr.get_or_add_ln()
                for child in list(ln_el):
                    if child.tag in (qn("a:solidFill"), qn("a:noFill")):
                        ln_el.remove(child)
                if isinstance(line_style, dict):
                    color = line_style.get("color")
                    if isinstance(color, dict) and color.get("type"):
                        _apply_color_spec(ln_el, color)
                    else:
                        ln_el.append(OxmlElement("a:noFill"))
                else:
                    ln_el.append(OxmlElement("a:noFill"))
            except Exception:
                pass

        # Fast path: preset geometry like rect/roundRect -> add_shape()
        prst = str(geom.get("prst") or "")
        prst_map = {
            "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            "roundRect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            "roundrect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
            "leftBrace": MSO_AUTO_SHAPE_TYPE.LEFT_BRACE,
            "leftbrace": MSO_AUTO_SHAPE_TYPE.LEFT_BRACE,
            "rightBrace": MSO_AUTO_SHAPE_TYPE.RIGHT_BRACE,
            "rightbrace": MSO_AUTO_SHAPE_TYPE.RIGHT_BRACE,
            "leftBracket": MSO_AUTO_SHAPE_TYPE.LEFT_BRACKET,
            "leftbracket": MSO_AUTO_SHAPE_TYPE.LEFT_BRACKET,
            "rightBracket": MSO_AUTO_SHAPE_TYPE.RIGHT_BRACKET,
            "rightbracket": MSO_AUTO_SHAPE_TYPE.RIGHT_BRACKET,
            "bracePair": MSO_AUTO_SHAPE_TYPE.DOUBLE_BRACE,
            "bracepair": MSO_AUTO_SHAPE_TYPE.DOUBLE_BRACE,
            "bracketPair": MSO_AUTO_SHAPE_TYPE.DOUBLE_BRACKET,
            "bracketpair": MSO_AUTO_SHAPE_TYPE.DOUBLE_BRACKET,
        }
        ast = prst_map.get(prst)
        if ast is not None:
            rot = float(layer.get("rotation_deg") or 0.0)
            if rot:
                cx = float(left.pt) + float(width.pt) / 2.0
                cy = float(top.pt) + float(height.pt) / 2.0
                r = rot % 360.0
                if abs(r - 90.0) < 1e-3 or abs(r - 270.0) < 1e-3:
                    w0, h0 = float(height.pt), float(width.pt)
                else:
                    w0, h0 = float(width.pt), float(height.pt)
                left0 = Pt(cx - w0 / 2.0)
                top0 = Pt(cy - h0 / 2.0)
                shp = slide.shapes.add_shape(ast, left0, top0, Pt(w0), Pt(h0))
                try:
                    shp.rotation = rot
                except Exception:
                    pass
            else:
                shp = slide.shapes.add_shape(ast, left, top, width, height)
            _apply_shape_flip(shp, layer)
            fill = shape_xml.get("fill")
            line_style = shape_xml.get("line") if isinstance(shape_xml.get("line"), dict) else None
            _apply_shape_fill_and_line(shp._element.spPr, fill, line_style)
            if prst.lower() in {"leftbrace", "rightbrace", "leftbracket", "rightbracket", "bracepair", "bracketpair"} and not fill:
                try:
                    sppr = shp._element.spPr
                    for child in list(sppr):
                        if child.tag in (qn("a:solidFill"), qn("a:noFill")):
                            sppr.remove(child)
                    sppr.append(OxmlElement("a:noFill"))
                except Exception:
                    pass
            if isinstance(line_style, dict):
                if line_style.get("width_pt") is not None:
                    try:
                        shp.line.width = Pt(float(line_style["width_pt"]))
                    except Exception:
                        pass
                if line_style.get("dash"):
                    dash_map = {
                        "solid": MSO_LINE_DASH_STYLE.SOLID,
                        "dash": MSO_LINE_DASH_STYLE.DASH,
                        "dashDot": MSO_LINE_DASH_STYLE.DASH_DOT,
                        "dashDotDot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
                        "lgDash": MSO_LINE_DASH_STYLE.LONG_DASH,
                        "lgDashDot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
                        "roundDot": MSO_LINE_DASH_STYLE.ROUND_DOT,
                        "sqDot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
                    }
                    try:
                        ds = dash_map.get(str(line_style["dash"]))
                        if ds is not None:
                            shp.line.dash_style = ds
                    except Exception:
                        pass
            text_content = layer.get("text_content")
            if isinstance(text_content, dict):
                try:
                    _fill_text_frame(shp.text_frame, text_content, _slide_size_pt_from_slide(slide))
                except Exception:
                    pass
            try:
                _dbg_report_method_ppt_style(
                    "A",
                    "combine.py:_add_graph_layer:prstGeom",
                    "[DEBUG] rebuild preset geom",
                    {
                        "shape_name": layer.get("shape_name"),
                        "geom": geom,
                        "fill": fill,
                        "line_style": line_style,
                        "sppr_xml": ET.tostring(shp._element.spPr, encoding="unicode"),
                    },
                )
            except Exception:
                pass
            return

        # Another fast path: many background panels are exported as custGeom
        # rectangles, but rebuilding them as true RECTANGLE shapes is more robust
        # in PowerPoint than custom geometry.
        if _is_rect_cust_geom(geom):
            shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            _apply_shape_flip(shp, layer)
            fill = shape_xml.get("fill")
            line_style = shape_xml.get("line") if isinstance(shape_xml.get("line"), dict) else None
            _apply_shape_fill_and_line(shp._element.spPr, fill, line_style)
            if isinstance(line_style, dict):
                if line_style.get("width_pt") is not None:
                    try:
                        shp.line.width = Pt(float(line_style["width_pt"]))
                    except Exception:
                        pass
                if line_style.get("dash"):
                    dash_map = {
                        "solid": MSO_LINE_DASH_STYLE.SOLID,
                        "dash": MSO_LINE_DASH_STYLE.DASH,
                        "dashDot": MSO_LINE_DASH_STYLE.DASH_DOT,
                        "dashDotDot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
                        "lgDash": MSO_LINE_DASH_STYLE.LONG_DASH,
                        "lgDashDot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
                        "roundDot": MSO_LINE_DASH_STYLE.ROUND_DOT,
                        "sqDot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
                    }
                    try:
                        ds = dash_map.get(str(line_style["dash"]))
                        if ds is not None:
                            shp.line.dash_style = ds
                    except Exception:
                        pass
            text_content = layer.get("text_content")
            if isinstance(text_content, dict):
                try:
                    _fill_text_frame(shp.text_frame, text_content, _slide_size_pt_from_slide(slide))
                except Exception:
                    pass
            return

        if geom.get("type") != "custGeom":
            return
        paths = geom.get("paths")
        if not isinstance(paths, list) or not paths:
            return

        def _pt_xy(rpt: Dict[str, Any], path_w: int, path_h: int) -> tuple[int, int]:
            x = float(rpt.get("x") or 0.0)
            y = float(rpt.get("y") or 0.0)
            return int(round(x * path_w)), int(round(y * path_h))

        rot_attr = ""
        try:
            rot = float(layer.get("rotation_deg") or 0.0)
            if rot:
                rot_attr = f' rot="{int(round(rot * 60000))}"'
        except Exception:
            rot_attr = ""
        flip_h = ' flipH="true"' if layer.get("flip_h") else ""
        flip_v = ' flipV="true"' if layer.get("flip_v") else ""

        fill_rgb = None
        fill = shape_xml.get("fill")
        if isinstance(fill, dict) and fill.get("type") == "srgb" and fill.get("rgb"):
            fill_rgb = str(fill["rgb"])
        line_style = shape_xml.get("line") if isinstance(shape_xml.get("line"), dict) else None

        # Build path xml
        path_chunks: List[str] = []
        for p in paths:
            if not isinstance(p, dict):
                continue
            try:
                path_w = int(p.get("w") or 1000000)
                path_h = int(p.get("h") or 1000000)
            except Exception:
                path_w, path_h = 1000000, 1000000
            if path_w <= 0 or path_h <= 0:
                path_w, path_h = 1000000, 1000000
            cmds = p.get("commands")
            if not isinstance(cmds, list):
                continue
            cmd_chunks: List[str] = []
            for cmd in cmds:
                if not isinstance(cmd, dict):
                    continue
                op = cmd.get("op")
                pts = cmd.get("pts") or []
                if op in {"moveTo", "lnTo"} and isinstance(pts, list) and len(pts) == 1:
                    x, y = _pt_xy(pts[0], path_w, path_h)
                    cmd_chunks.append(f'<a:{op}><a:pt x="{x}" y="{y}"/></a:{op}>')
                elif op == "cubicBezTo" and isinstance(pts, list) and len(pts) == 3:
                    pts_xml = "".join(
                        f'<a:pt x="{_pt_xy(pt, path_w, path_h)[0]}" y="{_pt_xy(pt, path_w, path_h)[1]}"/>'
                        for pt in pts
                    )
                    cmd_chunks.append(f"<a:cubicBezTo>{pts_xml}</a:cubicBezTo>")
                elif op == "close":
                    cmd_chunks.append("<a:close/>")
                else:
                    return
            path_chunks.append(
                f'<a:path w="{path_w}" h="{path_h}">' + "".join(cmd_chunks) + "</a:path>"
            )
        if not path_chunks:
            return

        fill_xml = ""
        def _color_xml(spec: Any) -> str:
            if not isinstance(spec, dict) or not spec.get("type"):
                return ""
            mods = ""
            for mod in spec.get("mods") or []:
                if isinstance(mod, dict) and mod.get("op") and mod.get("val") is not None:
                    mods += f'<a:{str(mod["op"])} val="{str(mod["val"])}"/>'
            if spec.get("type") == "srgb" and spec.get("rgb"):
                return f'<a:solidFill><a:srgbClr val="{str(spec.get("rgb"))}">{mods}</a:srgbClr></a:solidFill>'
            if spec.get("type") == "scheme" and spec.get("scheme"):
                return f'<a:solidFill><a:schemeClr val="{str(spec.get("scheme"))}">{mods}</a:schemeClr></a:solidFill>'
            return ""
        if isinstance(fill, dict) and fill.get("type"):
            fill_xml = _color_xml(fill)
        else:
            fill_xml = "<a:noFill/>"
        ln_xml = ""
        if isinstance(line_style, dict):
            parts: List[str] = []
            w_pt = line_style.get("width_pt")
            if w_pt is not None:
                try:
                    parts.append(f' w="{int(round(float(w_pt) * EMU_PER_PT))}"')
                except Exception:
                    pass
            ln_color = line_style.get("color") if isinstance(line_style.get("color"), dict) else None
            ln_fill = ""
            if isinstance(ln_color, dict) and ln_color.get("type"):
                ln_fill = _color_xml(ln_color)
            else:
                ln_fill = "<a:noFill/>"
            ln_xml = f"<a:ln{''.join(parts)}>{ln_fill}</a:ln>" if (parts or ln_fill) else ""
        else:
            ln_xml = "<a:ln><a:noFill/></a:ln>"

        xml = (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            "<p:nvSpPr><p:cNvPr id=\"1\" name=\"shape\"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>"
            "<p:spPr>"
            f'<a:xfrm{flip_h}{flip_v}{rot_attr}><a:off x="{off_x}" y="{off_y}"/><a:ext cx="{ext_cx}" cy="{ext_cy}"/></a:xfrm>'
            "<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/><a:rect l=\"l\" t=\"t\" r=\"r\" b=\"b\"/>"
            "<a:pathLst>"
            + "".join(path_chunks)
            + "</a:pathLst></a:custGeom>"
            + fill_xml
            + ln_xml
            + "</p:spPr></p:sp>"
        )
        try:
            shape_el = parse_xml(xml.encode("utf-8"))
        except Exception:
            return
    else:
        try:
            shape_el = parse_xml(str(shape_xml).encode("utf-8"))
        except Exception:
            return

    c_nv_pr = shape_el.find(".//p:cNvPr", namespaces=XML_NS)
    if c_nv_pr is not None:
        c_nv_pr.set("id", str(_next_shape_id(slide)))
        if layer.get("shape_name"):
            c_nv_pr.set("name", str(layer.get("shape_name")))

    xfrm = shape_el.find(".//a:xfrm", namespaces=XML_NS)
    if xfrm is not None and box is not None:
        left, top, width, height = box
        off = xfrm.find("a:off", namespaces=XML_NS)
        ext = xfrm.find("a:ext", namespaces=XML_NS)
        if off is None:
            off = OxmlElement("a:off")
            xfrm.insert(0, off)
        off.set("x", str(int(round(float(left.pt) * EMU_PER_PT))))
        off.set("y", str(int(round(float(top.pt) * EMU_PER_PT))))
        if ext is None:
            ext = OxmlElement("a:ext")
            xfrm.insert(1, ext)
        ext.set("cx", str(int(round(float(width.pt) * EMU_PER_PT))))
        ext.set("cy", str(int(round(float(height.pt) * EMU_PER_PT))))

        # For raw XML shapes, keep original rot/flip attrs untouched; only off/ext is box-driven.
        if not is_raw_xml:
            try:
                rot = float(layer.get("rotation_deg") or 0.0)
                if rot:
                    xfrm.set("rot", str(int(round(rot * 60000))))
                elif "rot" in xfrm.attrib:
                    del xfrm.attrib["rot"]
            except Exception:
                pass
            if layer.get("flip_h"):
                xfrm.set("flipH", "true")
            elif "flipH" in xfrm.attrib:
                del xfrm.attrib["flipH"]
            if layer.get("flip_v"):
                xfrm.set("flipV", "true")
            elif "flipV" in xfrm.attrib:
                del xfrm.attrib["flipV"]

    try:
        slide.shapes._spTree.insert_element_before(shape_el, "p:extLst")
    except Exception:
        slide.shapes._spTree.append(shape_el)


def _add_image_layer(slide, layer: Dict[str, Any], assets_dir: Path, shift_x: float, shift_y: float) -> None:
    saved_path = layer.get("saved_path")
    if not saved_path:
        return
    img_path = Path(saved_path)
    if not img_path.is_absolute():
        img_path = (assets_dir / img_path).resolve()
    if not img_path.exists():
        # Try by basename under assets_dir (handles moved output folders)
        img_path2 = (assets_dir / Path(saved_path).name).resolve()
        if img_path2.exists():
            img_path = img_path2
        else:
            return

    base_box = layer.get("box")
    box = _shift_box_pt(base_box, shift_x, shift_y, _slide_size_pt_from_slide(slide))
    if box:
        left, top, width, height = box
    else:
        # Background with no box: fill slide
        left, top = Pt(0), Pt(0)
        slide_w, slide_h = _slide_size_pt_from_slide(slide)
        width = Pt(slide_w)
        height = Pt(slide_h)

    slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)


def rebuild_ppt(
    layers_json: Path,
    assets_dir: Path,
    out_pptx: Path,
    normalize: bool,
) -> None:
    raw_payload = _decode_scaled_protocol_payload(json.loads(layers_json.read_text(encoding="utf-8")))
    layers: List[Dict[str, Any]] = _normalize_layers_payload(raw_payload)
    canvas_size = _read_canvas_size_pt(layers)

    # Read slide_canvas first and establish the presentation page size up front.
    if canvas_size is not None:
        initial_slide_w, initial_slide_h = canvas_size
    else:
        initial_slide_w, initial_slide_h = 720.0, 540.0

    # Group by slide, keep top->bottom order by layer_index_in_slide if present.
    by_slide: Dict[int, List[Dict[str, Any]]] = defaultdict(list)
    for layer in layers:
        by_slide[int(layer.get("slide", 1))].append(layer)
    for s in by_slide:
        by_slide[s].sort(key=lambda x: int(x.get("layer_index_in_slide", x.get("layer_index_global", 0))))

    # If slide_canvas exists, treat it as the authoritative page size and preserve absolute
    # coordinates exactly, including negative/off-canvas positions.
    if canvas_size is not None:
        slide_w, slide_h = initial_slide_w, initial_slide_h
        shift_x, shift_y = 0, 0
    elif normalize:
        shift_x, shift_y, slide_w, slide_h = _compute_canvas_shift(layers, canvas_size)
    else:
        slide_w, slide_h = initial_slide_w, initial_slide_h
        shift_x, shift_y = 0, 0

    prs = Presentation()
    prs.slide_width = Pt(initial_slide_w)
    prs.slide_height = Pt(initial_slide_h)
    prs.slide_width = Pt(slide_w)
    prs.slide_height = Pt(slide_h)

    blank_layout = prs.slide_layouts[6]
    max_slide = max(by_slide.keys()) if by_slide else 1
    for si in range(1, max_slide + 1):
        slide = prs.slides.add_slide(blank_layout)
        slide_layers = by_slide.get(si, [])
        _apply_slide_background_fill(slide, _extract_slide_background_fill(slide_layers))

        # Add bottom->top so visual order matches JSON top->bottom.
        for layer in reversed(slide_layers):
            kind = layer.get("kind")
            if kind == "slide_canvas":
                continue
            if kind == "text":
                _add_text_layer(slide, layer, shift_x, shift_y)
            elif kind == "ppt_graph_table":
                _add_table_layer(slide, layer, shift_x, shift_y)
            elif kind in {"ppt_graph_line", "ppt_graph_geo"}:
                _add_graph_layer(slide, layer, shift_x, shift_y)
            elif kind in {"image_png_rgba", "image_raw", "svg_image_png_rgba"}:
                _add_image_layer(slide, layer, assets_dir, shift_x, shift_y)
            else:
                continue

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))


def main() -> None:
    ap = argparse.ArgumentParser(description="Rebuild a PPTX from layers.json + assets/")
    ap.add_argument("layers_json", help="Path to layers.json")
    ap.add_argument("--assets", default=None, help="Assets folder (default: <layers_json_dir>/assets)")
    ap.add_argument("-o", "--out", default=None, help="Output pptx path (default: <layers_json_dir>/rebuilt.pptx)")
    ap.add_argument("--no-normalize", action="store_true", help="Do not shift negative coords into positive space")
    ap.add_argument("--rel-scale", type=float, default=1.0, help="Scale factor applied to rel-only font sizes (default: 1.0)")
    ap.add_argument("--rel-base", choices=["height", "min", "geom"], default="height", help="Baseline for rel-only font sizes (default: height)")
    args = ap.parse_args()

    global _REL_SCALE
    try:
        _REL_SCALE = float(args.rel_scale)
    except Exception:
        _REL_SCALE = 1.0
    global _REL_BASE
    _REL_BASE = str(args.rel_base or "height")

    layers_json = Path(args.layers_json).expanduser().resolve()
    assets_dir = _infer_assets_dir(layers_json, args.assets)
    out_pptx = (
        Path(args.out).expanduser().resolve()
        if args.out
        else (layers_json.parent / "rebuilt.pptx").resolve()
    )

    rebuild_ppt(layers_json, assets_dir, out_pptx, normalize=(not args.no_normalize))
    injected = _embed_extracted_fonts_into_pptx(layers_json, out_pptx)
    if injected:
        print(f"Embedded fonts from extracted assets: {injected} font part(s)")
    print(f"Saved: {out_pptx}")


if __name__ == "__main__":
    main()

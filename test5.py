#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Export a PPTX into the current layer JSON protocol.

- Text layers keep paragraph/run style data.
- Image layers are baked to final RGBA PNG files plus final ratio boxes.
- Shape layers keep enough geometry/XML to rebuild the slide.

Usage:
  pip install python-pptx pillow
  python test.py /path/to/input.pptx -o /path/to/out_dir
"""

import argparse
import colorsys
import io
import json
import math
import os
import re
import shutil
import subprocess
import time
import zipfile
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from PIL import Image, ImageDraw, ImageOps


EMU_PER_PT = 12700  # PowerPoint uses EMU; 1 pt = 12700 EMU
RASTER_TARGET_SIDE = 4096
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "asvg": "http://schemas.microsoft.com/office/drawing/2016/SVG/main",
}

TONGJI_TEMPLATE_DIR = Path("/Users/bytedance/Downloads/tongjiall")
_KIND_TEMPLATE_CACHE: Dict[str, Optional[Dict[str, Any]]] = {}
IMAGE_LIKE_EXPORT_KINDS = {"image_png_rgba", "svg_image_png_rgba", "image_raw"}
_CURRENT_PPT_THEME_COLORS: Dict[str, str] = {}


def _null_object_from_template(template: Dict[str, Any]) -> Dict[str, Any]:
    out: Dict[str, Any] = {}
    for key, value in template.items():
        if isinstance(value, dict):
            out[key] = _null_object_from_template(value)
        else:
            out[key] = None
    return out


def _compact_nested_template_value(value: Any) -> Any:
    if isinstance(value, dict):
        compacted: Dict[str, Any] = {}
        for key, child in value.items():
            child_value = _compact_nested_template_value(child)
            if child_value is None:
                continue
            if isinstance(child_value, dict) and not child_value:
                continue
            if isinstance(child_value, list) and not child_value:
                continue
            compacted[key] = child_value
        return compacted or None
    if isinstance(value, list):
        compacted_list = []
        for item in value:
            item_value = _compact_nested_template_value(item)
            if item_value is None:
                continue
            if isinstance(item_value, dict) and not item_value:
                continue
            if isinstance(item_value, list) and not item_value:
                continue
            compacted_list.append(item_value)
        return compacted_list or None
    return value


def _theme_scheme_alias(name: Any) -> str:
    key = str(name or "").strip()
    alias = {
        "bg1": "lt1",
        "tx1": "dk1",
        "bg2": "lt2",
        "tx2": "dk2",
    }
    return alias.get(key, key)


def _extract_theme_colors_from_pptx(pptx_path: Path) -> Dict[str, str]:
    out: Dict[str, str] = {}
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            theme_names = [name for name in zf.namelist() if name.startswith("ppt/theme/") and name.endswith(".xml")]
            for theme_name in theme_names:
                try:
                    root = ET.fromstring(zf.read(theme_name))
                except Exception:
                    continue
                clr_scheme = root.find(".//a:clrScheme", NS)
                if clr_scheme is None:
                    continue
                for child in list(clr_scheme):
                    tag = child.tag.rsplit("}", 1)[-1]
                    rgb = None
                    srgb = child.find("a:srgbClr", NS)
                    if srgb is not None and srgb.get("val"):
                        rgb = str(srgb.get("val")).upper()
                    else:
                        sys_clr = child.find("a:sysClr", NS)
                        if sys_clr is not None:
                            rgb = str(sys_clr.get("lastClr") or sys_clr.get("val") or "").upper()
                    if rgb and len(rgb) == 6:
                        out[_theme_scheme_alias(tag)] = rgb
                if out:
                    break
    except Exception:
        return {}
    return out


def _load_kind_template(kind: Any) -> Optional[Dict[str, Any]]:
    kind_name = str(kind or "").strip()
    if not kind_name:
        return None
    if kind_name in _KIND_TEMPLATE_CACHE:
        return _KIND_TEMPLATE_CACHE[kind_name]
    path = TONGJI_TEMPLATE_DIR / f"{kind_name}.json"
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
        templates = payload.get("templates")
        template = templates[0] if isinstance(templates, list) and isinstance(templates[0], dict) else None
    except Exception:
        template = None
    _KIND_TEMPLATE_CACHE[kind_name] = template
    return template


def _first_object_template(items: Any) -> Optional[Dict[str, Any]]:
    if not isinstance(items, list):
        return None
    for item in items:
        if isinstance(item, dict) and not set(item.keys()).issubset({"min", "max"}):
            return item
        if isinstance(item, list):
            nested = _first_object_template(item)
            if nested is not None:
                return nested
    return None


def _fill_value_from_template(template_value: Any, value: Any) -> Any:
    if isinstance(template_value, dict):
        source = value if isinstance(value, dict) else {}
        return {key: _fill_value_from_template(child_template, source.get(key)) for key, child_template in template_value.items()}

    item_template = _first_object_template(template_value)
    if item_template is not None:
        if isinstance(value, list):
            filled_items = [
                _fill_value_from_template(item_template, item if isinstance(item, dict) else {})
                for item in value
            ]
            return filled_items or [_fill_value_from_template(item_template, {})]
        if isinstance(value, dict):
            return _fill_value_from_template(item_template, value)
        return [_fill_value_from_template(item_template, {})]

    return value if value is not None else None


def _apply_kind_template(layer: Dict[str, Any]) -> Dict[str, Any]:
    template = _load_kind_template(layer.get("kind"))
    if not isinstance(template, dict):
        return layer
    out = {key: _fill_value_from_template(template_value, layer.get(key)) for key, template_value in template.items()}
    if layer.get("saved_path"):
        for extra_key in ("structure_info", "shape_xml", "fill", "line", "rotation_deg", "flip_h", "flip_v"):
            if extra_key in layer:
                out[extra_key] = layer.get(extra_key)
    return out


def _kind_template_has_key(kind: Any, key: str) -> bool:
    template = _load_kind_template(kind)
    return isinstance(template, dict) and key in template


def _minimal_text_schema_template() -> Dict[str, Any]:
    return {
        "shape_name": None,
        "kind": "text",
        "rotation_deg": None,
        "flip_h": None,
        "flip_v": None,
        "box": {
            "left": None,
            "top": None,
            "width": None,
            "height": None,
        },
        "paragraphs": [
            {
                "runs": [
                    {
                        "text": None,
                        "font_name": None,
                        "font_size_rel": None,
                        "bold": False,
                        "italic": False,
                        "underline": False,
                        "strike": False,
                        "color": {
                            "type": None,
                            "rgb": None,
                            "alpha_val": 100000,
                        },
                        "char_spacing_rel": 0,
                    }
                ],
                "ppr_attrs": {
                    "algn": None,
                    "marL_rel": 0,
                    "indent_rel": 0,
                    "lvl": 0,
                },
                "line_spacing": {
                    "rel": 0,
                },
                "space_before": {
                    "rel": 0,
                },
            }
        ],
        "body_pr_anchor": None,
    }


def _normalize_box_template(box: Any) -> Optional[Dict[str, Any]]:
    template = {
        "left": None,
        "top": None,
        "width": None,
        "height": None,
    }
    if not isinstance(box, dict):
        return dict(template)
    return {
        "left": box.get("left"),
        "top": box.get("top"),
        "width": box.get("width"),
        "height": box.get("height"),
    }


def _alpha_val_from_text_color(color: Any) -> int:
    if not isinstance(color, dict):
        return 100000
    alpha_val = color.get("alpha_val")
    if alpha_val is not None:
        try:
            return int(alpha_val)
        except Exception:
            return 100000
    for item in color.get("mods") or []:
        if not isinstance(item, dict):
            continue
        if str(item.get("op") or "").lower() != "alpha":
            continue
        try:
            return int(item.get("val"))
        except Exception:
            return 100000
    return 100000


def _color_mods_without_alpha(color: Any) -> Optional[List[Dict[str, Any]]]:
    if not isinstance(color, dict):
        return None
    mods: List[Dict[str, Any]] = []
    for item in color.get("mods") or []:
        if not isinstance(item, dict):
            continue
        op = item.get("op")
        val = item.get("val")
        if not op or val is None:
            continue
        if str(op).lower() == "alpha":
            continue
        mods.append({"op": str(op), "val": val})
    return mods or None


def _normalize_color_template(color: Any) -> Optional[Dict[str, Any]]:
    template = {
        "type": color.get("type"),
        "rgb": color.get("rgb"),
        "alpha_val": _alpha_val_from_text_color(color),
    } if isinstance(color, dict) else {
        "type": None,
        "rgb": None,
        "alpha_val": 100000,
    }
    mods = _color_mods_without_alpha(color)
    if mods:
        template["mods"] = mods
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_font_template(font: Any) -> Optional[Dict[str, Any]]:
    template = {"typeface": (font.get("typeface") if isinstance(font, dict) else None)}
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_bullet_template(bullet: Any) -> Optional[Dict[str, Any]]:
    template = {
        "type": bullet.get("type"),
        "char": bullet.get("char"),
        "font": _normalize_font_template(bullet.get("font")),
        "startAt": bullet.get("startAt"),
    } if isinstance(bullet, dict) else {
        "type": None,
        "char": None,
        "font": _null_object_from_template({"typeface": None}),
        "startAt": None,
    }
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _inline_simple_bullet_into_runs(runs: List[Dict[str, Any]], bullet: Any) -> Optional[Dict[str, Any]]:
    if not (isinstance(bullet, dict) and isinstance(runs, list) and runs):
        return bullet if isinstance(bullet, dict) else None
    if str(bullet.get("type") or "").lower() != "char":
        return bullet
    bullet_char = _normalize_bullet_char(bullet.get("char"), bullet.get("font"))
    if not bullet_char:
        return None
    first_run = runs[0]
    if not isinstance(first_run, dict):
        return None
    first_text = str(first_run.get("text") or "")
    if first_text.startswith(bullet_char):
        return None
    spacer = "" if not first_text or first_text[:1].isspace() else " "
    first_run["text"] = f"{bullet_char}{spacer}{first_text}"
    return None


def _normalize_bullet_char(raw_char: Any, font: Any = None) -> str:
    bullet_char = str(raw_char or "")
    if not bullet_char:
        return ""
    typeface = ""
    if isinstance(font, dict):
        typeface = str(font.get("typeface") or "").strip().lower()
    if typeface == "wingdings":
        return {
            "l": "•",
            "n": "■",
            "u": "◆",
        }.get(bullet_char, bullet_char)
    return bullet_char


def _style_flag(value: Any, *, strike: bool = False) -> bool:
    if isinstance(value, bool):
        return value
    if value is None:
        return False
    token = str(value).strip().lower()
    if not token:
        return False
    if strike:
        return token not in {"nostrike", "false", "0", "none"}
    if token in {"none", "false", "0"}:
        return False
    return True


def _number_to_alpha(value: int, uppercase: bool = False) -> str:
    if value <= 0:
        value = 1
    chars: List[str] = []
    n = value
    while n > 0:
        n -= 1
        chars.append(chr((n % 26) + (65 if uppercase else 97)))
        n //= 26
    return "".join(reversed(chars))


def _number_to_roman(value: int, uppercase: bool = False) -> str:
    if value <= 0:
        value = 1
    table = [
        (1000, "M"),
        (900, "CM"),
        (500, "D"),
        (400, "CD"),
        (100, "C"),
        (90, "XC"),
        (50, "L"),
        (40, "XL"),
        (10, "X"),
        (9, "IX"),
        (5, "V"),
        (4, "IV"),
        (1, "I"),
    ]
    out = []
    n = value
    for arabic, roman in table:
        while n >= arabic:
            out.append(roman)
            n -= arabic
    result = "".join(out)
    return result if uppercase else result.lower()


def _autonum_prefix_text(scheme: Any, ordinal: int) -> str:
    token = str(scheme or "arabicPeriod")
    if token == "alphaLcPeriod":
        return f"{_number_to_alpha(ordinal, uppercase=False)}."
    if token == "alphaUcPeriod":
        return f"{_number_to_alpha(ordinal, uppercase=True)}."
    if token == "alphaLcParenR":
        return f"{_number_to_alpha(ordinal, uppercase=False)})"
    if token == "alphaUcParenR":
        return f"{_number_to_alpha(ordinal, uppercase=True)})"
    if token == "romanLcPeriod":
        return f"{_number_to_roman(ordinal, uppercase=False)}."
    if token == "romanUcPeriod":
        return f"{_number_to_roman(ordinal, uppercase=True)}."
    if token == "romanLcParenR":
        return f"{_number_to_roman(ordinal, uppercase=False)})"
    if token == "romanUcParenR":
        return f"{_number_to_roman(ordinal, uppercase=True)})"
    if token == "arabicParenR":
        return f"{ordinal})"
    if token == "arabicPlain":
        return str(ordinal)
    return f"{ordinal}."


def _inline_bullet_into_runs(
    runs: List[Dict[str, Any]],
    bullet: Any,
    autonum_counters: Optional[Dict[Tuple[Any, ...], int]] = None,
    counter_key: Optional[Tuple[Any, ...]] = None,
) -> Optional[Dict[str, Any]]:
    if not (isinstance(bullet, dict) and isinstance(runs, list) and runs):
        return bullet if isinstance(bullet, dict) else None
    bullet_type = str(bullet.get("type") or "").lower()
    if bullet_type == "char":
        return _inline_simple_bullet_into_runs(runs, bullet)
    if bullet_type != "autonum":
        return bullet
    if autonum_counters is None or counter_key is None:
        return bullet
    current = autonum_counters.get(counter_key)
    if current is None:
        start_at = bullet.get("startAt")
        try:
            current = int(start_at) if start_at is not None else 1
        except Exception:
            current = 1
    prefix = _autonum_prefix_text(bullet.get("scheme"), current)
    autonum_counters[counter_key] = current + 1
    first_run = runs[0]
    if not isinstance(first_run, dict):
        return None
    first_text = str(first_run.get("text") or "")
    if first_text.startswith(prefix):
        return None
    spacer = "" if not first_text or first_text[:1].isspace() else " "
    first_run["text"] = f"{prefix}{spacer}{first_text}"
    return None


def _normalize_ppr_attrs_template(attrs: Any) -> Optional[Dict[str, Any]]:
    template = {
        "algn": attrs.get("algn"),
        "marL_rel": attrs.get("marL_rel", 0),
        "indent_rel": attrs.get("indent_rel", 0),
        "lvl": attrs.get("lvl", 0),
    } if isinstance(attrs, dict) else {
        "algn": None,
        "marL_rel": 0,
        "indent_rel": 0,
        "lvl": 0,
    }
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_spacing_template(spec: Any) -> Optional[Dict[str, Any]]:
    template = {} if isinstance(spec, dict) else {"rel": 0}
    if isinstance(spec, dict):
        mode = spec.get("mode")
        if mode not in {None, "points"}:
            template["mode"] = mode
        if spec.get("raw") is not None:
            template["raw"] = spec.get("raw")
        template["rel"] = spec.get("rel", 0)
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_min_text_spacing_template(spec: Any) -> Optional[Dict[str, Any]]:
    template = {
        "mode": spec.get("mode") if isinstance(spec, dict) else None,
        "rel": spec.get("rel") if isinstance(spec, dict) else None,
    }
    return _compact_nested_template_value(template)


def _normalize_end_para_rpr_template(spec: Any) -> Optional[Dict[str, Any]]:
    template = {
        "color": _normalize_color_template(spec.get("color")) if isinstance(spec, dict) else _null_object_from_template(
            {"type": None, "rgb": None, "alpha_val": 100000}
        )
    }
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_run_template(run: Any) -> Optional[Dict[str, Any]]:
    if not isinstance(run, dict):
        return None
    return _compact_nested_template_value(
        {
            "text": run.get("text"),
            "font_name": run.get("font_name"),
            "font_size_rel": run.get("font_size_rel"),
            "bold": bool(run.get("bold")),
            "italic": bool(run.get("italic")),
            "underline": bool(run.get("underline")),
            "strike": bool(run.get("strike")),
            "color": _normalize_color_template(run.get("color")),
            "char_spacing_rel": 0 if run.get("char_spacing_rel") is None else run.get("char_spacing_rel"),
        }
    )


def _normalize_paragraph_template(paragraph: Any) -> Optional[Dict[str, Any]]:
    if not isinstance(paragraph, dict):
        return None
    runs = paragraph.get("runs")
    return _compact_nested_template_value(
        {
            "runs": ([_normalize_run_template(item) for item in runs if isinstance(item, dict)] if isinstance(runs, list) else None),
            "ppr_attrs": _normalize_ppr_attrs_template(paragraph.get("ppr_attrs")),
            "line_spacing": _normalize_spacing_template(paragraph.get("line_spacing")),
            "space_before": _normalize_spacing_template(paragraph.get("space_before")),
        }
    )


def _normalize_body_pr_attrs_template(attrs: Any, *, include_textbox_insets: bool) -> Optional[Dict[str, Any]]:
    if isinstance(attrs, dict):
        out = {
            "anchor": attrs.get("anchor"),
            "rtlCol": attrs.get("rtlCol"),
        }
    else:
        out = {
            "anchor": None,
            "rtlCol": None,
        }
    if include_textbox_insets:
        out.update(
            {
                "tIns": attrs.get("tIns") if isinstance(attrs, dict) else None,
                "lIns": attrs.get("lIns") if isinstance(attrs, dict) else None,
                "bIns": attrs.get("bIns") if isinstance(attrs, dict) else None,
                "rIns": attrs.get("rIns") if isinstance(attrs, dict) else None,
                "wrap": attrs.get("wrap") if isinstance(attrs, dict) else None,
            }
        )
    return _compact_nested_template_value(out) or _null_object_from_template(out)


def _normalize_body_pr_template(body_pr: Any, *, include_textbox_insets: bool) -> Optional[Dict[str, Any]]:
    template = {
        "attrs": _normalize_body_pr_attrs_template(
            body_pr.get("attrs") if isinstance(body_pr, dict) else None,
            include_textbox_insets=include_textbox_insets,
        ),
        "autofit": body_pr.get("autofit") if isinstance(body_pr, dict) else None,
    }
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_body_pr_compact(body_pr: Any, *, include_textbox_insets: bool) -> Optional[Dict[str, Any]]:
    if not isinstance(body_pr, dict):
        return None
    attrs = body_pr.get("attrs")
    template = {
        "attrs": _compact_nested_template_value(
            {
                "anchor": attrs.get("anchor") if isinstance(attrs, dict) else None,
                "rtlCol": attrs.get("rtlCol") if isinstance(attrs, dict) else None,
                "tIns": attrs.get("tIns") if include_textbox_insets and isinstance(attrs, dict) else None,
                "lIns": attrs.get("lIns") if include_textbox_insets and isinstance(attrs, dict) else None,
                "bIns": attrs.get("bIns") if include_textbox_insets and isinstance(attrs, dict) else None,
                "rIns": attrs.get("rIns") if include_textbox_insets and isinstance(attrs, dict) else None,
                "wrap": attrs.get("wrap") if include_textbox_insets and isinstance(attrs, dict) else None,
            }
        ),
        "autofit": body_pr.get("autofit"),
    }
    return _compact_nested_template_value(template)


def _prune_text_run_defaults(run: Dict[str, Any]) -> Dict[str, Any]:
    color = run.get("color")
    if isinstance(color, dict):
        if not color:
            run.pop("color", None)
    for key in ("char_spacing_rel",):
        if run.get(key) is None:
            run.pop(key, None)
    return run


def _prune_text_paragraph_defaults(paragraph: Dict[str, Any]) -> Dict[str, Any]:
    ppr_attrs = paragraph.get("ppr_attrs")
    if isinstance(ppr_attrs, dict) and all(v is None for v in ppr_attrs.values()):
        paragraph.pop("ppr_attrs", None)
    for key in ("space_before", "space_after"):
        spacing = paragraph.get(key)
        if isinstance(spacing, dict) and all(v is None for v in spacing.values()):
            paragraph.pop(key, None)
    return paragraph


def _prune_text_layer_defaults(layer: Dict[str, Any]) -> Dict[str, Any]:
    paragraphs = layer.get("paragraphs")
    if isinstance(paragraphs, list):
        for paragraph in paragraphs:
            if not isinstance(paragraph, dict):
                continue
            runs = paragraph.get("runs")
            if isinstance(runs, list):
                for run in runs:
                    if isinstance(run, dict):
                        _prune_text_run_defaults(run)
            _prune_text_paragraph_defaults(paragraph)
    body_pr = layer.get("body_pr")
    if isinstance(body_pr, dict):
        attrs = body_pr.get("attrs")
        if isinstance(attrs, dict):
            if attrs.get("anchor") == "t":
                attrs.pop("anchor", None)
            rtl_col = attrs.get("rtlCol")
            if rtl_col in {None, 0, "0", False, "false"}:
                attrs.pop("rtlCol", None)
            if attrs.get("wrap") is None:
                attrs.pop("wrap", None)
            if not attrs:
                body_pr.pop("attrs", None)
        attrs = body_pr.get("attrs")
        if isinstance(attrs, dict):
            only_zero_insets = set(attrs.keys()).issubset({"tIns", "lIns", "bIns", "rIns"}) and all(
                attrs.get(key) in {None, 0, 0.0, "0"} for key in ("tIns", "lIns", "bIns", "rIns")
            )
            if only_zero_insets and body_pr.get("autofit") == "spAutoFit":
                layer.pop("body_pr", None)
                return layer
        if not body_pr:
            layer.pop("body_pr", None)
    return layer


def _normalize_line_end_template(end_spec: Any) -> Optional[Dict[str, Any]]:
    template = {
        "type": end_spec.get("type"),
        "w": end_spec.get("w"),
        "len": end_spec.get("len"),
    } if isinstance(end_spec, dict) else {
        "type": None,
        "w": None,
        "len": None,
    }
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_line_style_template(line: Any, *, include_arrow_ends: bool) -> Optional[Dict[str, Any]]:
    if isinstance(line, dict):
        out = {
            "width_pt": line.get("width_pt"),
            "color": _normalize_color_template(line.get("color")),
            "dash": line.get("dash"),
        }
    else:
        out = {
            "width_pt": None,
            "color": _null_object_from_template({"type": None, "rgb": None, "alpha_val": 100000}),
            "dash": None,
        }
    if include_arrow_ends:
        out["head_end"] = _normalize_line_end_template(line.get("head_end") if isinstance(line, dict) else None)
        out["tail_end"] = _normalize_line_end_template(line.get("tail_end") if isinstance(line, dict) else None)
    return _compact_nested_template_value(out) or _null_object_from_template(out)


def _normalize_line_point_template(point: Any) -> Optional[Dict[str, Any]]:
    template = {
        "x": point.get("x") if isinstance(point, dict) else None,
        "y": point.get("y") if isinstance(point, dict) else None,
    }
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_line_shape_xml_template(shape_xml: Any) -> Optional[Dict[str, Any]]:
    template = {
        "p1": _normalize_line_point_template(shape_xml.get("p1")) if isinstance(shape_xml, dict) else _null_object_from_template({"x": None, "y": None}),
        "p2": _normalize_line_point_template(shape_xml.get("p2")) if isinstance(shape_xml, dict) else _null_object_from_template({"x": None, "y": None}),
        "line": _normalize_line_style_template(shape_xml.get("line"), include_arrow_ends=True) if isinstance(shape_xml, dict) else _null_object_from_template(
            {
                "width_pt": None,
                "color": {"type": None, "rgb": None, "alpha_val": 100000},
                "dash": None,
                "head_end": {"type": None, "w": None, "len": None},
                "tail_end": {"type": None, "w": None, "len": None},
            }
        ),
    }
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_table_borders_template(borders: Any) -> Optional[Dict[str, Any]]:
    template = {
        "l": _normalize_line_style_template(borders.get("l"), include_arrow_ends=False) if isinstance(borders, dict) else _null_object_from_template(
            {"width_pt": None, "color": {"type": None, "rgb": None, "mods": None}, "dash": None}
        ),
        "r": _normalize_line_style_template(borders.get("r"), include_arrow_ends=False) if isinstance(borders, dict) else _null_object_from_template(
            {"width_pt": None, "color": {"type": None, "rgb": None, "mods": None}, "dash": None}
        ),
        "t": _normalize_line_style_template(borders.get("t"), include_arrow_ends=False) if isinstance(borders, dict) else _null_object_from_template(
            {"width_pt": None, "color": {"type": None, "rgb": None, "mods": None}, "dash": None}
        ),
        "b": _normalize_line_style_template(borders.get("b"), include_arrow_ends=False) if isinstance(borders, dict) else _null_object_from_template(
            {"width_pt": None, "color": {"type": None, "rgb": None, "mods": None}, "dash": None}
        ),
    }
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_table_cell_template(cell: Any) -> Optional[Dict[str, Any]]:
    if not isinstance(cell, dict):
        return None
    paragraphs = cell.get("paragraphs")
    return _compact_nested_template_value(
        {
        "row": cell.get("row"),
        "col": cell.get("col"),
        "text": cell.get("text"),
        "paragraphs": (
            [_normalize_paragraph_template(item) for item in paragraphs if isinstance(item, dict)]
            if isinstance(paragraphs, list)
            else None
        ),
        "margin_left_pt": cell.get("margin_left_pt"),
        "margin_top_pt": cell.get("margin_top_pt"),
        "margin_right_pt": cell.get("margin_right_pt"),
        "margin_bottom_pt": cell.get("margin_bottom_pt"),
        "vertical_anchor": cell.get("vertical_anchor"),
        "body_pr": _normalize_body_pr_template(cell.get("body_pr"), include_textbox_insets=False),
        "fill": _normalize_color_template(cell.get("fill")),
        "borders": _normalize_table_borders_template(cell.get("borders")),
        }
    )


def _normalize_2d_table_cells_template(cells: Any) -> Optional[List[List[Optional[Dict[str, Any]]]]]:
    if not isinstance(cells, list):
        return None
    out: List[List[Optional[Dict[str, Any]]]] = []
    for row in cells:
        if not isinstance(row, list):
            out.append([])
            continue
        out.append([_normalize_table_cell_template(cell) for cell in row if isinstance(cell, dict)])
    return out


def _normalize_image_like_layer(layer: Dict[str, Any]) -> Dict[str, Any]:
    out = {
        "slide": layer.get("slide"),
        "shape_name": layer.get("shape_name"),
        "kind": layer.get("kind"),
        "box": _normalize_box_template(layer.get("box")),
        "saved_path": layer.get("saved_path"),
        "caption": layer.get("caption"),
        "short_caption": layer.get("short_caption"),
    }
    kind = str(layer.get("kind") or "")
    structure_info = layer.get("structure_info")
    if kind in {"ppt_graph_geo", "ppt_graph_line", "ppt_graph_table"} and isinstance(structure_info, dict):
        info_out: Dict[str, Any] = {
            "box": _normalize_box_template(structure_info.get("box")),
            "rotation_deg": structure_info.get("rotation_deg"),
            "flip_h": structure_info.get("flip_h"),
            "flip_v": structure_info.get("flip_v"),
        }
        if kind == "ppt_graph_line":
            info_out["shape_xml"] = _normalize_line_shape_xml_template(structure_info.get("shape_xml"))
        elif kind == "ppt_graph_geo":
            info_out["shape_xml"] = structure_info.get("shape_xml")
            info_out["fill"] = _normalize_color_template(structure_info.get("fill"))
            info_out["line"] = _normalize_line_style_template(structure_info.get("line"), include_arrow_ends=True)
        elif kind == "ppt_graph_table":
            col_widths = structure_info.get("col_widths_pt")
            row_heights = structure_info.get("row_heights_pt")
            info_out["rows"] = structure_info.get("rows")
            info_out["cols"] = structure_info.get("cols")
            info_out["table_style_id"] = structure_info.get("table_style_id")
            info_out["col_widths_pt"] = list(col_widths) if isinstance(col_widths, list) else None
            info_out["row_heights_pt"] = list(row_heights) if isinstance(row_heights, list) else None
            info_out["cells"] = _normalize_2d_table_cells_template(structure_info.get("cells"))
        out["structure_info"] = info_out
    return out


def _normalize_text_layer(layer: Dict[str, Any]) -> Dict[str, Any]:
    paragraphs = layer.get("paragraphs")
    normalized = {
        "slide": layer.get("slide"),
        "shape_name": layer.get("shape_name"),
        "kind": layer.get("kind"),
        "rotation_deg": layer.get("rotation_deg"),
        "flip_h": layer.get("flip_h"),
        "flip_v": layer.get("flip_v"),
        "box": _normalize_box_template(layer.get("box")),
        "paragraphs": (
            [_normalize_paragraph_template(item) for item in paragraphs if isinstance(item, dict)]
            if isinstance(paragraphs, list)
            else None
        ),
        "body_pr_anchor": layer.get("body_pr_anchor"),
    }
    filled = _fill_value_from_template(_minimal_text_schema_template(), normalized)
    compact_body_pr = _normalize_body_pr_compact(layer.get("body_pr"), include_textbox_insets=True)
    if compact_body_pr is not None:
        filled["body_pr"] = compact_body_pr
    return filled


def _normalize_graph_line_layer(layer: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "slide": layer.get("slide"),
        "shape_name": layer.get("shape_name"),
        "kind": layer.get("kind"),
        "rotation_deg": layer.get("rotation_deg"),
        "flip_h": layer.get("flip_h"),
        "flip_v": layer.get("flip_v"),
        "box": _normalize_box_template(layer.get("box")),
        "shape_xml": _normalize_line_shape_xml_template(layer.get("shape_xml")),
    }


def _prune_graph_line_protocol_defaults(layer: Dict[str, Any]) -> Dict[str, Any]:
    shape_xml = layer.get("shape_xml")
    if not isinstance(shape_xml, dict):
        return layer

    if shape_xml.get("format") in {None, "shape_spec_v1"}:
        shape_xml.pop("format", None)
    if shape_xml.get("type") in {None, "line"}:
        shape_xml.pop("type", None)

    line = shape_xml.get("line")
    if not isinstance(line, dict):
        return layer
    color = line.get("color")
    if not isinstance(color, dict):
        return layer

    alpha_val = color.get("alpha_val")
    kept_mods: List[Dict[str, Any]] = []
    for item in color.get("mods") or []:
        if not isinstance(item, dict):
            continue
        op = item.get("op")
        val = item.get("val")
        if not op or val is None:
            continue
        if str(op).lower() == "alpha":
            if alpha_val is None:
                alpha_val = val
            continue
        kept_mods.append({"op": str(op), "val": val})

    if alpha_val is None:
        alpha_val = 100000
    try:
        color["alpha_val"] = int(alpha_val)
    except Exception:
        color["alpha_val"] = 100000
    if kept_mods:
        color["mods"] = kept_mods
    else:
        color.pop("mods", None)
    return layer


def _normalize_graph_geo_layer(layer: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "slide": layer.get("slide"),
        "shape_name": layer.get("shape_name"),
        "kind": layer.get("kind"),
        "rotation_deg": layer.get("rotation_deg"),
        "flip_h": layer.get("flip_h"),
        "flip_v": layer.get("flip_v"),
        "box": _normalize_box_template(layer.get("box")),
        "shape_xml": layer.get("shape_xml"),
        "fill": _normalize_color_template(layer.get("fill")),
        "line": _normalize_line_style_template(layer.get("line"), include_arrow_ends=True),
    }


def _normalize_table_layer(layer: Dict[str, Any]) -> Dict[str, Any]:
    col_widths = layer.get("col_widths_pt")
    row_heights = layer.get("row_heights_pt")
    out = {
        "slide": layer.get("slide"),
        "shape_name": layer.get("shape_name"),
        "kind": layer.get("kind"),
        "rotation_deg": layer.get("rotation_deg"),
        "flip_h": layer.get("flip_h"),
        "flip_v": layer.get("flip_v"),
        "box": _normalize_box_template(layer.get("box")),
        "rows": layer.get("rows"),
        "cols": layer.get("cols"),
        "table_style_id": layer.get("table_style_id"),
        "col_widths_pt": list(col_widths) if isinstance(col_widths, list) else None,
        "row_heights_pt": list(row_heights) if isinstance(row_heights, list) else None,
        "cells": _normalize_2d_table_cells_template(layer.get("cells")),
    }
    if layer.get("saved_path"):
        out["saved_path"] = layer.get("saved_path")
    return out


def _validate_dual_export_layer(layer: Dict[str, Any]) -> Dict[str, Any]:
    kind = str(layer.get("kind") or "")
    if kind not in {"ppt_graph_geo", "ppt_graph_line", "ppt_graph_table"}:
        return layer

    missing: List[str] = []
    if not layer.get("saved_path"):
        missing.append("saved_path")

    structure_info = layer.get("structure_info")
    if not isinstance(structure_info, dict):
        missing.append("structure_info")
    else:
        if not isinstance(structure_info.get("box"), dict):
            missing.append("structure_info.box")
        for key in ("rotation_deg", "flip_h", "flip_v"):
            if key not in structure_info:
                missing.append(f"structure_info.{key}")
        if kind in {"ppt_graph_geo", "ppt_graph_line"}:
            if not isinstance(structure_info.get("shape_xml"), dict):
                missing.append("structure_info.shape_xml")
        elif kind == "ppt_graph_table":
            if structure_info.get("rows") is None:
                missing.append("structure_info.rows")
            if structure_info.get("cols") is None:
                missing.append("structure_info.cols")
            if not isinstance(structure_info.get("cells"), list):
                missing.append("structure_info.cells")

    if missing:
        raise RuntimeError(
            f"{kind} dual export failed for shape={layer.get('shape_name')!r}: missing {', '.join(missing)}"
        )
    return layer


def _normalize_background_fill_template(fill: Any) -> Optional[Dict[str, Any]]:
    template = {
        "type": fill.get("type") if isinstance(fill, dict) else None,
        "color_type": fill.get("color_type") if isinstance(fill, dict) else None,
        "rgb": fill.get("rgb") if isinstance(fill, dict) else None,
    }
    return _compact_nested_template_value(template) or _null_object_from_template(template)


def _normalize_slide_canvas_layer(layer: Dict[str, Any]) -> Dict[str, Any]:
    out = {
        "slide": layer.get("slide"),
        "shape_name": layer.get("shape_name"),
        "kind": layer.get("kind"),
        "box": _normalize_box_template(layer.get("box")),
        "canvas_width_emu": layer.get("canvas_width_emu"),
        "canvas_height_emu": layer.get("canvas_height_emu"),
        "caption": layer.get("caption"),
        "short_caption": layer.get("short_caption"),
    }
    if layer.get("saved_path"):
        out["saved_path"] = layer.get("saved_path")
    return out


def normalize_export_layer_to_template(layer: Dict[str, Any]) -> Dict[str, Any]:
    layer = _validate_dual_export_layer(dict(layer))
    kind = layer.get("kind")
    if kind in IMAGE_LIKE_EXPORT_KINDS:
        normalized = _normalize_image_like_layer(layer)
        return _apply_kind_template(normalized)
    if kind in {"ppt_graph_geo", "ppt_graph_line", "ppt_graph_table"} and layer.get("saved_path"):
        normalized = _normalize_image_like_layer(layer)
        if _kind_template_has_key(kind, "saved_path"):
            return _apply_kind_template(normalized)
        return _compact_nested_template_value(normalized) or normalized
    if kind == "text":
        return _normalize_text_layer(layer)
    if kind == "ppt_graph_line":
        normalized = _normalize_graph_line_layer(layer)
        return _prune_graph_line_protocol_defaults(_apply_kind_template(normalized))
    if kind == "ppt_graph_geo":
        normalized = _normalize_graph_geo_layer(layer)
        return _apply_kind_template(normalized)
    if kind == "ppt_graph_table":
        normalized = _normalize_table_layer(layer)
        return _apply_kind_template(normalized)
    if kind == "slide_canvas":
        normalized = _normalize_slide_canvas_layer(layer)
        return normalized
    return _apply_kind_template(dict(layer))


def _canvas_base_emu(canvas_emu: Optional[tuple[int, int]]) -> Optional[float]:
    if not canvas_emu:
        return None
    try:
        w = float(canvas_emu[0] or 0.0)
        h = float(canvas_emu[1] or 0.0)
        if w <= 0 or h <= 0:
            return None
        return math.sqrt(w * h)
    except Exception:
        return None


def _emu_rel(value_emu: Optional[float], canvas_emu: Optional[tuple[int, int]]) -> Optional[float]:
    base = _canvas_base_emu(canvas_emu)
    if value_emu is None or not base:
        return None
    try:
        return float(value_emu) / float(base)
    except Exception:
        return None


def _pt_rel(value_pt: Optional[float], canvas_emu: Optional[tuple[int, int]]) -> Optional[float]:
    if value_pt is None:
        return None
    try:
        return _emu_rel(float(value_pt) * EMU_PER_PT, canvas_emu)
    except Exception:
        return None


# #region debug-point G:report-table-gradient-fill
def _dbg_report_table_gradient_fill(hypothesis_id: str, location: str, msg: str, data: Optional[Dict[str, Any]] = None) -> None:
    try:
        import json as _json
        import urllib.request as _urlreq

        _u, _s, _r = "http://127.0.0.1:7777/event", "table-gradient-fill", os.environ.get("DEBUG_RUN_ID", "pre-fix")
        try:
            with open(".dbg/table-gradient-fill.env", "r", encoding="utf-8") as _f:
                for _line in _f.read().splitlines():
                    if _line.startswith("DEBUG_SERVER_URL="):
                        _u = _line.split("=", 1)[1] or _u
                    elif _line.startswith("DEBUG_SESSION_ID="):
                        _s = _line.split("=", 1)[1] or _s
                    elif _line.startswith("DEBUG_RUN_ID="):
                        _r = _line.split("=", 1)[1] or _r
        except Exception:
            pass

        _payload = {
            "sessionId": _s,
            "runId": _r,
            "hypothesisId": hypothesis_id,
            "location": location,
            "msg": msg,
            "data": data or {},
            "ts": int(time.time() * 1000),
        }
        _urlreq.urlopen(
            _urlreq.Request(
                _u,
                data=_json.dumps(_payload).encode("utf-8"),
                headers={"Content-Type": "application/json"},
            ),
            timeout=1.0,
        ).read()
    except Exception:
        pass
# #endregion


def _raw_100pt_rel(raw_100pt: Optional[int], canvas_emu: Optional[tuple[int, int]]) -> Optional[float]:
    if raw_100pt is None:
        return None
    try:
        return _emu_rel((float(raw_100pt) / 100.0) * EMU_PER_PT, canvas_emu)
    except Exception:
        return None


# #region debug-point A:report
def _dbg_report_method_ppt_style(hypothesis_id: str, location: str, msg: str, data: Optional[Dict[str, Any]] = None) -> None:
    try:
        import json as _json
        import urllib.request as _urlreq

        _u, _s, _r = "http://127.0.0.1:7777/event", "method-ppt-style", os.environ.get("DEBUG_RUN_ID", "pre-fix")
        try:
            with open(".dbg/method-ppt-style.env", "r", encoding="utf-8") as _f:
                for _line in _f.read().splitlines():
                    if _line.startswith("DEBUG_SERVER_URL="):
                        _u = _line.split("=", 1)[1] or _u
                    elif _line.startswith("DEBUG_SESSION_ID="):
                        _s = _line.split("=", 1)[1] or _s
                    elif _line.startswith("DEBUG_RUN_ID="):
                        _r = _line.split("=", 1)[1] or _r
        except Exception:
            pass

        _payload = {
            "sessionId": _s,
            "runId": _r,
            "hypothesisId": hypothesis_id,
            "location": location,
            "msg": msg,
            "data": data or {},
            "ts": int(time.time() * 1000),
        }
        _urlreq.urlopen(
            _urlreq.Request(
                _u,
                data=_json.dumps(_payload).encode("utf-8"),
                headers={"Content-Type": "application/json"},
            ),
            timeout=1.0,
        ).read()
    except Exception:
        pass
# #endregion


def identity_transform() -> Dict[str, float]:
    return {"a": 1.0, "b": 0.0, "c": 0.0, "d": 1.0, "tx": 0.0, "ty": 0.0}


def compose_transform(parent: Dict[str, float], child: Dict[str, float]) -> Dict[str, float]:
    return {
        "a": parent["a"] * child["a"] + parent["b"] * child["c"],
        "b": parent["a"] * child["b"] + parent["b"] * child["d"],
        "c": parent["c"] * child["a"] + parent["d"] * child["c"],
        "d": parent["c"] * child["b"] + parent["d"] * child["d"],
        "tx": parent["a"] * child["tx"] + parent["b"] * child["ty"] + parent["tx"],
        "ty": parent["c"] * child["tx"] + parent["d"] * child["ty"] + parent["ty"],
    }


def transform_point(transform: Dict[str, float], x: float, y: float) -> tuple[float, float]:
    return (
        transform["a"] * x + transform["b"] * y + transform["tx"],
        transform["c"] * x + transform["d"] * y + transform["ty"],
    )


def transform_vector(transform: Dict[str, float], x: float, y: float) -> tuple[float, float]:
    return (
        transform["a"] * x + transform["b"] * y,
        transform["c"] * x + transform["d"] * y,
    )


def transform_rotation_deg(transform: Dict[str, float]) -> float:
    return math.degrees(math.atan2(transform["c"], transform["a"])) % 360.0


def ensure_dir(p: Path) -> None:
    p.mkdir(parents=True, exist_ok=True)


def emu_to_pt(v: Optional[int]) -> Optional[float]:
    if v is None:
        return None
    try:
        return float(v) / EMU_PER_PT
    except Exception:
        return None


def _ratio_box(
    left: Optional[int],
    top: Optional[int],
    width: Optional[int],
    height: Optional[int],
    canvas_emu: Optional[tuple[int, int]],
) -> Dict[str, Any]:
    def _round_coord(v: float) -> float:
        r = round(float(v), 5)
        return 0.0 if r == -0.0 else r

    if None in (left, top, width, height) or not canvas_emu:
        return {"left": None, "top": None, "width": None, "height": None}
    canvas_width, canvas_height = canvas_emu
    if canvas_width <= 0 or canvas_height <= 0:
        return {"left": None, "top": None, "width": None, "height": None}
    return {
        "left": _round_coord(float(left) / float(canvas_width)),
        "top": _round_coord(float(top) / float(canvas_height)),
        "width": _round_coord(float(width) / float(canvas_width)),
        "height": _round_coord(float(height) / float(canvas_height)),
    }


def _round_rotation_deg(v: float) -> float:
    """Round exported rotation degrees to 2 decimals (protocol requirement)."""
    try:
        r = round(float(v), 2)
    except Exception:
        return 0.0
    return 0.0 if r == -0.0 else r


def _encode_scaled_numbers_for_json(x: Any, scale: int = 1000) -> Any:
    """Encode protocol floats as scaled integers for JSON output."""
    if isinstance(x, float):
        try:
            v = int(round(float(x) * float(scale)))
        except Exception:
            v = 0
        return 0 if v == -0 else v
    if isinstance(x, list):
        return [_encode_scaled_numbers_for_json(i, scale) for i in x]
    if isinstance(x, dict):
        return {k: _encode_scaled_numbers_for_json(v, scale) for k, v in x.items()}
    return x


def _shape_bounds_box_emu(shp, transform: Optional[Dict[str, float]] = None) -> Optional[tuple[int, int, int, int]]:
    # left/top/width/height might not exist for some shape types; keep it defensive.
    # For shapes inside groups, `python-pptx` exposes local coordinates; apply the accumulated
    # group transform so every exported box is in slide absolute coordinates.
    transform = transform or identity_transform()
    left = getattr(shp, "left", None)
    top = getattr(shp, "top", None)
    width = getattr(shp, "width", None)
    height = getattr(shp, "height", None)
    if None not in (left, top, width, height):
        cx, cy = transform_point(transform, float(left) + float(width) / 2.0, float(top) + float(height) / 2.0)
        vx = transform_vector(transform, float(width), 0.0)
        vy = transform_vector(transform, 0.0, float(height))
        abs_width = math.hypot(vx[0], vx[1])
        abs_height = math.hypot(vy[0], vy[1])
        left = int(round(cx - abs_width / 2.0))
        top = int(round(cy - abs_height / 2.0))
        width = int(round(abs_width))
        height = int(round(abs_height))
    if None in (left, top, width, height):
        return None
    return int(left), int(top), int(width), int(height)


def shape_bounds_box(
    shp, transform: Optional[Dict[str, float]] = None, canvas_emu: Optional[tuple[int, int]] = None
) -> Dict[str, Any]:
    box_emu = _shape_bounds_box_emu(shp, transform)
    if box_emu is None:
        return _ratio_box(None, None, None, None, canvas_emu)
    left, top, width, height = box_emu
    return _ratio_box(left, top, width, height, canvas_emu)


def _shape_box_emu(shp, transform: Optional[Dict[str, float]] = None) -> Optional[tuple[int, int, int, int]]:
    bounds_box = _shape_bounds_box_emu(shp, transform)
    if bounds_box is None:
        return None
    left, top, width, height = bounds_box

    # Export the visual axis-aligned bounding box after rotation.
    rotation_deg = shape_rotation_deg(shp, transform)
    if rotation_deg:
        theta = math.radians(rotation_deg)
        bbox_w = abs(width * math.cos(theta)) + abs(height * math.sin(theta))
        bbox_h = abs(width * math.sin(theta)) + abs(height * math.cos(theta))
        cx = left + width / 2.0
        cy = top + height / 2.0
        left = int(round(cx - bbox_w / 2.0))
        top = int(round(cy - bbox_h / 2.0))
        width = int(round(bbox_w))
        height = int(round(bbox_h))
    return int(left), int(top), int(width), int(height)


def shape_box(shp, transform: Optional[Dict[str, float]] = None, canvas_emu: Optional[tuple[int, int]] = None) -> Dict[str, Any]:
    box_emu = _shape_box_emu(shp, transform)
    if box_emu is None:
        return _ratio_box(None, None, None, None, canvas_emu)
    left, top, width, height = box_emu
    return _ratio_box(left, top, width, height, canvas_emu)


def child_transform_for_group(group_shp, parent_transform: Optional[Dict[str, float]] = None) -> Dict[str, float]:
    """
    Build a transform that maps coordinates in this group's child space into slide coordinates.

    OOXML group transform:
      parent_space = off + (child_space - chOff) * (ext / chExt)
    We compose that with the parent transform so nested groups are handled correctly.
    """
    parent_transform = parent_transform or identity_transform()
    xfrm = getattr(getattr(getattr(group_shp, "_element", None), "grpSpPr", None), "xfrm", None)
    if xfrm is None:
        # Fallback: if internals differ, at least offset by the group's own box.
        left = getattr(group_shp, "left", 0) or 0
        top = getattr(group_shp, "top", 0) or 0
        return compose_transform(
            parent_transform,
            {"a": 1.0, "b": 0.0, "c": 0.0, "d": 1.0, "tx": float(left), "ty": float(top)},
        )

    off_x = float(getattr(xfrm.off, "x", 0) or 0)
    off_y = float(getattr(xfrm.off, "y", 0) or 0)
    ext_cx = float(getattr(xfrm.ext, "cx", 0) or 0)
    ext_cy = float(getattr(xfrm.ext, "cy", 0) or 0)
    ch_off_x = float(getattr(xfrm.chOff, "x", 0) or 0)
    ch_off_y = float(getattr(xfrm.chOff, "y", 0) or 0)
    ch_ext_cx = float(getattr(xfrm.chExt, "cx", 0) or 0)
    ch_ext_cy = float(getattr(xfrm.chExt, "cy", 0) or 0)

    scale_x = (ext_cx / ch_ext_cx) if ch_ext_cx else 1.0
    scale_y = (ext_cy / ch_ext_cy) if ch_ext_cy else 1.0
    rot_deg = float(xfrm.get("rot", "0") or 0) / 60000.0
    theta = math.radians(rot_deg)
    cos_t = math.cos(theta)
    sin_t = math.sin(theta)

    tx0 = off_x - ch_off_x * scale_x
    ty0 = off_y - ch_off_y * scale_y
    center_x = off_x + ext_cx / 2.0
    center_y = off_y + ext_cy / 2.0

    group_transform = {
        "a": cos_t * scale_x,
        "b": -sin_t * scale_y,
        "c": sin_t * scale_x,
        "d": cos_t * scale_y,
        "tx": cos_t * tx0 - sin_t * ty0 + center_x - cos_t * center_x + sin_t * center_y,
        "ty": sin_t * tx0 + cos_t * ty0 + center_y - sin_t * center_x - cos_t * center_y,
    }
    return compose_transform(parent_transform, group_transform)


def font_color_to_dict(font) -> Optional[Dict[str, Any]]:
    c = getattr(font, "color", None)
    if c is None:
        return None
    try:
        ctype = c.type
    except Exception:
        return None

    out: Dict[str, Any] = {"type": None, "rgb": None}
    if ctype == MSO_COLOR_TYPE.RGB and c.rgb is not None:
        out["type"] = "RGB"
        out["rgb"] = str(c.rgb)  # e.g. 'FF00AA'
    elif (str(ctype).endswith("THEME") or str(ctype) == "SCHEME") and getattr(c, "theme_color", None) is not None:
        out["type"] = "THEME"
    else:
        out["type"] = str(ctype)
    return out


def _text_color_from_run(run) -> Optional[Dict[str, Any]]:
    try:
        r_pr = run._r.find(qn("a:rPr"))
    except Exception:
        r_pr = None
    if r_pr is None:
        return None
    spec = _solid_fill_spec_from_node(r_pr)
    if not spec:
        return None
    out: Dict[str, Any] = {"type": None, "rgb": None}
    if spec.get("type") == "srgb" and spec.get("rgb"):
        out["type"] = "RGB"
        out["rgb"] = spec.get("rgb")
    elif spec.get("type") == "scheme" and spec.get("scheme"):
        out["type"] = "SCHEME"
        out["scheme"] = spec.get("scheme")
    if spec.get("mods"):
        out["mods"] = spec.get("mods")
    return out


def shape_rotation_deg(shp, transform: Optional[Dict[str, float]] = None) -> float:
    try:
        local_rot = float(getattr(shp, "rotation", 0.0) or 0.0)
    except Exception:
        local_rot = 0.0
    base_rot = transform_rotation_deg(transform or identity_transform())
    return (base_rot + local_rot) % 360.0


def shape_flip_flags(shp) -> Dict[str, bool]:
    try:
        xfrm = getattr(getattr(shp, "_element", None), "spPr", None)
        xfrm = getattr(xfrm, "xfrm", None)
        if xfrm is None:
            xfrm = getattr(getattr(getattr(shp, "_element", None), "grpSpPr", None), "xfrm", None)
        if xfrm is None:
            return {"flip_h": False, "flip_v": False}
        return {
            "flip_h": str(xfrm.get("flipH", "false")).lower() in {"1", "true"},
            "flip_v": str(xfrm.get("flipV", "false")).lower() in {"1", "true"},
        }
    except Exception:
        return {"flip_h": False, "flip_v": False}


def _resolve_image_parts(slide, blip):
    def _rel_to_part(r_id: Optional[str]):
        if not r_id:
            return None
        try:
            rel = slide.part.rels[r_id]
        except KeyError:
            return None
        if getattr(rel, "is_external", False):
            return None
        part = getattr(rel, "target_part", None)
        if part is None or not hasattr(part, "blob"):
            return None
        ctype = getattr(part, "content_type", None)
        if not (ctype and str(ctype).startswith("image/")):
            return None
        return part

    svg_rids: List[str] = []
    try:
        blip_et = ET.fromstring(ET.tostring(blip, encoding="utf-8"))
        svg_blip = blip_et.find(".//asvg:svgBlip", NS)
        if svg_blip is not None:
            rid = svg_blip.get(qn("r:embed"))
            if rid:
                svg_rids.append(rid)
    except Exception:
        svg_rids = []

    primary_part = None
    for svg_rid in svg_rids:
        part = _rel_to_part(svg_rid)
        if part is not None:
            primary_part = part
            break

    preview_part = _rel_to_part(blip.get(qn("r:embed")))
    if primary_part is not None:
        return primary_part, preview_part
    return preview_part, None


def _xml_spacing_to_dict(parent, tag_local: str) -> Optional[Dict[str, Any]]:
    if parent is None:
        return None
    node = parent.find(qn(f"a:{tag_local}"))
    if node is None:
        return None
    pts = node.find(qn("a:spcPts"))
    if pts is not None:
        try:
            raw = int(pts.get("val"))
            return {"mode": "points", "raw": raw}
        except Exception:
            return None
    pct = node.find(qn("a:spcPct"))
    if pct is not None:
        try:
            raw = int(pct.get("val"))
            return {"mode": "percent", "raw": raw}
        except Exception:
            return None
    return None


def _xml_bullet_to_dict(p_pr) -> Optional[Dict[str, Any]]:
    if p_pr is None:
        return None
    bu_none = p_pr.find(qn("a:buNone"))
    if bu_none is not None:
        return {"type": "none"}
    bu_char = p_pr.find(qn("a:buChar"))
    if bu_char is not None:
        out: Dict[str, Any] = {"type": "char", "char": bu_char.get("char")}
        bu_font = p_pr.find(qn("a:buFont"))
        if bu_font is not None and bu_font.get("typeface"):
            out["font"] = {"typeface": bu_font.get("typeface")}
        return out
    bu_auto = p_pr.find(qn("a:buAutoNum"))
    if bu_auto is not None:
        out = {"type": "autoNum"}
        if bu_auto.get("type"):
            out["scheme"] = bu_auto.get("type")
        if bu_auto.get("startAt"):
            try:
                out["startAt"] = int(bu_auto.get("startAt"))
            except Exception:
                pass
        return out
    return None


def _xml_paragraph_attrs_to_dict(p_pr, canvas_emu: Optional[tuple[int, int]]) -> Optional[Dict[str, Any]]:
    if p_pr is None:
        return None
    out: Dict[str, Any] = {}
    for key in ("marL", "marR", "indent", "lvl", "algn"):
        val = p_pr.get(key)
        if val is None:
            continue
        if key in {"marL", "marR", "indent"}:
            try:
                rel = _emu_rel(int(val), canvas_emu)
                if rel is not None:
                    out[f"{key}_rel"] = rel
            except Exception:
                continue
        elif key == "lvl":
            try:
                out[key] = int(val)
            except Exception:
                continue
        else:
            out[key] = val
    return out or None


def _normalize_spacing_dict(spec: Optional[Dict[str, Any]], canvas_emu: Optional[tuple[int, int]]) -> Optional[Dict[str, Any]]:
    if not isinstance(spec, dict):
        return spec
    mode = str(spec.get("mode") or "").lower()
    if mode != "points":
        return spec
    out = {"mode": "points"}
    rel = _raw_100pt_rel(spec.get("raw"), canvas_emu)
    if rel is not None:
        out["rel"] = rel
    return out


def _xml_int_attr(node, attr_name: str) -> Optional[int]:
    if node is None:
        return None
    val = node.get(attr_name)
    if val is None:
        return None
    try:
        return int(val)
    except Exception:
        return None


def _run_spacing_info(run, canvas_emu: Optional[tuple[int, int]]) -> Dict[str, Any]:
    r_pr = run._r.find(qn("a:rPr"))
    char_spacing_raw = _xml_int_attr(r_pr, "spc")
    kern_raw = _xml_int_attr(r_pr, "kern")
    return {
        "char_spacing_raw": char_spacing_raw,
        "char_spacing_rel": _raw_100pt_rel(char_spacing_raw, canvas_emu),
        "kern_raw": kern_raw,
    }


def _run_merge_key(run_info: Dict[str, Any]) -> Tuple[Any, ...]:
    color = run_info.get("color")
    if isinstance(color, dict):
        color_key = json.dumps(color, ensure_ascii=False, sort_keys=True)
    else:
        color_key = color
    return (
        run_info.get("font_name"),
        run_info.get("font_size_rel"),
        run_info.get("bold"),
        run_info.get("italic"),
        run_info.get("underline"),
        run_info.get("strike"),
        color_key,
        run_info.get("char_spacing_raw"),
        run_info.get("char_spacing_rel"),
        run_info.get("kern_raw"),
    )


def _merge_adjacent_runs(runs: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    merged: List[Dict[str, Any]] = []
    for run in runs:
        if not merged:
            merged.append(dict(run))
            continue
        prev = merged[-1]
        if _run_merge_key(prev) == _run_merge_key(run):
            prev["text"] = f'{prev.get("text") or ""}{run.get("text") or ""}'
            continue
        merged.append(dict(run))
    for idx, run in enumerate(merged):
        run["run_index"] = idx
    return merged


def _normalize_text_breaks(text: Optional[str]) -> str:
    if not text:
        return ""
    # PowerPoint/python-pptx may use vertical-tab for soft line breaks inside a paragraph.
    return str(text).replace("\r\n", "\n").replace("\r", "\n").replace("\v", "\n")


def _xml_el_to_spec(el: Optional[ET.Element]) -> Optional[Dict[str, Any]]:
    if el is None:
        return None
    tag = el.tag
    ns_uri = None
    local = tag
    if tag.startswith("{") and "}" in tag:
        ns_uri, local = tag[1:].split("}", 1)
    ns_name = None
    if ns_uri:
        for k, v in NS.items():
            if v == ns_uri:
                ns_name = k
                break
    spec: Dict[str, Any] = {
        "tag": local,
        "attrs": dict(el.attrib),
        "children": [],
    }
    if ns_name:
        spec["ns"] = ns_name
    if el.text:
        spec["text"] = el.text
    for child in list(el):
        child_spec = _xml_el_to_spec(child)
        if child_spec is not None:
            spec["children"].append(child_spec)
    return spec


def _fallback_text_from_shape_xml(shp, canvas_emu: Optional[tuple[int, int]] = None) -> Optional[List[Dict[str, Any]]]:
    """Best-effort text extraction from raw shape XML.

    Some shapes contain text runs that python-pptx doesn't surface via TextFrame,
    e.g. a:fld or certain drawing constructs. This fallback keeps data minimal:
    paragraphs with runs + basic rPr-derived styling.
    """
    try:
        raw = ET.tostring(getattr(shp, "_element", None), encoding="utf-8")
        if not raw:
            return None
        root = ET.fromstring(raw)
    except Exception:
        return None

    tx = root.find(".//a:txBody", NS)
    if tx is None:
        return None

    paragraphs: List[Dict[str, Any]] = []
    for p in tx.findall("a:p", NS):
        runs: List[Dict[str, Any]] = []
        text_parts: List[str] = []
        for child in list(p):
            tag = child.tag.rsplit("}", 1)[-1]
            if tag in {"r", "fld"}:
                r_pr = child.find("a:rPr", NS)
                t_nodes = child.findall("a:t", NS)
                txt = "".join((t.text or "") for t in t_nodes)
                if not txt:
                    continue
                txt = _normalize_text_breaks(txt)
                run_info: Dict[str, Any] = {
                    "run_index": len(runs),
                    "text": txt,
                    "font_name": None,
                    "font_size_rel": None,
                    "bold": False,
                    "italic": False,
                    "underline": False,
                    "strike": False,
                    "color": None,
                    "char_spacing_raw": None,
                    "char_spacing_rel": None,
                    "kern_raw": None,
                }
                if r_pr is not None:
                    # Font name from latin/rFonts.
                    latin = r_pr.find("a:latin", NS)
                    if latin is not None and latin.get("typeface"):
                        run_info["font_name"] = latin.get("typeface")
                    r_fonts = r_pr.find("a:rFonts", NS)
                    if run_info["font_name"] is None and r_fonts is not None:
                        run_info["font_name"] = r_fonts.get("ascii") or r_fonts.get("hAnsi")
                    # Size/bold/italic/underline.
                    if r_pr.get("sz"):
                        try:
                            run_info["font_size_rel"] = _pt_rel(int(r_pr.get("sz")) / 100.0, canvas_emu)
                        except Exception:
                            pass
                    if r_pr.get("b") is not None:
                        run_info["bold"] = _style_flag(r_pr.get("b"))
                    if r_pr.get("i") is not None:
                        run_info["italic"] = _style_flag(r_pr.get("i"))
                    if r_pr.get("u") is not None:
                        run_info["underline"] = _style_flag(r_pr.get("u"))
                    if r_pr.get("strike") is not None:
                        run_info["strike"] = _style_flag(r_pr.get("strike"), strike=True)
                    # Color
                    color_spec = _solid_fill_spec_from_node(r_pr)
                    if color_spec:
                        if color_spec.get("type") == "srgb":
                            run_info["color"] = {"type": "RGB", "rgb": color_spec.get("rgb"), "mods": color_spec.get("mods")}
                        elif color_spec.get("type") == "scheme":
                            run_info["color"] = {"type": "SCHEME", "scheme": color_spec.get("scheme"), "mods": color_spec.get("mods")}
                runs.append(run_info)
                text_parts.append(txt)
            elif tag == "br":
                text_parts.append("\n")
            else:
                continue
        joined = "".join(text_parts).strip("\n")
        if not joined and not runs:
            continue
        runs = _merge_adjacent_runs(runs)
        paragraphs.append(
            {
                "runs": runs,
                "line_spacing": None,
                "space_before": None,
            }
        )
    return paragraphs or None


def _paragraph_info(para, canvas_emu: Optional[tuple[int, int]] = None) -> Dict[str, Any]:
    p_pr = para._p.find(qn("a:pPr"))
    end_rpr = para._p.find(qn("a:endParaRPr"))
    runs_info: List[Dict[str, Any]] = []
    paragraph_text_parts: List[str] = []

    for ri, run in enumerate(para.runs):
        if not run.text:
            continue
        f = run.font
        size_pt = None
        try:
            if f.size is not None:
                size_pt = f.size.pt
        except Exception:
            size_pt = None

        run_info = {
            "run_index": ri,
            "text": _normalize_text_breaks(run.text),
            "font_name": f.name,
            "font_size_rel": _pt_rel(size_pt, canvas_emu),
            "bold": _style_flag(f.bold),
            "italic": _style_flag(f.italic),
            "underline": _style_flag(f.underline),
            "strike": False,
            "color": _text_color_from_run(run) or font_color_to_dict(f),
        }
        r_pr = run._r.find(qn("a:rPr"))
        if r_pr is not None and r_pr.get("strike") is not None:
            run_info["strike"] = _style_flag(r_pr.get("strike"), strike=True)
        run_info.update(_run_spacing_info(run, canvas_emu))
        # #region debug-point C:text-token-export
        if "token" in (run_info.get("text") or "").lower():
            try:
                _dbg_report_method_ppt_style(
                    "C",
                    "test.py:_paragraph_info",
                    "[DEBUG] export token run",
                    {
                        "text": run_info.get("text"),
                        "font_name": run_info.get("font_name"),
                        "font_size_rel": run_info.get("font_size_rel"),
                        "bold": run_info.get("bold"),
                        "italic": run_info.get("italic"),
                        "underline": run_info.get("underline"),
                        "color": run_info.get("color"),
                        "rpr_xml": ET.tostring(run._r.find(qn("a:rPr")), encoding="unicode") if run._r.find(qn("a:rPr")) is not None else None,
                    },
                )
            except Exception:
                pass
        # #endregion
        runs_info.append(run_info)
        paragraph_text_parts.append(_normalize_text_breaks(run.text))

    runs_info = _merge_adjacent_runs(runs_info)
    bullet_info = _xml_bullet_to_dict(p_pr)

    line_spacing = _normalize_spacing_dict(_xml_spacing_to_dict(p_pr, "lnSpc"), canvas_emu)
    space_before = _normalize_spacing_dict(_xml_spacing_to_dict(p_pr, "spcBef"), canvas_emu)

    return {
        "runs": runs_info,
        "bullet": bullet_info,
        "ppr_attrs": _xml_paragraph_attrs_to_dict(p_pr, canvas_emu),
        "line_spacing": line_spacing,
        "space_before": space_before,
    }


def extract_text_layer(
    shp,
    slide_index: int,
    canvas_emu: Optional[tuple[int, int]],
    transform: Optional[Dict[str, float]] = None,
) -> Optional[Dict[str, Any]]:
    tf = shp.text_frame
    paragraphs: List[Dict[str, Any]] = []
    autonum_counters: Dict[Tuple[Any, ...], int] = {}
    for pi, para in enumerate(tf.paragraphs):
        para_info = _paragraph_info(para, canvas_emu)
        bullet = para_info.get("bullet")
        ppr_attrs = para_info.get("ppr_attrs") if isinstance(para_info.get("ppr_attrs"), dict) else {}
        counter_key = (
            ppr_attrs.get("lvl"),
            bullet.get("type") if isinstance(bullet, dict) else None,
            bullet.get("scheme") if isinstance(bullet, dict) else None,
        )
        para_info["bullet"] = _inline_bullet_into_runs(
            para_info.get("runs") or [],
            bullet,
            autonum_counters=autonum_counters,
            counter_key=counter_key,
        )
        paragraphs.append(para_info)

    # Many shapes technically "have_text_frame" but contain no editable text.
    has_runs = any(para_info.get("runs") for para_info in paragraphs)
    if not has_runs and not (tf.text or "").strip():
        # Fallback to raw XML extraction for cases python-pptx doesn't surface.
        fallback_paras = _fallback_text_from_shape_xml(shp, canvas_emu)
        if not fallback_paras:
            return None
        paragraphs = fallback_paras

    # Capture txBody/a:bodyPr. It affects line breaking and vertical alignment.
    # Keep it as a compact spec (attrs + autofit mode) so rebuild can reproduce it faithfully.
    body_pr_anchor = None
    try:
        root = ET.fromstring(ET.tostring(getattr(shp, "_element", None), encoding="utf-8"))
        bp = root.find(".//p:txBody/a:bodyPr", NS)
        if bp is not None:
            body_pr_anchor = bp.attrib.get("anchor")
    except Exception:
        body_pr_anchor = None

    out = {
        "slide": slide_index,
        "shape_name": getattr(shp, "name", None),
        "kind": "text",
        "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
        **shape_flip_flags(shp),
        # Use the un-rotated bbox (xfrm off/ext) as the placement source of truth.
        # The visual bbox after rotation would "double count" rotation on rebuild.
        "box": shape_bounds_box(shp, transform, canvas_emu),
        "paragraphs": paragraphs,
        "body_pr_anchor": body_pr_anchor,
    }
    return out


def save_image_rgba_png(blob: bytes, out_path: Path) -> None:
    im = Image.open(io.BytesIO(blob))
    im = im.convert("RGBA")
    im.save(out_path, format="PNG")


def export_asset_image(
    *,
    blob: bytes,
    ext: str,
    content_type: Optional[str],
    out_base: str,
    assets_dir: Path,
) -> Dict[str, Any]:
    ext = (ext or "").lower()
    is_svg = (ext == "svg") or (content_type == "image/svg+xml")

    if is_svg:
        out_path = assets_dir / f"{out_base}.svg"
        out_path.write_bytes(blob)
        return {
            "kind": "image_svg",
            "saved_path": str(out_path),
        }

    out_path = assets_dir / f"{out_base}_image_png_rgba.png"
    try:
        save_image_rgba_png(blob, out_path)
        return {
            "kind": "image_png_rgba",
            "saved_path": str(out_path),
        }
    except Exception:
        # Some PPT images can be EMF/WMF and not readable by Pillow.
        fallback_ext = ext if ext else "bin"
        fallback = assets_dir / f"{out_base}.{fallback_ext}"
        fallback.write_bytes(blob)
        return {
            "kind": "image_raw",
            "saved_path": str(fallback),
        }


def _load_blob_as_rgba(
    *,
    blob: bytes,
    ext: str,
    content_type: Optional[str],
    box_emu: Optional[tuple[int, int, int, int]],
) -> Optional[Image.Image]:
    ext = (ext or "").lower()
    is_svg = (ext == "svg") or (content_type == "image/svg+xml")
    out_w, out_h = _raster_size_2048_for_box(box_emu)

    if is_svg:
        png_bytes = _svg_bytes_to_png_bytes(blob, out_w, out_h)
        if png_bytes is None:
            return None
        try:
            return Image.open(io.BytesIO(png_bytes)).convert("RGBA")
        except Exception:
            return None

    try:
        im = Image.open(io.BytesIO(blob)).convert("RGBA")
    except Exception:
        return None
    if im.size != (out_w, out_h):
        im = im.resize((out_w, out_h), resample=Image.LANCZOS)
    return im


def _load_fill_blob_as_rgba(blob: bytes, ext: str, content_type: Optional[str]) -> Optional[Image.Image]:
    """Load fill images without pre-warping them to the destination box aspect ratio."""
    ext = (ext or "").lower()
    is_svg = (ext == "svg") or (content_type == "image/svg+xml")

    if is_svg:
        aspect = _svg_aspect_ratio(blob)
        if aspect is None or aspect <= 0:
            aspect = 1.0

        if aspect >= 1.0:
            out_w = RASTER_TARGET_SIDE
            out_h = max(1, int(round(out_w / aspect)))
        else:
            out_h = RASTER_TARGET_SIDE
            out_w = max(1, int(round(out_h * aspect)))
        png_bytes = _svg_bytes_to_png_bytes(blob, out_w, out_h)
        if png_bytes is None:
            return None
        try:
            return Image.open(io.BytesIO(png_bytes)).convert("RGBA")
        except Exception:
            return None

    try:
        return Image.open(io.BytesIO(blob)).convert("RGBA")
    except Exception:
        return None


def _save_final_image_layer(
    *,
    im: Image.Image,
    out_path: Path,
    final_box_emu: Optional[tuple[int, int, int, int]],
    canvas_emu: Optional[tuple[int, int]],
    rotation_deg: float,
    flip_h: bool,
    flip_v: bool,
) -> Optional[tuple[str, Optional[tuple[int, int, int, int]]]]:
    if flip_h:
        im = ImageOps.mirror(im)
    if flip_v:
        im = ImageOps.flip(im)
    if rotation_deg % 360.0:
        im = im.rotate(-rotation_deg, expand=True, resample=Image.BICUBIC)

    placed_box_emu = final_box_emu
    if final_box_emu and canvas_emu:
        im, placed_box_emu = _clip_raster_to_canvas(im, final_box_emu, canvas_emu)

    # Match the saved raster aspect ratio to the final placement box so rebuild
    # can place it 1:1 without an extra stretch step.
    target_box_emu = placed_box_emu or final_box_emu
    if target_box_emu is not None:
        target_w, target_h = _raster_size_2048_for_box(target_box_emu)
        if im.size != (target_w, target_h):
            im = im.resize((target_w, target_h), resample=Image.LANCZOS)

    try:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        im.save(out_path, format="PNG")
        return str(out_path), placed_box_emu
    except Exception:
        return None


def _fill_rect_pct(rect: Optional[Dict[str, int]]) -> Dict[str, float]:
    if not rect:
        return {"l": 0.0, "t": 0.0, "r": 0.0, "b": 0.0}
    return {k: float(rect.get(k, 0) or 0) / 100000.0 for k in ("l", "t", "r", "b")}


def _bake_fill_image_ops(
    im: Image.Image,
    src_rect: Optional[Dict[str, int]],
    fill_rect: Optional[Dict[str, int]],
    mask_oval: bool,
    mask_paths: Optional[List[Dict[str, Any]]] = None,
) -> Image.Image:
    s = _fill_rect_pct(src_rect)
    f = _fill_rect_pct(fill_rect)
    has_crop = not (all(abs(s[k]) < 1e-9 for k in s) and all(abs(f[k]) < 1e-9 for k in f))

    if has_crop:
        src_l = min(max(s["l"], 0.0), 1.0)
        src_t = min(max(s["t"], 0.0), 1.0)
        src_r = min(max(s["r"], 0.0), 1.0)
        src_b = min(max(s["b"], 0.0), 1.0)
        src_w = max(1e-9, 1.0 - src_l - src_r)
        src_h = max(1e-9, 1.0 - src_t - src_b)

        fill_l = f["l"]
        fill_t = f["t"]
        fill_r = f["r"]
        fill_b = f["b"]
        fill_w = max(1e-9, 1.0 - fill_l - fill_r)
        fill_h = max(1e-9, 1.0 - fill_t - fill_b)

        vis_u0 = max(0.0, min(1.0, (0.0 - fill_l) / fill_w))
        vis_v0 = max(0.0, min(1.0, (0.0 - fill_t) / fill_h))
        vis_u1 = max(0.0, min(1.0, (1.0 - fill_l) / fill_w))
        vis_v1 = max(0.0, min(1.0, (1.0 - fill_t) / fill_h))

        orig_u0 = src_l + vis_u0 * src_w
        orig_v0 = src_t + vis_v0 * src_h
        orig_u1 = src_l + vis_u1 * src_w
        orig_v1 = src_t + vis_v1 * src_h

        x0 = max(0, min(im.width, int(round(orig_u0 * im.width))))
        y0 = max(0, min(im.height, int(round(orig_v0 * im.height))))
        x1 = max(x0 + 1, min(im.width, int(round(orig_u1 * im.width))))
        y1 = max(y0 + 1, min(im.height, int(round(orig_v1 * im.height))))
        im = im.crop((x0, y0, x1, y1))

    if mask_oval:
        mask = Image.new("L", im.size, 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse((0, 0, im.size[0] - 1, im.size[1] - 1), fill=255)
        im.putalpha(mask)
    elif mask_paths:
        def _sample_cubic(p0, p1, p2, p3, steps: int = 48) -> List[tuple[float, float]]:
            pts: List[tuple[float, float]] = []
            for i in range(1, steps + 1):
                t = i / float(steps)
                mt = 1.0 - t
                x = (mt ** 3) * p0[0] + 3 * (mt ** 2) * t * p1[0] + 3 * mt * (t ** 2) * p2[0] + (t ** 3) * p3[0]
                y = (mt ** 3) * p0[1] + 3 * (mt ** 2) * t * p1[1] + 3 * mt * (t ** 2) * p2[1] + (t ** 3) * p3[1]
                pts.append((x, y))
            return pts

        scale = 4
        mask = Image.new("L", (max(1, im.size[0] * scale), max(1, im.size[1] * scale)), 0)
        draw = ImageDraw.Draw(mask)
        for path_spec in mask_paths:
            if not isinstance(path_spec, dict):
                continue
            commands = path_spec.get("commands") or []
            if not isinstance(commands, list) or not commands:
                continue
            poly: List[tuple[float, float]] = []
            start_pt: Optional[tuple[float, float]] = None
            cur: Optional[tuple[float, float]] = None
            for cmd in commands:
                op = str(cmd.get("op") or "").lower()
                pts = cmd.get("pts") or []
                if op == "moveto" and pts:
                    if len(poly) >= 3:
                        draw.polygon(poly, fill=255)
                    x = float(pts[0].get("x") or 0.0) * (mask.size[0] - 1)
                    y = float(pts[0].get("y") or 0.0) * (mask.size[1] - 1)
                    cur = (x, y)
                    start_pt = cur
                    poly = [cur]
                elif op == "lnto" and pts and cur is not None:
                    x = float(pts[0].get("x") or 0.0) * (mask.size[0] - 1)
                    y = float(pts[0].get("y") or 0.0) * (mask.size[1] - 1)
                    cur = (x, y)
                    poly.append(cur)
                elif op == "cubicbezto" and len(pts) == 3 and cur is not None:
                    p1 = (float(pts[0].get("x") or 0.0) * (mask.size[0] - 1), float(pts[0].get("y") or 0.0) * (mask.size[1] - 1))
                    p2 = (float(pts[1].get("x") or 0.0) * (mask.size[0] - 1), float(pts[1].get("y") or 0.0) * (mask.size[1] - 1))
                    p3 = (float(pts[2].get("x") or 0.0) * (mask.size[0] - 1), float(pts[2].get("y") or 0.0) * (mask.size[1] - 1))
                    poly.extend(_sample_cubic(cur, p1, p2, p3))
                    cur = p3
                elif op == "close":
                    if start_pt is not None and poly and poly[-1] != start_pt:
                        poly.append(start_pt)
                    if len(poly) >= 3:
                        draw.polygon(poly, fill=255)
                    poly = []
                    cur = start_pt
            if len(poly) >= 3:
                draw.polygon(poly, fill=255)
        if scale > 1:
            mask = mask.resize(im.size, resample=Image.LANCZOS)
        im = im.convert("RGBA")
        im.putalpha(mask)

    return im


def _apply_blip_alpha_mod_fix(im: Image.Image, alpha_amt: Optional[int]) -> Image.Image:
    """
    Apply DrawingML alpha modifications from a:alphaModFix/a:alpha to the rendered bitmap.

    DrawingML uses a 0..100000 scale (100000 == 100%).
    """
    if alpha_amt is None:
        return im
    try:
        amt = int(alpha_amt)
    except Exception:
        return im
    amt = max(0, min(100000, amt))
    if amt == 100000:
        return im
    if amt == 0:
        return Image.new("RGBA", im.size, (0, 0, 0, 0))
    try:
        im = im.convert("RGBA")
        r, g, b, a = im.split()
        # Multiply existing alpha by amt/100000.
        a = a.point(lambda v: int(round(v * (amt / 100000.0))))
        return Image.merge("RGBA", (r, g, b, a))
    except Exception:
        return im


def _slide_canvas_emu(slide) -> Optional[tuple[int, int]]:
    """Best-effort slide size in EMU."""
    try:
        pres = slide.part.package.presentation_part.presentation
        w = int(pres.slide_width)
        h = int(pres.slide_height)
        return w, h
    except Exception:
        return None


def _clip_box_emu_to_canvas(
    box_emu: Optional[tuple[int, int, int, int]], canvas_emu: Optional[tuple[int, int]]
) -> Optional[tuple[int, int, int, int]]:
    """Clamp an axis-aligned box in EMU to the slide canvas."""
    if not box_emu or not canvas_emu:
        return box_emu
    cw, ch = canvas_emu
    l, t, w, h = box_emu
    if w <= 0 or h <= 0:
        return None
    vis_l = max(0, l)
    vis_t = max(0, t)
    vis_r = min(cw, l + w)
    vis_b = min(ch, t + h)
    if vis_r <= vis_l or vis_b <= vis_t:
        return None
    return vis_l, vis_t, vis_r - vis_l, vis_b - vis_t


def _clamp_line_box_emu_to_canvas(
    box_emu: Optional[tuple[int, int, int, int]], canvas_emu: Optional[tuple[int, int]]
) -> Optional[tuple[int, int, int, int]]:
    """Clamp a line/connector box in EMU to the slide canvas.

    Connector boxes commonly have width==0 (vertical line) or height==0
    (horizontal line). Those are valid and must not be discarded.
    """
    if not box_emu or not canvas_emu:
        return box_emu
    cw, ch = canvas_emu
    l, t, w, h = box_emu
    if w < 0 or h < 0 or (w == 0 and h == 0):
        return None

    # Vertical line
    if w == 0:
        if l < 0 or l > cw:
            return None
        vis_t = max(0, t)
        vis_b = min(ch, t + h)
        if vis_b <= vis_t:
            return None
        return l, vis_t, 0, vis_b - vis_t

    # Horizontal line
    if h == 0:
        if t < 0 or t > ch:
            return None
        vis_l = max(0, l)
        vis_r = min(cw, l + w)
        if vis_r <= vis_l:
            return None
        return vis_l, t, vis_r - vis_l, 0

    return _clip_box_emu_to_canvas(box_emu, canvas_emu)


def _raster_size_2048_for_box(box_emu: Optional[tuple[int, int, int, int]]) -> tuple[int, int]:
    # Keep the raster area close to RASTER_TARGET_SIDE^2 while preserving aspect ratio.
    if not box_emu:
        return RASTER_TARGET_SIDE, RASTER_TARGET_SIDE
    _, _, w, h = box_emu
    if w <= 0 or h <= 0:
        return RASTER_TARGET_SIDE, RASTER_TARGET_SIDE
    target_area = float(RASTER_TARGET_SIDE * RASTER_TARGET_SIDE)
    aspect = float(w) / float(h)
    out_w = max(1, int(round(math.sqrt(target_area * aspect))))
    out_h = max(1, int(round(target_area / float(out_w))))
    max_side = max(1, RASTER_TARGET_SIDE * 2)
    if out_w > max_side:
        out_w = max_side
        out_h = max(1, int(round(float(out_w) / aspect)))
    if out_h > max_side:
        out_h = max_side
        out_w = max(1, int(round(float(out_h) * aspect)))
    return out_w, out_h


def _clip_raster_to_canvas(
    im: Image.Image, box_emu: Optional[tuple[int, int, int, int]], canvas_emu: Optional[tuple[int, int]]
) -> tuple[Image.Image, Optional[tuple[int, int, int, int]]]:
    if not box_emu or not canvas_emu:
        return im, box_emu
    l, t, w, h = box_emu
    cw, ch = canvas_emu
    # Compute intersection between [l,l+w]x[t,t+h] and slide canvas [0,cw]x[0,ch]
    vis_l = max(0, l)
    vis_t = max(0, t)
    vis_r = min(cw, l + w)
    vis_b = min(ch, t + h)
    if vis_r <= vis_l or vis_b <= vis_t:
        # Fully off-canvas; return a 1x1 transparent pixel to avoid crashes downstream.
        return Image.new("RGBA", (1, 1), (0, 0, 0, 0)), None
    # Map visible region into raster pixel space.
    x0 = int(round((vis_l - l) / float(w) * im.width))
    y0 = int(round((vis_t - t) / float(h) * im.height))
    x1 = int(round((vis_r - l) / float(w) * im.width))
    y1 = int(round((vis_b - t) / float(h) * im.height))
    x0 = max(0, min(im.width - 1, x0))
    y0 = max(0, min(im.height - 1, y0))
    x1 = max(x0 + 1, min(im.width, x1))
    y1 = max(y0 + 1, min(im.height, y1))
    cropped = im.crop((x0, y0, x1, y1))
    clipped_box = (vis_l, vis_t, vis_r - vis_l, vis_b - vis_t)
    return cropped, clipped_box


def extract_picture_asset(
    slide, shp, slide_index: int, assets_dir: Path, transform: Optional[Dict[str, float]] = None
) -> Dict[str, Any]:
    shape_id = getattr(shp, "shape_id", None) or 0
    base = f"slide{slide_index:03d}_shape{shape_id}"
    blip = None
    try:
        blip = shp._element.find(".//a:blip", NS)
    except Exception:
        blip = None
    part = None
    if blip is not None:
        part, _ = _resolve_image_parts(slide, blip)
    if part is not None:
        blob = part.blob
        ext = (getattr(part, "ext", "") or "").lower()
        ctype = getattr(part, "content_type", None)
    else:
        img = shp.image
        blob = img.blob
        ext = (img.ext or "").lower()
        ctype = getattr(img, "content_type", None)
    canvas_emu = _slide_canvas_emu(slide)
    base_box_emu = _shape_bounds_box_emu(shp, transform)
    final_box_emu = _shape_box_emu(shp, transform) or base_box_emu
    flips = shape_flip_flags(shp)
    rotation_deg = shape_rotation_deg(shp, transform)
    source_is_svg = (ext == "svg") or (ctype == "image/svg+xml")
    rendered_kind = "svg_image_png_rgba" if source_is_svg else "image_png_rgba"
    src_rect = None
    fill_rect = None
    alpha_amt = None

    def _rect_from_el(rect_el) -> Optional[Dict[str, int]]:
        if rect_el is None:
            return None
        out_rect: Dict[str, int] = {}
        found = False
        for key in ("l", "t", "r", "b"):
            val = rect_el.get(key)
            if val is None:
                out_rect[key] = 0
                continue
            try:
                out_rect[key] = int(val)
                found = True
            except Exception:
                out_rect[key] = 0
        return out_rect if found else None

    if blip is not None:
        try:
            parent = blip.getparent()
            blip_fill = None
            while parent is not None:
                if parent.tag == qn("p:blipFill"):
                    blip_fill = parent
                    break
                parent = parent.getparent()
            if blip_fill is not None:
                src_rect = _rect_from_el(blip_fill.find(qn("a:srcRect")))
                stretch = blip_fill.find(qn("a:stretch"))
                if stretch is not None:
                    fill_rect = _rect_from_el(stretch.find(qn("a:fillRect")))
        except Exception:
            src_rect = None
            fill_rect = None
        try:
            amf = blip.find(qn("a:alphaModFix"))
            if amf is not None and amf.get("amt") is not None:
                alpha_amt = int(amf.get("amt"))
            else:
                ael = blip.find(qn("a:alpha"))
                if ael is not None and ael.get("val") is not None:
                    alpha_amt = int(ael.get("val"))
        except Exception:
            alpha_amt = None

    has_picture_crop = (src_rect is not None) or (fill_rect is not None)
    if has_picture_crop:
        im = _load_fill_blob_as_rgba(blob=blob, ext=ext, content_type=ctype)
        if im is not None:
            im = _bake_fill_image_ops(im, src_rect, fill_rect, False)
            im = _apply_blip_alpha_mod_fix(im, alpha_amt)
    else:
        im = _load_blob_as_rgba(blob=blob, ext=ext, content_type=ctype, box_emu=base_box_emu or final_box_emu)
        if im is not None:
            im = _apply_blip_alpha_mod_fix(im, alpha_amt)
    if im is not None and final_box_emu is not None:
        rendered = _save_final_image_layer(
            im=im,
            out_path=(assets_dir / f"{base}_{rendered_kind}_rendered.png").resolve(),
            final_box_emu=final_box_emu,
            canvas_emu=canvas_emu,
            rotation_deg=rotation_deg,
            flip_h=bool(flips.get("flip_h")),
            flip_v=bool(flips.get("flip_v")),
        )
        if rendered is not None:
            rendered_path, rendered_box_emu = rendered
            placed_box_emu = rendered_box_emu or final_box_emu
            return {
                "slide": slide_index,
                "shape_name": getattr(shp, "name", None),
                "kind": rendered_kind,
                "box": _ratio_box(*placed_box_emu, canvas_emu),
                "saved_path": rendered_path,
            }

    exported = export_asset_image(blob=blob, ext=ext, content_type=ctype, out_base=base, assets_dir=assets_dir)
    return {
        "slide": slide_index,
        "shape_name": getattr(shp, "name", None),
        "kind": exported["kind"],
        "box": shape_box(shp, transform, canvas_emu),
        "saved_path": exported["saved_path"],
    }


def extract_fill_image_assets(
    slide, shp, slide_index: int, assets_dir: Path, transform: Optional[Dict[str, float]] = None
) -> List[Dict[str, Any]]:
    """
    Extract images used as shape fills (and similar) by scanning for a:blip/@r:embed in the shape XML.
    Many templates use picture-fill shapes instead of PICTURE shapes.
    """
    el = getattr(shp, "_element", None)
    if el is None or not hasattr(el, "xpath"):
        return []

    out: List[Dict[str, Any]] = []
    # python-pptx oxml elements provide .xpath() with built-in namespace mapping
    blips = el.xpath(".//a:blip[@r:embed]")
    for idx, blip in enumerate(blips, start=1):
        part, preview_part = _resolve_image_parts(slide, blip)
        if part is None:
            continue
        ctype = getattr(part, "content_type", None)

        blip_fill = None
        try:
            parent = blip.getparent()
            while parent is not None:
                if parent.tag == qn("a:blipFill"):
                    blip_fill = parent
                    break
                parent = parent.getparent()
        except Exception:
            blip_fill = None

        def _rect_from_el(rect_el) -> Optional[Dict[str, int]]:
            if rect_el is None:
                return None
            out_rect: Dict[str, int] = {}
            found = False
            for key in ("l", "t", "r", "b"):
                val = rect_el.get(key)
                if val is None:
                    out_rect[key] = 0
                    continue
                try:
                    out_rect[key] = int(val)
                    found = True
                except Exception:
                    out_rect[key] = 0
            return out_rect if found else None

        src_rect = None
        fill_rect = None
        alpha_amt = None
        if blip_fill is not None:
            src_rect = _rect_from_el(blip_fill.find(qn("a:srcRect")))
            stretch = blip_fill.find(qn("a:stretch"))
            if stretch is not None:
                fill_rect = _rect_from_el(stretch.find(qn("a:fillRect")))
        # Some templates apply transparency to picture fills via a:alphaModFix/a:alpha on the blip.
        try:
            amf = blip.find(qn("a:alphaModFix"))
            if amf is not None and amf.get("amt") is not None:
                alpha_amt = int(amf.get("amt"))
            else:
                ael = blip.find(qn("a:alpha"))
                if ael is not None and ael.get("val") is not None:
                    alpha_amt = int(ael.get("val"))
        except Exception:
            alpha_amt = None

        shape_id = getattr(shp, "shape_id", None) or 0
        base = f"slide{slide_index:03d}_shape{shape_id}_fill{idx}"
        ext = ""
        try:
            ext = (getattr(part, "ext", "") or "").lower()
        except Exception:
            ext = ""

        mask_oval = False
        mask_paths = None
        try:
            mask_oval = (
                getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE
                and getattr(shp, "auto_shape_type", None) == MSO_AUTO_SHAPE_TYPE.OVAL
            )
        except Exception:
            mask_oval = False

        # Many templates use a custom-geometry "freeform circle" instead of a preset oval.
        # Detect the common cubic-bezier circle approximation:
        # - a:custGeom -> a:path w==h
        # - moveTo at (w/2, 0)
        # - 4 cubicBezTo segments
        def _is_custgeom_circle_like() -> bool:
            try:
                shp_el = getattr(shp, "_element", None)
                if shp_el is None:
                    return False
                root = ET.fromstring(ET.tostring(shp_el, encoding="utf-8"))
                path = root.find(".//a:custGeom/a:pathLst/a:path", NS)
                if path is None:
                    return False
                w = int(path.get("w") or "0")
                h = int(path.get("h") or "0")
                tol = max(2, max(w, h) // 500)
                if w <= 0 or h <= 0 or abs(w - h) > tol:
                    return False
                move = path.find("a:moveTo/a:pt", NS)
                if move is None:
                    return False
                mx = int(move.get("x") or "0")
                my = int(move.get("y") or "0")
                # Circle paths may start at any cardinal midpoint:
                # top-center, right-center, bottom-center, or left-center.
                cardinal_midpoints = (
                    (w // 2, 0),
                    (w, h // 2),
                    (w // 2, h),
                    (0, h // 2),
                )
                if not any(abs(mx - ex) <= tol and abs(my - ey) <= tol for ex, ey in cardinal_midpoints):
                    return False
                cubics = path.findall("a:cubicBezTo", NS)
                if len(cubics) != 4:
                    return False
                close = path.find("a:close", NS)
                if close is None:
                    return False
                return True
            except Exception:
                return False

        def _is_custgeom_rect_like() -> bool:
            try:
                shp_el = getattr(shp, "_element", None)
                if shp_el is None:
                    return False
                root = ET.fromstring(ET.tostring(shp_el, encoding="utf-8"))
                path = root.find(".//a:custGeom/a:pathLst/a:path", NS)
                if path is None:
                    return False
                w = int(path.get("w") or "0")
                h = int(path.get("h") or "0")
                if w <= 0 or h <= 0:
                    return False
                tol = max(2, max(w, h) // 500)
                cmds = list(path)
                ops = [c.tag.rsplit("}", 1)[-1] for c in cmds]
                if ops not in (["moveTo", "lnTo", "lnTo", "lnTo", "lnTo", "close"], ["moveTo", "lnTo", "lnTo", "lnTo", "close"]):
                    return False
                pts = []
                for child in cmds:
                    tag = child.tag.rsplit("}", 1)[-1]
                    if tag not in {"moveTo", "lnTo"}:
                        continue
                    pt = child.find("a:pt", NS)
                    if pt is None:
                        return False
                    pts.append((int(pt.get("x") or "0"), int(pt.get("y") or "0")))
                allowed = {(0, 0), (w, 0), (w, h), (0, h)}
                for x, y in pts:
                    if not any(abs(x - ax) <= tol and abs(y - ay) <= tol for ax, ay in allowed):
                        return False
                return True
            except Exception:
                return False

        if not mask_oval and _is_custgeom_circle_like():
            mask_oval = True

        if not mask_oval:
            try:
                shp_el = getattr(shp, "_element", None)
                if shp_el is not None:
                    root = ET.fromstring(ET.tostring(shp_el, encoding="utf-8"))
                    cust = _cust_geom_paths_spec_from_root(root)
                    if (not _is_custgeom_rect_like()) and isinstance(cust, dict) and isinstance(cust.get("paths"), list) and cust.get("paths"):
                        mask_paths = cust.get("paths")
            except Exception:
                mask_paths = None

        canvas_emu = _slide_canvas_emu(slide)
        base_box_emu = _shape_bounds_box_emu(shp, transform)
        final_box_emu = _shape_box_emu(shp, transform) or base_box_emu
        flips = shape_flip_flags(shp)
        rotation_deg = shape_rotation_deg(shp, transform)
        source_is_svg = (ext == "svg") or (ctype == "image/svg+xml")
        rendered_kind = "svg_image_png_rgba" if source_is_svg else "image_png_rgba"

        im = _load_fill_blob_as_rgba(blob=part.blob, ext=ext, content_type=ctype)
        if im is None and preview_part is not None:
            preview_ext = (getattr(preview_part, "ext", "") or "").lower()
            preview_ctype = getattr(preview_part, "content_type", None)
            im = _load_fill_blob_as_rgba(blob=preview_part.blob, ext=preview_ext, content_type=preview_ctype)

        if im is not None and final_box_emu is not None:
            im = _bake_fill_image_ops(im, src_rect, fill_rect, mask_oval, mask_paths)
            im = _apply_blip_alpha_mod_fix(im, alpha_amt)
            rendered = _save_final_image_layer(
                im=im,
                out_path=(assets_dir / f"{base}_{rendered_kind}_rendered.png").resolve(),
                final_box_emu=final_box_emu,
                canvas_emu=canvas_emu,
                rotation_deg=rotation_deg,
                flip_h=bool(flips.get("flip_h")),
                flip_v=bool(flips.get("flip_v")),
            )
            if rendered is not None:
                rendered_path, rendered_box_emu = rendered
                placed_box_emu = rendered_box_emu or final_box_emu
                out.append(
                    {
                        "slide": slide_index,
                        "shape_name": getattr(shp, "name", None),
                        "kind": rendered_kind,
                        "box": _ratio_box(*placed_box_emu, canvas_emu),
                        "saved_path": rendered_path,
                    }
                )
                continue

        exported = export_asset_image(blob=part.blob, ext=ext, content_type=ctype, out_base=base, assets_dir=assets_dir)
        out.append(
            {
                "slide": slide_index,
                "shape_name": getattr(shp, "name", None),
                "kind": exported["kind"],
                "box": shape_box(shp, transform, canvas_emu),
                "saved_path": exported["saved_path"],
            }
        )
    return out


def _shape_has_blip_fill(shp) -> bool:
    try:
        return getattr(shp, "_element", None).find(".//a:blip", NS) is not None
    except Exception:
        return False


def _remove_direct_child(root: ET.Element, tag: str) -> None:
    for child in list(root):
        if child.tag == tag:
            root.remove(child)


def _shape_xml_without_text(shp, transform: Optional[Dict[str, float]] = None) -> Optional[str]:
    try:
        root = ET.fromstring(ET.tostring(shp._element, encoding="utf-8"))
    except Exception:
        return None

    # Text is exported separately; remove it from the geometry XML to avoid duplication on rebuild.
    _remove_direct_child(root, f"{{{NS['p']}}}txBody")

    xfrm = root.find(".//a:xfrm", NS)
    if xfrm is not None:
        off = xfrm.find("a:off", NS)
        ext = xfrm.find("a:ext", NS)
        bounds_box = _shape_bounds_box_emu(shp, transform)
        if bounds_box is not None:
            left, top, width, height = bounds_box
        else:
            left = top = width = height = None
        if off is not None:
            if left is not None:
                off.set("x", str(int(left)))
            if top is not None:
                off.set("y", str(int(top)))
        if ext is not None:
            if width is not None:
                ext.set("cx", str(int(width)))
            if height is not None:
                ext.set("cy", str(int(height)))
        try:
            rot = int(round(shape_rotation_deg(shp, transform) * 60000))
            if rot:
                xfrm.set("rot", str(rot))
            elif "rot" in xfrm.attrib:
                del xfrm.attrib["rot"]
        except Exception:
            pass

    return ET.tostring(root, encoding="unicode")


def _color_spec_from_color_el(color_el: ET.Element) -> Optional[Dict[str, Any]]:
    tag = color_el.tag.rsplit("}", 1)[-1]
    if tag == "srgbClr" and color_el.get("val"):
        spec: Dict[str, Any] = {"type": "srgb", "rgb": str(color_el.get("val")).upper()}
    elif tag == "schemeClr" and color_el.get("val"):
        spec = {"type": "scheme", "scheme": str(color_el.get("val"))}
    else:
        return None
    mods: List[Dict[str, Any]] = []
    for child in list(color_el):
        tag = child.tag.rsplit("}", 1)[-1]
        if tag in {"tint", "shade", "lumMod", "lumOff", "alpha"} and child.get("val"):
            try:
                mods.append({"op": tag, "val": int(child.get("val"))})
            except Exception:
                mods.append({"op": tag, "val": child.get("val")})
    if mods:
        spec["mods"] = mods
    return spec


def _solid_fill_spec_from_node(node: ET.Element) -> Optional[Dict[str, Any]]:
    solid = node.find("a:solidFill", NS)
    if solid is None and node.tag.rsplit("}", 1)[-1] == "solidFill":
        solid = node
    if solid is None:
        return None
    srgb = solid.find("a:srgbClr", NS)
    if srgb is not None and srgb.get("val"):
        return _color_spec_from_color_el(srgb)
    scheme = solid.find("a:schemeClr", NS)
    if scheme is not None and scheme.get("val"):
        return _color_spec_from_color_el(scheme)
    return None


def _background_fill_from_color_spec(color_spec: Optional[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    if not isinstance(color_spec, dict):
        return None
    out: Dict[str, Any] = {"type": "solid"}
    color_type = str(color_spec.get("type") or "").lower()
    if color_type in {"rgb", "srgb"} and color_spec.get("rgb"):
        out.update({"color_type": "RGB", "rgb": str(color_spec.get("rgb")).upper()})
    elif color_type == "scheme" and color_spec.get("scheme"):
        out.update({"color_type": "SCHEME", "scheme": color_spec.get("scheme")})
    if color_spec.get("mods"):
        out["mods"] = list(color_spec.get("mods") or [])
    return out


def _gradient_fill_spec_from_node(node: ET.Element) -> Optional[Dict[str, Any]]:
    grad = node.find("a:gradFill", NS)
    if grad is None and node.tag.rsplit("}", 1)[-1] == "gradFill":
        grad = node
    if grad is None:
        return None

    out: Dict[str, Any] = {"type": "gradient"}
    stops: List[Dict[str, Any]] = []
    gs_lst = grad.find("a:gsLst", NS)
    if gs_lst is not None:
        for gs in gs_lst.findall("a:gs", NS):
            stop: Dict[str, Any] = {}
            if gs.get("pos") is not None:
                try:
                    stop["pos"] = int(gs.get("pos"))
                except Exception:
                    stop["pos"] = gs.get("pos")

            color_spec: Optional[Dict[str, Any]] = None
            for child in list(gs):
                color_spec = _color_spec_from_color_el(child)
                if color_spec is not None:
                    break

            if color_spec is not None:
                color_type = str(color_spec.get("type") or "").lower()
                if color_type in {"rgb", "srgb"} and color_spec.get("rgb"):
                    stop.update({"color_type": "RGB", "rgb": str(color_spec.get("rgb")).upper()})
                elif color_type == "scheme" and color_spec.get("scheme"):
                    stop.update({"color_type": "SCHEME", "scheme": color_spec.get("scheme")})
                if color_spec.get("mods"):
                    stop["mods"] = list(color_spec.get("mods") or [])
            if stop:
                stops.append(stop)

    if stops:
        out["stops"] = stops
        first_stop = next((stop for stop in stops if isinstance(stop, dict)), None)
        if isinstance(first_stop, dict):
            if first_stop.get("color_type") == "RGB" and first_stop.get("rgb"):
                out.update({"color_type": "RGB", "rgb": first_stop.get("rgb")})
            elif first_stop.get("color_type") == "SCHEME" and first_stop.get("scheme"):
                out.update({"color_type": "SCHEME", "scheme": first_stop.get("scheme")})
            if first_stop.get("mods"):
                out["mods"] = list(first_stop.get("mods") or [])

    lin = grad.find("a:lin", NS)
    if lin is not None and lin.get("ang") is not None:
        try:
            out["angle"] = int(lin.get("ang"))
        except Exception:
            out["angle"] = lin.get("ang")

    path = grad.find("a:path", NS)
    if path is not None and path.get("path"):
        out["path"] = path.get("path")
        fill_to_rect = path.find("a:fillToRect", NS)
        if fill_to_rect is not None:
            rect_spec: Dict[str, int] = {}
            for key in ("l", "t", "r", "b"):
                if fill_to_rect.get(key) is None:
                    continue
                try:
                    rect_spec[key] = int(fill_to_rect.get(key))
                except Exception:
                    continue
            if rect_spec:
                out["fill_to_rect"] = rect_spec

    return out if len(out) > 1 else None


def _gradient_stop_rgba(stop: Dict[str, Any]) -> Optional[tuple[int, int, int, float]]:
    if not isinstance(stop, dict):
        return None
    spec: Dict[str, Any] = {"mods": list(stop.get("mods") or [])}
    color_type = str(stop.get("color_type") or "").upper()
    if color_type == "RGB" and stop.get("rgb"):
        spec.update({"type": "srgb", "rgb": stop.get("rgb")})
    elif color_type == "SCHEME" and stop.get("scheme"):
        spec.update({"type": "scheme", "scheme": stop.get("scheme")})
    else:
        return None
    rgba = _rgba_from_color_spec(spec)
    if rgba is None:
        return None
    r, g, b, alpha = rgba
    return r, g, b, max(0.0, min(1.0, float(alpha) / 255.0))


def _gradient_fill_to_svg(
    gradient_spec: Dict[str, Any], out_w: int, out_h: int
) -> Optional[str]:
    if out_w <= 0 or out_h <= 0:
        return None

    stops_svg: List[str] = []
    for stop in gradient_spec.get("stops") or []:
        rgba = _gradient_stop_rgba(stop)
        if rgba is None:
            continue
        r, g, b, alpha = rgba
        try:
            pos = max(0.0, min(1.0, float(stop.get("pos", 0)) / 100000.0))
        except Exception:
            pos = 0.0
        stops_svg.append(
            f'<stop offset="{pos * 100:.4f}%" stop-color="rgb({r},{g},{b})" stop-opacity="{alpha:.6f}"/>'
        )
    if not stops_svg:
        return None

    gradient_markup = ""
    if str(gradient_spec.get("path") or "").lower() == "circle":
        gradient_markup = (
            '<radialGradient id="bggrad" cx="50%" cy="50%" r="70.710678%" '
            'fx="50%" fy="50%">'
            + "".join(stops_svg)
            + "</radialGradient>"
        )
    else:
        try:
            angle_deg = float(gradient_spec.get("angle") or 0.0) / 60000.0
        except Exception:
            angle_deg = 0.0
        theta = math.radians(angle_deg)
        dx = math.cos(theta)
        dy = -math.sin(theta)
        x1 = 50.0 - (dx * 50.0)
        y1 = 50.0 - (dy * 50.0)
        x2 = 50.0 + (dx * 50.0)
        y2 = 50.0 + (dy * 50.0)
        gradient_markup = (
            f'<linearGradient id="bggrad" x1="{x1:.4f}%" y1="{y1:.4f}%" '
            f'x2="{x2:.4f}%" y2="{y2:.4f}%">'
            + "".join(stops_svg)
            + "</linearGradient>"
        )

    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{out_w}" height="{out_h}" '
        f'viewBox="0 0 {out_w} {out_h}">'
        "<defs>"
        f"{gradient_markup}"
        "</defs>"
        f'<rect x="0" y="0" width="{out_w}" height="{out_h}" fill="url(#bggrad)"/>'
        "</svg>"
    )


def _interpolate_rgba_stops(
    stops: List[tuple[float, tuple[int, int, int, int]]],
    t: float,
) -> tuple[int, int, int, int]:
    if not stops:
        return (0, 0, 0, 0)
    t = max(0.0, min(1.0, float(t)))
    if t <= stops[0][0]:
        return stops[0][1]
    if t >= stops[-1][0]:
        return stops[-1][1]
    for idx in range(1, len(stops)):
        p0, c0 = stops[idx - 1]
        p1, c1 = stops[idx]
        if t > p1:
            continue
        span = p1 - p0
        ratio = 0.0 if span <= 0 else (t - p0) / span
        return tuple(
            int(round((float(c0[i]) * (1.0 - ratio)) + (float(c1[i]) * ratio)))
            for i in range(4)
        )
    return stops[-1][1]


def _paint_gradient_fill_rect(
    im: Image.Image,
    rect: tuple[int, int, int, int],
    gradient_spec: Dict[str, Any],
) -> None:
    x0, y0, x1, y1 = rect
    width = max(0, int(x1) - int(x0))
    height = max(0, int(y1) - int(y0))
    if width <= 0 or height <= 0:
        return
    stops: List[tuple[float, tuple[int, int, int, int]]] = []
    for stop in gradient_spec.get("stops") or []:
        rgba = _gradient_stop_rgba(stop)
        if rgba is None:
            continue
        r, g, b, alpha = rgba
        try:
            pos = max(0.0, min(1.0, float(stop.get("pos", 0)) / 100000.0))
        except Exception:
            pos = 0.0
        stops.append((pos, (r, g, b, int(round(alpha * 255.0)))))
    if not stops:
        return
    stops.sort(key=lambda item: item[0])
    if len(stops) == 1:
        tile = Image.new("RGBA", (width, height), stops[0][1])
        im.alpha_composite(tile, (int(x0), int(y0)))
        return

    if str(gradient_spec.get("path") or "").lower() == "circle":
        fill_to_rect = gradient_spec.get("fill_to_rect") if isinstance(gradient_spec, dict) else None
        focus_left = 0.5 * float(width)
        focus_right = focus_left
        focus_top = 0.5 * float(height)
        focus_bottom = focus_top
        if isinstance(fill_to_rect, dict):
            try:
                focus_left = (float(fill_to_rect.get("l", 0.0)) / 100000.0) * float(width)
                focus_right = float(width) - ((float(fill_to_rect.get("r", 0.0)) / 100000.0) * float(width))
                focus_top = (float(fill_to_rect.get("t", 0.0)) / 100000.0) * float(height)
                focus_bottom = float(height) - ((float(fill_to_rect.get("b", 0.0)) / 100000.0) * float(height))
            except Exception:
                focus_left = focus_right = 0.5 * float(width)
                focus_top = focus_bottom = 0.5 * float(height)

        left_edge = min(focus_left, focus_right)
        right_edge = max(focus_left, focus_right)
        top_edge = min(focus_top, focus_bottom)
        bottom_edge = max(focus_top, focus_bottom)

        pixels: List[tuple[int, int, int, int]] = []
        for py in range(height):
            fy = min(max(float(py), top_edge), bottom_edge)
            for px in range(width):
                fx = min(max(float(px), left_edge), right_edge)
                dx = float(px) - fx
                dy = float(py) - fy
                if abs(dx) < 1e-9 and abs(dy) < 1e-9:
                    pixels.append(_interpolate_rgba_stops(stops, 0.0))
                    continue
                scales: List[float] = []
                if dx > 0:
                    scales.append((float(width - 1) - fx) / dx)
                elif dx < 0:
                    scales.append((0.0 - fx) / dx)
                if dy > 0:
                    scales.append((float(height - 1) - fy) / dy)
                elif dy < 0:
                    scales.append((0.0 - fy) / dy)
                scales = [val for val in scales if val > 0]
                if not scales:
                    pixels.append(_interpolate_rgba_stops(stops, 1.0))
                    continue
                scale = min(scales)
                boundary_dist = math.hypot(dx * scale, dy * scale)
                if boundary_dist <= 1e-9:
                    pixels.append(_interpolate_rgba_stops(stops, 1.0))
                    continue
                t = math.hypot(dx, dy) / boundary_dist
                pixels.append(_interpolate_rgba_stops(stops, t))
        tile = Image.new("RGBA", (width, height))
        tile.putdata(pixels)
        im.alpha_composite(tile, (int(x0), int(y0)))
        return

    try:
        angle_deg = float(gradient_spec.get("angle") or 0.0) / 60000.0
    except Exception:
        angle_deg = 0.0
    theta = math.radians(angle_deg)
    dx = math.cos(theta)
    dy = -math.sin(theta)
    scale = max(1.0, ((width - 1) * abs(dx)) + ((height - 1) * abs(dy)))
    cx = (width - 1) / 2.0
    cy = (height - 1) / 2.0
    pixels: List[tuple[int, int, int, int]] = []
    for py in range(height):
        ry = py - cy
        for px in range(width):
            rx = px - cx
            t = 0.5 + (((rx * dx) + (ry * dy)) / scale)
            pixels.append(_interpolate_rgba_stops(stops, t))
    tile = Image.new("RGBA", (width, height))
    tile.putdata(pixels)
    im.alpha_composite(tile, (int(x0), int(y0)))


def _svg_aspect_ratio(svg_bytes: bytes) -> Optional[float]:
    text = svg_bytes.decode("utf-8", errors="replace")
    m = re.search(r'viewBox="([^"]+)"', text)
    if m:
        try:
            _, _, w, h = [float(x) for x in m.group(1).replace(",", " ").split()]
            if w > 0 and h > 0:
                return w / h
        except Exception:
            pass
    mw = re.search(r'width="([0-9.]+)', text)
    mh = re.search(r'height="([0-9.]+)', text)
    try:
        if mw and mh:
            w = float(mw.group(1))
            h = float(mh.group(1))
            if w > 0 and h > 0:
                return w / h
    except Exception:
        pass
    return None


def _svg_bytes_to_png_via_rsvg_convert(svg_bytes: bytes, out_w: int, out_h: int) -> Optional[bytes]:
    exe = shutil.which("rsvg-convert")
    if not exe:
        return None
    try:
        proc = subprocess.run(
            [
                exe,
                "--format",
                "png",
                "--width",
                str(max(1, int(out_w))),
                "--height",
                str(max(1, int(out_h))),
            ],
            input=svg_bytes,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=False,
        )
    except Exception:
        return None
    if proc.returncode != 0:
        return None
    return proc.stdout or None


def _svg_bytes_to_png_via_cairosvg(svg_bytes: bytes, out_w: int, out_h: int) -> Optional[bytes]:
    try:
        import cairosvg  # type: ignore
    except Exception:
        return None
    try:
        return cairosvg.svg2png(
            bytestring=svg_bytes,
            output_width=max(1, int(out_w)),
            output_height=max(1, int(out_h)),
        )
    except Exception:
        return None


def _svg_bytes_to_png_bytes(svg_bytes: bytes, out_w: int, out_h: int) -> Optional[bytes]:
    png_bytes = _svg_bytes_to_png_via_rsvg_convert(svg_bytes, out_w, out_h)
    if png_bytes:
        return png_bytes
    return _svg_bytes_to_png_via_cairosvg(svg_bytes, out_w, out_h)


def _sample_cubic_points(
    p0: tuple[float, float],
    p1: tuple[float, float],
    p2: tuple[float, float],
    p3: tuple[float, float],
    *,
    steps: int = 24,
) -> List[tuple[float, float]]:
    pts: List[tuple[float, float]] = []
    for i in range(1, max(1, steps) + 1):
        t = i / float(max(1, steps))
        mt = 1.0 - t
        x = (mt ** 3) * p0[0] + 3 * (mt ** 2) * t * p1[0] + 3 * mt * (t ** 2) * p2[0] + (t ** 3) * p3[0]
        y = (mt ** 3) * p0[1] + 3 * (mt ** 2) * t * p1[1] + 3 * mt * (t ** 2) * p2[1] + (t ** 3) * p3[1]
        pts.append((x, y))
    return pts


def _line_width_px_for_box(line: Optional[Dict[str, Any]], box_emu: tuple[int, int, int, int], out_w: int, out_h: int) -> int:
    try:
        width_emu = max(float((line or {}).get("width_pt") or 0.75) * EMU_PER_PT, 1.0)
    except Exception:
        width_emu = 0.75 * EMU_PER_PT
    box_w = max(1.0, float(box_emu[2]))
    box_h = max(1.0, float(box_emu[3]))
    scale = min(float(out_w) / box_w, float(out_h) / box_h)
    return max(1, int(round(width_emu * scale)))


def _draw_poly_commands(
    draw: ImageDraw.ImageDraw,
    commands: List[Dict[str, Any]],
    *,
    op_key: str,
    fill_rgba: Optional[tuple[int, int, int, int]],
    stroke_rgba: Optional[tuple[int, int, int, int]],
    stroke_width: int,
    out_w: int,
    out_h: int,
) -> bool:
    drew = False
    poly: List[tuple[float, float]] = []
    start_pt: Optional[tuple[float, float]] = None
    cur: Optional[tuple[float, float]] = None
    closed = False
    for cmd in commands:
        if not isinstance(cmd, dict):
            continue
        op = str(cmd.get(op_key) or "").lower()
        pts = cmd.get("pts") or []
        if op == "moveto" and pts:
            if len(poly) >= 2 and stroke_rgba is not None:
                draw.line(poly, fill=stroke_rgba, width=stroke_width)
                drew = True
            poly = []
            x = float(pts[0].get("x") or 0.0) * float(out_w - 1)
            y = float(pts[0].get("y") or 0.0) * float(out_h - 1)
            cur = (x, y)
            start_pt = cur
            poly = [cur]
            closed = False
        elif op == "lnto" and pts and cur is not None:
            x = float(pts[0].get("x") or 0.0) * float(out_w - 1)
            y = float(pts[0].get("y") or 0.0) * float(out_h - 1)
            cur = (x, y)
            poly.append(cur)
        elif op == "cubicbezto" and len(pts) == 3 and cur is not None:
            p1 = (float(pts[0].get("x") or 0.0) * float(out_w - 1), float(pts[0].get("y") or 0.0) * float(out_h - 1))
            p2 = (float(pts[1].get("x") or 0.0) * float(out_w - 1), float(pts[1].get("y") or 0.0) * float(out_h - 1))
            p3 = (float(pts[2].get("x") or 0.0) * float(out_w - 1), float(pts[2].get("y") or 0.0) * float(out_h - 1))
            poly.extend(_sample_cubic_points(cur, p1, p2, p3))
            cur = p3
        elif op == "close":
            if start_pt is not None and poly and poly[-1] != start_pt:
                poly.append(start_pt)
            closed = True
            if len(poly) >= 3 and fill_rgba is not None:
                draw.polygon(poly, fill=fill_rgba)
                drew = True
            if len(poly) >= 2 and stroke_rgba is not None:
                draw.line(poly, fill=stroke_rgba, width=stroke_width)
                drew = True
            poly = []
            cur = start_pt
    if poly:
        if closed and len(poly) >= 3 and fill_rgba is not None:
            draw.polygon(poly, fill=fill_rgba)
            drew = True
        if len(poly) >= 2 and stroke_rgba is not None:
            draw.line(poly, fill=stroke_rgba, width=stroke_width)
            drew = True
    return drew


def _fallback_geo_layer_image(layer: Dict[str, Any], out_w: int, out_h: int, box_emu: tuple[int, int, int, int]) -> Image.Image:
    im = Image.new("RGBA", (max(1, out_w), max(1, out_h)), (0, 0, 0, 0))
    draw = ImageDraw.Draw(im, "RGBA")
    shape_xml = layer.get("shape_xml") or {}
    fill_spec = layer.get("fill")
    line_spec = layer.get("line")
    if isinstance(shape_xml, dict) and shape_xml.get("format") == "shape_spec_v1":
        if fill_spec is None:
            fill_spec = shape_xml.get("fill")
        if line_spec is None and isinstance(shape_xml.get("line"), dict):
            line_spec = shape_xml.get("line")
    fill_rgba = _rgba_from_color_spec(fill_spec)
    stroke_rgba = _rgba_from_color_spec((line_spec or {}).get("color") if isinstance(line_spec, dict) else None)
    stroke_width = _line_width_px_for_box(line_spec if isinstance(line_spec, dict) else None, box_emu, out_w, out_h)
    inset = max(0.0, stroke_width / 2.0)
    bounds = (inset, inset, max(inset, float(out_w - 1) - inset), max(inset, float(out_h - 1) - inset))
    drew = False

    if isinstance(shape_xml, dict) and shape_xml.get("format") == "shape_spec_v1" and shape_xml.get("type") == "geo":
        geom = shape_xml.get("geom") or {}
        if isinstance(geom, dict) and geom.get("type") == "prstGeom":
            prst = str(geom.get("prst") or "").lower()
            if prst == "rect":
                draw.rectangle(bounds, fill=fill_rgba, outline=stroke_rgba, width=stroke_width if stroke_rgba else 0)
                drew = True
            elif prst in {"roundrect", "roundrectangle"}:
                radius = max(1, int(round(min(out_w, out_h) * 0.125)))
                draw.rounded_rectangle(bounds, radius=radius, fill=fill_rgba, outline=stroke_rgba, width=stroke_width if stroke_rgba else 0)
                drew = True
            elif prst == "ellipse":
                draw.ellipse(bounds, fill=fill_rgba, outline=stroke_rgba, width=stroke_width if stroke_rgba else 0)
                drew = True
        elif isinstance(geom, dict) and geom.get("type") == "custGeom":
            for path in geom.get("paths") or []:
                if not isinstance(path, dict):
                    continue
                drew = _draw_poly_commands(
                    draw,
                    path.get("commands") or [],
                    op_key="op",
                    fill_rgba=fill_rgba,
                    stroke_rgba=stroke_rgba,
                    stroke_width=stroke_width,
                    out_w=out_w,
                    out_h=out_h,
                ) or drew

    if not drew and isinstance(shape_xml, dict) and isinstance(shape_xml.get("spPr"), dict):
        cust = (shape_xml.get("spPr") or {}).get("custGeom") or {}
        for path in cust.get("pathLst") or []:
            if not isinstance(path, dict):
                continue
            drew = _draw_poly_commands(
                draw,
                path.get("commands") or [],
                op_key="type",
                fill_rgba=fill_rgba,
                stroke_rgba=stroke_rgba,
                stroke_width=stroke_width,
                out_w=out_w,
                out_h=out_h,
            ) or drew

    if not drew:
        draw.rectangle((0, 0, max(0, out_w - 1), max(0, out_h - 1)), fill=fill_rgba)
    return im


def _draw_line_marker_fallback(
    draw: ImageDraw.ImageDraw,
    marker_type: Optional[str],
    anchor: tuple[float, float],
    direction: tuple[float, float],
    size: float,
    color: tuple[int, int, int, int],
) -> None:
    if not marker_type:
        return
    dx, dy = direction
    norm = math.hypot(dx, dy)
    if norm <= 1e-6:
        return
    ux, uy = dx / norm, dy / norm
    px, py = -uy, ux
    if marker_type == "oval":
        r = max(1.5, size * 0.5)
        draw.ellipse((anchor[0] - r, anchor[1] - r, anchor[0] + r, anchor[1] + r), fill=color)
        return
    tip = anchor
    base = (anchor[0] - ux * size, anchor[1] - uy * size)
    half = size * 0.5
    p1 = (base[0] + px * half, base[1] + py * half)
    p2 = (base[0] - px * half, base[1] - py * half)
    draw.polygon([tip, p1, p2], fill=color)


def _fallback_line_layer_image(
    layer: Dict[str, Any],
    out_w: int,
    out_h: int,
    line_box_emu: tuple[int, int, int, int],
    canvas_emu: tuple[int, int],
) -> Image.Image:
    im = Image.new("RGBA", (max(1, out_w), max(1, out_h)), (0, 0, 0, 0))
    draw = ImageDraw.Draw(im, "RGBA")
    shape_xml = layer.get("shape_xml") or {}
    p1 = shape_xml.get("p1") or {}
    p2 = shape_xml.get("p2") or {}
    line = shape_xml.get("line") or layer.get("line") or {}
    stroke = _rgba_from_color_spec((line or {}).get("color") if isinstance(line, dict) else None, (0, 0, 0, 255))
    stroke_width = _line_width_px_for_box(line if isinstance(line, dict) else None, line_box_emu, out_w, out_h)
    box_l, box_t, box_w, box_h = line_box_emu
    cw, ch = canvas_emu
    x1 = ((float(p1.get("x") or 0.0) * float(cw)) - box_l) * (float(out_w) / max(1.0, float(box_w)))
    y1 = ((float(p1.get("y") or 0.0) * float(ch)) - box_t) * (float(out_h) / max(1.0, float(box_h)))
    x2 = ((float(p2.get("x") or 0.0) * float(cw)) - box_l) * (float(out_w) / max(1.0, float(box_w)))
    y2 = ((float(p2.get("y") or 0.0) * float(ch)) - box_t) * (float(out_h) / max(1.0, float(box_h)))
    draw.line((x1, y1, x2, y2), fill=stroke, width=stroke_width)
    marker_size = max(6.0, float(stroke_width) * 3.0)
    _draw_line_marker_fallback(draw, _svg_marker_type(line.get("tail_end") if isinstance(line, dict) else None), (x1, y1), (x1 - x2, y1 - y2), marker_size, stroke)
    _draw_line_marker_fallback(draw, _svg_marker_type(line.get("head_end") if isinstance(line, dict) else None), (x2, y2), (x2 - x1, y2 - y1), marker_size, stroke)
    return im


def _css_color_from_spec(spec: Optional[Dict[str, Any]], default: str = "transparent") -> str:
    if not isinstance(spec, dict):
        return default
    rgba = _rgba_from_color_spec(spec)
    if rgba is not None:
        r, g, b, alpha = rgba
        if alpha >= 255:
            return f"#{r:02X}{g:02X}{b:02X}"
        return f"rgba({r}, {g}, {b}, {alpha / 255.0:.4f})"
    kind = str(spec.get("type") or spec.get("color_type") or "").lower()
    alpha = 1.0
    for mod in spec.get("mods") or []:
        if not isinstance(mod, dict):
            continue
        if str(mod.get("op") or "").lower() != "alpha":
            continue
        try:
            alpha = max(0.0, min(1.0, float(mod.get("val")) / 100000.0))
        except Exception:
            alpha = 1.0
        break

    def _to_css(rgb_value: Any) -> str:
        rgb = str(rgb_value or "").strip().lstrip("#")
        if len(rgb) != 6:
            return default
        if alpha >= 0.999:
            return f"#{rgb}"
        try:
            r = int(rgb[0:2], 16)
            g = int(rgb[2:4], 16)
            b = int(rgb[4:6], 16)
        except Exception:
            return default
        return f"rgba({r}, {g}, {b}, {alpha:.4f})"

    if kind in {"rgb", "srgb"} and spec.get("rgb"):
        return _to_css(spec.get("rgb"))
    if kind == "solid" and str(spec.get("color_type") or "").lower() in {"rgb", "srgb"} and spec.get("rgb"):
        return _to_css(spec.get("rgb"))
    if spec.get("rgb"):
        return _to_css(spec.get("rgb"))
    return default


def _dash_array(dash: Any, stroke_width: float) -> Optional[str]:
    val = str(dash or "").lower()
    if not val or val in {"solid", "none"}:
        return None
    unit = max(1.0, float(stroke_width or 1.0))
    if val in {"dash", "dashdot", "lgdash", "lgdashdot", "lgdashdotdot"}:
        return f"{unit * 4:.2f},{unit * 3:.2f}"
    if val in {"sysdot", "dot"}:
        return f"{unit:.2f},{unit * 2:.2f}"
    if val in {"sysdash", "sysdashdot"}:
        return f"{unit * 3:.2f},{unit * 2:.2f}"
    return f"{unit * 4:.2f},{unit * 3:.2f}"


def _svg_path_from_geo_commands(commands: List[Dict[str, Any]], out_w: int, out_h: int, *, op_key: str) -> str:
    parts: List[str] = []
    for cmd in commands:
        if not isinstance(cmd, dict):
            continue
        ctype = str(cmd.get(op_key) or "")
        if ctype in {"moveTo", "lnTo"}:
            pts = cmd.get("pts") if op_key == "op" else None
            pt = (pts[0] if isinstance(pts, list) and pts else None) or cmd.get("pt")
            if not isinstance(pt, dict):
                continue
            x = float(pt.get("x") or 0.0) * out_w
            y = float(pt.get("y") or 0.0) * out_h
            parts.append(("M" if ctype == "moveTo" else "L") + f" {x:.4f} {y:.4f}")
        elif ctype == "cubicBezTo":
            pts = cmd.get("pts")
            if not (isinstance(pts, list) and len(pts) == 3):
                continue
            nums: List[str] = []
            for pt in pts:
                if not isinstance(pt, dict):
                    continue
                x = float(pt.get("x") or 0.0) * out_w
                y = float(pt.get("y") or 0.0) * out_h
                nums.append(f"{x:.4f} {y:.4f}")
            if len(nums) == 3:
                parts.append("C " + " ".join(nums))
        elif ctype == "close":
            parts.append("Z")
    return " ".join(parts)


def _geo_layer_to_svg_text(layer: Dict[str, Any], out_w: int, out_h: int) -> Optional[str]:
    shape_xml = layer.get("shape_xml") or {}
    fill_spec = layer.get("fill")
    line = layer.get("line") or {}
    if isinstance(shape_xml, dict) and shape_xml.get("format") == "shape_spec_v1":
        if fill_spec is None:
            fill_spec = shape_xml.get("fill")
        if not line and isinstance(shape_xml.get("line"), dict):
            line = shape_xml.get("line") or {}

    fill = _css_color_from_spec(fill_spec)
    stroke = _css_color_from_spec(line.get("color") if isinstance(line, dict) else None, "transparent")
    try:
        stroke_width = max(0.0, float((line or {}).get("width_pt") or 0.0) * (96.0 / 72.0))
    except Exception:
        stroke_width = 0.0
    fill_attr = "none" if fill == "transparent" else fill
    stroke_attr = "none" if stroke == "transparent" or stroke_width <= 0 else stroke
    dash_array = _dash_array((line or {}).get("dash"), stroke_width)
    dash_attr = f' stroke-dasharray="{dash_array}"' if dash_array else ""
    shape_tags: List[str] = []

    if isinstance(shape_xml, dict) and shape_xml.get("format") == "shape_spec_v1" and shape_xml.get("type") == "geo":
        geom = shape_xml.get("geom") or {}
        if isinstance(geom, dict) and geom.get("type") == "prstGeom":
            prst = str(geom.get("prst") or "").lower()
            inset = max(0.0, stroke_width / 2.0)
            x = inset
            y = inset
            w = max(0.0, float(out_w) - inset * 2.0)
            h = max(0.0, float(out_h) - inset * 2.0)
            common = f'fill="{fill_attr}" stroke="{stroke_attr}" stroke-width="{stroke_width:.4f}" stroke-linecap="round" stroke-linejoin="round"{dash_attr}'
            if prst == "ellipse":
                shape_tags.append(
                    f'<ellipse cx="{out_w / 2.0:.4f}" cy="{out_h / 2.0:.4f}" rx="{w / 2.0:.4f}" ry="{h / 2.0:.4f}" {common}/>'
                )
            elif prst == "roundrect":
                radius = max(0.0, min(w, h) * 0.12)
                shape_tags.append(f'<rect x="{x:.4f}" y="{y:.4f}" width="{w:.4f}" height="{h:.4f}" rx="{radius:.4f}" ry="{radius:.4f}" {common}/>')
            else:
                shape_tags.append(f'<rect x="{x:.4f}" y="{y:.4f}" width="{w:.4f}" height="{h:.4f}" {common}/>')
        elif isinstance(geom, dict) and geom.get("type") == "custGeom":
            for path in geom.get("paths") or []:
                if not isinstance(path, dict):
                    continue
                d = _svg_path_from_geo_commands(path.get("commands") or [], out_w, out_h, op_key="op")
                if d:
                    shape_tags.append(
                        f'<path d="{d}" fill="{fill_attr}" stroke="{stroke_attr}" stroke-width="{stroke_width:.4f}" stroke-linecap="round" stroke-linejoin="round"{dash_attr}/>'
                    )

    if not shape_tags and isinstance(shape_xml, dict) and isinstance(shape_xml.get("spPr"), dict):
        cust = (shape_xml.get("spPr") or {}).get("custGeom") or {}
        for path in cust.get("pathLst") or []:
            if not isinstance(path, dict):
                continue
            d = _svg_path_from_geo_commands(path.get("commands") or [], out_w, out_h, op_key="type")
            if d:
                shape_tags.append(
                    f'<path d="{d}" fill="{fill_attr}" stroke="{stroke_attr}" stroke-width="{stroke_width:.4f}" stroke-linecap="round" stroke-linejoin="round"{dash_attr}/>'
                )

    if not shape_tags:
        return None
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{out_w}" height="{out_h}" '
        f'viewBox="0 0 {out_w} {out_h}" preserveAspectRatio="none">'
        + "".join(shape_tags)
        + "</svg>"
    )


def _svg_marker_type(end_spec: Any) -> Optional[str]:
    if not isinstance(end_spec, dict):
        return None
    marker_type = str(end_spec.get("type") or "").strip().lower()
    if not marker_type or marker_type in {"none", "nil"}:
        return None
    return marker_type


def _svg_stroke_linecap(line: Optional[Dict[str, Any]]) -> str:
    cap = str(line.get("cap") or "").strip().lower() if isinstance(line, dict) else ""
    if cap == "rnd":
        return "round"
    if cap == "sq":
        return "square"
    return "butt"


def _line_end_marker_defs(line: Dict[str, Any], stroke: str) -> tuple[str, str, str]:
    defs: List[str] = []
    marker_start = ""
    marker_end = ""
    for key, marker_id, attr, orient in (
        ("tail_end", "line_tail", "marker-start", "auto-start-reverse"),
        ("head_end", "line_head", "marker-end", "auto"),
    ):
        marker_type = _svg_marker_type(line.get(key))
        if marker_type is None:
            continue
        if marker_type == "oval":
            body = f'<circle cx="3" cy="3" r="3" fill="{stroke}"/>'
            view_box = "0 0 6 6"
            ref_x = "3"
            ref_y = "3"
        else:
            body = f'<path d="M 0 0 L 8 4 L 0 8 z" fill="{stroke}"/>'
            view_box = "0 0 8 8"
            ref_x = "8" if attr == "marker-end" else "0"
            ref_y = "4"
        defs.append(
            f'<marker id="{marker_id}" viewBox="{view_box}" refX="{ref_x}" refY="{ref_y}" '
            f'markerWidth="6" markerHeight="6" orient="{orient}" markerUnits="strokeWidth">'
            f"{body}</marker>"
        )
        if attr == "marker-start":
            marker_start = f' marker-start="url(#{marker_id})"'
        else:
            marker_end = f' marker-end="url(#{marker_id})"'
    defs_text = f"<defs>{''.join(defs)}</defs>" if defs else ""
    return defs_text, marker_start, marker_end


def _line_layer_to_svg_text(
    layer: Dict[str, Any],
    out_w: int,
    out_h: int,
    line_box_emu: tuple[int, int, int, int],
    canvas_emu: tuple[int, int],
) -> Optional[str]:
    shape_xml = layer.get("shape_xml") or {}
    if not isinstance(shape_xml, dict):
        return None
    p1 = shape_xml.get("p1")
    p2 = shape_xml.get("p2")
    line = shape_xml.get("line") or layer.get("line") or {}
    if not (isinstance(p1, dict) and isinstance(p2, dict) and isinstance(line, dict)):
        return None
    stroke = _css_color_from_spec(line.get("color"), "transparent")
    if stroke == "transparent":
        return None
    try:
        stroke_width_emu = max(float(line.get("width_pt") or 0.75) * EMU_PER_PT, 1.0)
    except Exception:
        stroke_width_emu = 0.75 * EMU_PER_PT
    box_l, box_t, box_w, box_h = line_box_emu
    if box_w <= 0 or box_h <= 0:
        return None
    cw, ch = canvas_emu
    x1_emu = float(p1.get("x") or 0.0) * float(cw)
    y1_emu = float(p1.get("y") or 0.0) * float(ch)
    x2_emu = float(p2.get("x") or 0.0) * float(cw)
    y2_emu = float(p2.get("y") or 0.0) * float(ch)
    scale_x = float(out_w) / float(box_w)
    scale_y = float(out_h) / float(box_h)
    stroke_width = max(1.0, stroke_width_emu * min(scale_x, scale_y))
    x1 = (x1_emu - box_l) * scale_x
    y1 = (y1_emu - box_t) * scale_y
    x2 = (x2_emu - box_l) * scale_x
    y2 = (y2_emu - box_t) * scale_y
    linecap = _svg_stroke_linecap(line)
    dash_array = _dash_array(line.get("dash"), stroke_width)
    dash_attr = f' stroke-dasharray="{dash_array}"' if dash_array else ""
    defs_text, marker_start, marker_end = _line_end_marker_defs(line, stroke)
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{out_w}" height="{out_h}" '
        f'viewBox="0 0 {out_w} {out_h}" preserveAspectRatio="none">'
        f"{defs_text}"
        f'<line x1="{x1:.4f}" y1="{y1:.4f}" x2="{x2:.4f}" y2="{y2:.4f}" '
        f'stroke="{stroke}" stroke-width="{stroke_width:.4f}" stroke-linecap="{linecap}" '
        f'stroke-linejoin="round"{dash_attr}{marker_start}{marker_end}/>'
        "</svg>"
    )


def _line_image_box_emu(
    layer: Dict[str, Any],
    canvas_emu: tuple[int, int],
) -> Optional[tuple[int, int, int, int]]:
    shape_xml = layer.get("shape_xml") or {}
    if not isinstance(shape_xml, dict):
        return None
    p1 = shape_xml.get("p1")
    p2 = shape_xml.get("p2")
    if not (isinstance(p1, dict) and isinstance(p2, dict)):
        return None
    cw, ch = canvas_emu
    x1 = float(p1.get("x") or 0.0) * float(cw)
    y1 = float(p1.get("y") or 0.0) * float(ch)
    x2 = float(p2.get("x") or 0.0) * float(cw)
    y2 = float(p2.get("y") or 0.0) * float(ch)
    line = shape_xml.get("line") or {}
    try:
        width_emu = max(float(line.get("width_pt") or 0.75) * EMU_PER_PT, 1.0) if isinstance(line, dict) else 0.75 * EMU_PER_PT
    except Exception:
        width_emu = 0.75 * EMU_PER_PT
    marker_pad_multiplier = 6.0 if (
        isinstance(line, dict)
        and (_svg_marker_type(line.get("head_end")) is not None or _svg_marker_type(line.get("tail_end")) is not None)
    ) else 0.75
    pad = max(width_emu * marker_pad_multiplier, min(float(cw), float(ch)) * 0.0005)
    left = int(math.floor(min(x1, x2) - pad))
    top = int(math.floor(min(y1, y2) - pad))
    right = int(math.ceil(max(x1, x2) + pad))
    bottom = int(math.ceil(max(y1, y2) + pad))
    return left, top, max(1, right - left), max(1, bottom - top)


def _save_ppt_graph_geo_svg_image_layer(
    *,
    layer: Dict[str, Any],
    shp,
    base_box_emu: tuple[int, int, int, int],
    final_box_emu: tuple[int, int, int, int],
    canvas_emu: tuple[int, int],
    assets_dir: Path,
    transform: Optional[Dict[str, float]] = None,
) -> Optional[Dict[str, Any]]:
    out_w, out_h = _raster_size_2048_for_box(base_box_emu)
    svg_text = _geo_layer_to_svg_text(layer, out_w, out_h)
    im: Optional[Image.Image] = None
    if svg_text:
        png_bytes = _svg_bytes_to_png_bytes(svg_text.encode("utf-8"), out_w, out_h)
        if png_bytes:
            try:
                im = Image.open(io.BytesIO(png_bytes)).convert("RGBA")
            except Exception:
                im = None
    if im is None:
        im = _fallback_geo_layer_image(layer, out_w, out_h, base_box_emu)
    shape_id = getattr(shp, "shape_id", None) or 0
    out_path = (assets_dir / f"slide{int(layer.get('slide') or 0):03d}_shape{shape_id}_ppt_graph_geo_rendered.png").resolve()
    flips = shape_flip_flags(shp)
    rendered = _save_final_image_layer(
        im=im,
        out_path=out_path,
        final_box_emu=final_box_emu,
        canvas_emu=canvas_emu,
        rotation_deg=shape_rotation_deg(shp, transform),
        flip_h=bool(flips.get("flip_h")),
        flip_v=bool(flips.get("flip_v")),
    )
    if rendered is None:
        return None
    rendered_path, rendered_box_emu = rendered
    placed_box_emu = rendered_box_emu or final_box_emu
    return {
        "slide": layer.get("slide"),
        "shape_name": layer.get("shape_name"),
        "kind": layer.get("kind"),
        "box": _ratio_box(*placed_box_emu, canvas_emu),
        "saved_path": rendered_path,
        "structure_info": {
            "box": layer.get("box"),
            "shape_xml": layer.get("shape_xml"),
            "fill": layer.get("fill"),
            "line": layer.get("line"),
            "rotation_deg": layer.get("rotation_deg"),
            "flip_h": layer.get("flip_h"),
            "flip_v": layer.get("flip_v"),
        },
    }


def _save_ppt_graph_line_svg_image_layer(
    *,
    layer: Dict[str, Any],
    shp,
    canvas_emu: tuple[int, int],
    assets_dir: Path,
) -> Optional[Dict[str, Any]]:
    line_box_emu = _line_image_box_emu(layer, canvas_emu)
    if line_box_emu is None:
        return None
    out_w, out_h = _raster_size_2048_for_box(line_box_emu)
    svg_text = _line_layer_to_svg_text(layer, out_w, out_h, line_box_emu, canvas_emu)
    im: Optional[Image.Image] = None
    if svg_text:
        png_bytes = _svg_bytes_to_png_bytes(svg_text.encode("utf-8"), out_w, out_h)
        if png_bytes:
            try:
                im = Image.open(io.BytesIO(png_bytes)).convert("RGBA")
            except Exception:
                im = None
    if im is None:
        im = _fallback_line_layer_image(layer, out_w, out_h, line_box_emu, canvas_emu)
    shape_id = getattr(shp, "shape_id", None) or 0
    out_path = (assets_dir / f"slide{int(layer.get('slide') or 0):03d}_shape{shape_id}_ppt_graph_line_rendered.png").resolve()
    rendered = _save_final_image_layer(
        im=im,
        out_path=out_path,
        final_box_emu=line_box_emu,
        canvas_emu=canvas_emu,
        rotation_deg=0.0,
        flip_h=False,
        flip_v=False,
    )
    if rendered is None:
        return None
    rendered_path, rendered_box_emu = rendered
    placed_box_emu = rendered_box_emu or line_box_emu
    return {
        "slide": layer.get("slide"),
        "shape_name": layer.get("shape_name"),
        "kind": layer.get("kind"),
        "box": _ratio_box(*placed_box_emu, canvas_emu),
        "saved_path": rendered_path,
        "structure_info": {
            "box": layer.get("box"),
            "shape_xml": layer.get("shape_xml"),
            "fill": layer.get("fill"),
            "line": layer.get("line"),
            "rotation_deg": layer.get("rotation_deg"),
            "flip_h": layer.get("flip_h"),
            "flip_v": layer.get("flip_v"),
        },
    }


def _export_gradient_background_layer(
    gradient_spec: Dict[str, Any],
    slide_index: int,
    slide_width_emu: int,
    slide_height_emu: int,
    assets_dir: Path,
    *,
    suffix: str = "gradient",
) -> Optional[Dict[str, Any]]:
    full_box_emu = (0, 0, int(slide_width_emu), int(slide_height_emu))
    out_w, out_h = _raster_size_2048_for_box(full_box_emu)
    out_path = assets_dir / f"slide{slide_index:03d}_background_{suffix}_svg_image_png_rgba.png"
    try:
        im = Image.new("RGBA", (out_w, out_h), (0, 0, 0, 0))
        _paint_gradient_fill_rect(im, (0, 0, out_w, out_h), gradient_spec)
        im.save(out_path, format="PNG")
    except Exception:
        return None

    return {
        "slide": slide_index,
        "shape_name": "slide_background_gradient",
        "kind": "svg_image_png_rgba",
        "box": {"left": 0.0, "top": 0.0, "width": 1.0, "height": 1.0},
        "saved_path": str(out_path),
        "source": "slide_background_gradient",
    }


def _style_ref_color_spec_from_root(root: ET.Element, ref_name: str) -> Optional[Dict[str, Any]]:
    ref = root.find(f".//p:style/a:{ref_name}", NS)
    if ref is None:
        return None
    for child in list(ref):
        spec = _color_spec_from_color_el(child)
        if spec:
            return spec
    return None


def _line_style_from_root(root: ET.Element) -> Optional[Dict[str, Any]]:
    ln = root.find(".//a:ln", NS)
    if ln is None:
        return None
    width_emu = None
    try:
        width_emu = int(ln.get("w") or 0)
    except Exception:
        width_emu = None
    color = _solid_fill_spec_from_node(ln) or _style_ref_color_spec_from_root(root, "lnRef")
    out: Dict[str, Any] = {}
    if width_emu is not None and width_emu > 0:
        out["width_pt"] = round(float(width_emu) / EMU_PER_PT, 4)
    if color is not None:
        out["color"] = color
    dash = ln.find("a:prstDash", NS)
    if dash is not None and dash.get("val"):
        out["dash"] = str(dash.get("val"))
    if ln.get("cap"):
        out["cap"] = str(ln.get("cap"))
    head = ln.find("a:headEnd", NS)
    if head is not None and head.get("type"):
        out["head_end"] = {
            "type": str(head.get("type")),
            "w": head.get("w"),
            "len": head.get("len"),
        }
    tail = ln.find("a:tailEnd", NS)
    if tail is not None and tail.get("type"):
        out["tail_end"] = {
            "type": str(tail.get("type")),
            "w": tail.get("w"),
            "len": tail.get("len"),
        }
    return out or None


def _has_visible_fill(fill: Optional[Dict[str, Any]]) -> bool:
    return isinstance(fill, dict) and bool(fill)


def _has_visible_line(line: Optional[Dict[str, Any]]) -> bool:
    if not isinstance(line, dict):
        return False
    try:
        width_pt = float(line.get("width_pt") or 0.0)
    except Exception:
        width_pt = 0.0
    return width_pt > 0 and line.get("color") is not None


def _is_fully_transparent_geo(fill: Optional[Dict[str, Any]], line: Optional[Dict[str, Any]]) -> bool:
    return (not _has_visible_fill(fill)) and (not _has_visible_line(line))


def _cust_geom_paths_spec_from_root(root: ET.Element) -> Optional[Dict[str, Any]]:
    cust = root.find(".//a:custGeom", NS)
    if cust is None:
        return None
    path_lst = cust.find("a:pathLst", NS)
    if path_lst is None:
        return None

    paths_out: List[Dict[str, Any]] = []
    for path in path_lst.findall("a:path", NS):
        try:
            w = int(path.get("w") or 0)
            h = int(path.get("h") or 0)
        except Exception:
            w, h = 0, 0
        if w <= 0 or h <= 0:
            return None

        commands: List[Dict[str, Any]] = []
        def _round_coord(v: float) -> float:
            r = round(float(v), 5)
            return 0.0 if r == -0.0 else r
        for child in list(path):
            tag = child.tag.rsplit("}", 1)[-1]
            if tag == "moveTo":
                pt = child.find("a:pt", NS)
                if pt is None:
                    continue
                x = int(pt.get("x") or 0)
                y = int(pt.get("y") or 0)
                commands.append({"op": "moveTo", "pts": [{"x": _round_coord(x / w), "y": _round_coord(y / h)}]})
            elif tag == "lnTo":
                pt = child.find("a:pt", NS)
                if pt is None:
                    continue
                x = int(pt.get("x") or 0)
                y = int(pt.get("y") or 0)
                commands.append({"op": "lnTo", "pts": [{"x": _round_coord(x / w), "y": _round_coord(y / h)}]})
            elif tag == "cubicBezTo":
                pts = child.findall("a:pt", NS)
                if len(pts) != 3:
                    continue
                out_pts = []
                for pt in pts:
                    x = int(pt.get("x") or 0)
                    y = int(pt.get("y") or 0)
                    out_pts.append({"x": _round_coord(x / w), "y": _round_coord(y / h)})
                commands.append({"op": "cubicBezTo", "pts": out_pts})
            elif tag == "close":
                commands.append({"op": "close", "pts": []})
            else:
                # Unknown ops are rare in current templates; keep it strict.
                return None

        paths_out.append({"w": w, "h": h, "commands": commands})

    return {"type": "custGeom", "paths": paths_out}


def _extract_sppr_custgeom_spec_from_shape_xml(shape_xml: str, canvas_emu: tuple[int, int]) -> Optional[Dict[str, Any]]:
    """Export a compact spPr spec with fully relative geometry.

    - xfrm.off/ext are omitted (set to None); rebuild uses `box` for placement.
    - path w/h are stored relative to canvas width/height.
    - path points are stored relative to the path's own w/h (0..1).
    """
    try:
        root = ET.fromstring(shape_xml)
    except Exception:
        return None

    cw, ch = canvas_emu
    if cw <= 0 or ch <= 0:
        return None

    def _round_coord(v: float) -> float:
        r = round(float(v), 5)
        return 0.0 if r == -0.0 else r

    xfrm = root.find(".//a:xfrm", NS)
    flip_h = False
    flip_v = False
    try:
        if xfrm is not None:
            flip_h = str(xfrm.get("flipH") or "").lower() in {"1", "true"}
            flip_v = str(xfrm.get("flipV") or "").lower() in {"1", "true"}
    except Exception:
        flip_h = False
        flip_v = False

    cust = root.find(".//a:custGeom", NS)
    if cust is None:
        return None
    path_lst = cust.find("a:pathLst", NS)
    if path_lst is None:
        return None

    path_items: List[Dict[str, Any]] = []
    for path in path_lst.findall("a:path", NS):
        try:
            w = int(path.get("w") or 0)
            h = int(path.get("h") or 0)
        except Exception:
            w, h = 0, 0
        if w <= 0 or h <= 0:
            continue

        commands: List[Dict[str, Any]] = []
        for child in list(path):
            tag = child.tag.rsplit("}", 1)[-1]
            if tag in {"moveTo", "lnTo"}:
                pt = child.find("a:pt", NS)
                if pt is None:
                    continue
                x = int(pt.get("x") or 0)
                y = int(pt.get("y") or 0)
                commands.append({"type": tag, "pt": {"x": _round_coord(x / w), "y": _round_coord(y / h)}})
            elif tag == "cubicBezTo":
                pts = child.findall("a:pt", NS)
                if len(pts) != 3:
                    continue
                out_pts = []
                for pt in pts:
                    x = int(pt.get("x") or 0)
                    y = int(pt.get("y") or 0)
                    out_pts.append({"x": _round_coord(x / w), "y": _round_coord(y / h)})
                commands.append({"type": "cubicBezTo", "pts": out_pts})
            elif tag == "close":
                commands.append({"type": "close"})
            else:
                continue

        # w/h relative to canvas
        path_items.append({"w": _round_coord(w / cw), "h": _round_coord(h / ch), "commands": commands})

    if not path_items:
        return None

    return {
        "spPr": {
            "xfrm": {
                "flipH": bool(flip_h),
                "flipV": bool(flip_v),
                "off": None,
                "ext": None,
            },
            "custGeom": {
                "avLst": [],
                "gdLst": [],
                "ahLst": [],
                "cxnLst": [],
                "rect": {"l": "l", "t": "t", "r": "r", "b": "b"},
                "pathLst": path_items,
            },
        }
    }


def _prst_geom_spec_from_root(root: ET.Element) -> Optional[Dict[str, Any]]:
    prst = root.find(".//a:prstGeom", NS)
    if prst is None:
        return None
    name = prst.get("prst")
    if not name:
        return None
    return {"type": "prstGeom", "prst": str(name)}

def _graph_shape_spec_from_shape_xml(shape_xml: str, kind: str, box_emu: tuple[int, int, int, int], canvas_emu: tuple[int, int]) -> Optional[Dict[str, Any]]:
    try:
        root = ET.fromstring(shape_xml)
    except Exception:
        return None

    if kind == "ppt_graph_geo":
        geom = _cust_geom_paths_spec_from_root(root) or _prst_geom_spec_from_root(root)
        if geom is None:
            return None
        sp_pr = root.find(".//p:spPr", NS)
        fill = (_solid_fill_spec_from_node(sp_pr) if sp_pr is not None else None) or _style_ref_color_spec_from_root(root, "fillRef")
        if isinstance(geom, dict) and geom.get("type") == "prstGeom" and geom.get("prst") in {
            "leftBrace",
            "rightBrace",
            "leftBracket",
            "rightBracket",
            "bracePair",
            "bracketPair",
        }:
            fill = None
        return {
            "format": "shape_spec_v1",
            "type": "geo",
            "geom": geom,
            "fill": fill,
            "line": _line_style_from_root(root),
        }

    return None


def _is_outer_page_border_graph_layer(layer: Optional[Dict[str, Any]]) -> bool:
    """
    Filter decorative page-edge frames at export time only.
    We target a nearly full-page no-fill rectangle with a single normalized rect path.
    """
    if not isinstance(layer, dict):
        return False
    if layer.get("kind") != "ppt_graph_geo":
        return False
    if layer.get("fill") is not None:
        return False
    if not isinstance(layer.get("line"), dict):
        return False

    box = layer.get("box") or {}
    try:
        left = float(box.get("left") or 0.0)
        top = float(box.get("top") or 0.0)
        width = float(box.get("width") or 0.0)
        height = float(box.get("height") or 0.0)
    except Exception:
        return False
    if not (width >= 0.90 and height >= 0.90 and left <= 0.06 and top <= 0.06):
        return False

    shape_xml = layer.get("shape_xml") or {}
    sppr = (shape_xml.get("spPr") or {}) if isinstance(shape_xml, dict) else {}
    cust = sppr.get("custGeom") if isinstance(sppr, dict) else None
    if not isinstance(cust, dict):
        return False
    path_lst = cust.get("pathLst") or []
    if not (isinstance(path_lst, list) and len(path_lst) == 1 and isinstance(path_lst[0], dict)):
        return False
    cmds = path_lst[0].get("commands") or []
    if not (isinstance(cmds, list) and len(cmds) >= 5):
        return False
    try:
        pts = []
        for cmd in cmds[:4]:
            if cmd.get("type") not in {"moveTo", "lnTo"}:
                return False
            pt = cmd.get("pt") or {}
            pts.append((float(pt.get("x") or 0.0), float(pt.get("y") or 0.0)))
        if str(cmds[4].get("type") or "").lower() != "close":
            return False
        if pts != [(0.0, 0.0), (1.0, 0.0), (1.0, 1.0), (0.0, 1.0)]:
            return False
    except Exception:
        return False
    return True


def extract_graph_layer(
    slide, shp, slide_index: int, assets_dir: Path, transform: Optional[Dict[str, float]] = None
) -> Optional[Dict[str, Any]]:
    st = getattr(shp, "shape_type", None)
    # TextBoxes are usually exported as editable text layers, but some templates use empty
    # TextBox shapes as colored rectangles/panels (no text). Preserve those as geo layers.
    allow_textbox_geo = False
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        has_text = False
        try:
            if getattr(shp, "has_text_frame", False) and getattr(shp, "text_frame", None) is not None:
                tf = shp.text_frame
                if (tf.text or "").strip():
                    has_text = True
                else:
                    for p in tf.paragraphs:
                        for r in getattr(p, "runs", []):
                            if (getattr(r, "text", "") or "").strip():
                                has_text = True
                                break
                        if has_text:
                            break
        except Exception:
            has_text = False
        allow_textbox_geo = not has_text
        if not allow_textbox_geo:
            return None
    # Connector/arrow lines are not consistently exposed as a named enum across
    # python-pptx versions. Detect them from the underlying OOXML tag as well.
    el = getattr(shp, "_element", None)
    is_connector = False
    try:
        is_connector = el is not None and str(getattr(el, "tag", "")).endswith("}cxnSp")
    except Exception:
        is_connector = False
    if st not in {MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.TEXT_BOX} and not is_connector:
        return None
    if _shape_has_blip_fill(shp):
        # Picture-fill shapes are already exported via fill-image extraction.
        return None

    shape_xml = _shape_xml_without_text(shp, transform)
    if not shape_xml:
        return None

    # Some templates use preset geometry `prstGeom prst="line"` which python-pptx
    # exposes as AUTO_SHAPE instead of LINE, so detect it from OOXML too.
    is_prst_line = 'prst="line"' in shape_xml
    kind = "ppt_graph_line" if (st == MSO_SHAPE_TYPE.LINE or is_connector or is_prst_line) else "ppt_graph_geo"
    canvas_emu = _slide_canvas_emu(slide)

    # Export the unrotated xfrm-like box. Rotation is exported separately and applied on rebuild.
    # Using a post-rotation AABB here would double-apply rotation in combine.py.
    box_emu = _shape_bounds_box_emu(shp, transform) or _shape_box_emu(shp, transform)
    if box_emu is None:
        return None

    # Normalize ET-generated prefixes (ns0/ns1) to stable p/a for readability and diffing.
    shape_xml_norm = (
        shape_xml.replace('xmlns:ns0="http://schemas.openxmlformats.org/presentationml/2006/main"', 'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"')
        .replace('xmlns:ns1="http://schemas.openxmlformats.org/drawingml/2006/main"', 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"')
        .replace("<ns0:", "<p:")
        .replace("</ns0:", "</p:")
        .replace("<ns1:", "<a:")
        .replace("</ns1:", "</a:")
    )

    # Lines/connectors: export as endpoints + line style (incl. arrowheads).
    if kind == "ppt_graph_line":
        cw, ch = canvas_emu
        if cw <= 0 or ch <= 0:
            return None
        p1 = None
        p2 = None
        point_source = "connector_endpoints"
        try:
            p1 = {"x": float(shp.begin_x) / cw, "y": float(shp.begin_y) / ch}
            p2 = {"x": float(shp.end_x) / cw, "y": float(shp.end_y) / ch}
        except Exception:
            point_source = "box_rotation_fallback"
            l, t, w, h = (float(v) for v in box_emu)
            if w == 0 and h == 0:
                p1 = {"x": l / cw, "y": t / ch}
                p2 = {"x": l / cw, "y": t / ch}
            else:
                flips = shape_flip_flags(shp)
                if w == 0:
                    lp1 = (l, t)
                    lp2 = (l, t + h)
                elif h == 0:
                    lp1 = (l, t)
                    lp2 = (l + w, t)
                else:
                    lp1 = (l, t)
                    lp2 = (l + w, t + h)
                    if bool(flips.get("flip_h")):
                        lp1 = (l + w, lp1[1])
                        lp2 = (l, lp2[1])
                    if bool(flips.get("flip_v")):
                        lp1 = (lp1[0], t + h if lp1[1] == t else t)
                        lp2 = (lp2[0], t + h if lp2[1] == t else t)

                cx = l + (w / 2.0)
                cy = t + (h / 2.0)
                theta = math.radians(float(shape_rotation_deg(shp, transform) or 0.0))

                def _rot_pt(px: float, py: float) -> tuple[float, float]:
                    dx = px - cx
                    dy = py - cy
                    rx = cx + (dx * math.cos(theta) - dy * math.sin(theta))
                    ry = cy + (dx * math.sin(theta) + dy * math.cos(theta))
                    return rx, ry

                x1, y1 = _rot_pt(*lp1)
                x2, y2 = _rot_pt(*lp2)
                p1 = {"x": x1 / cw, "y": y1 / ch}
                p2 = {"x": x2 / cw, "y": y2 / ch}
        try:
            root = ET.fromstring(shape_xml_norm)
        except Exception:
            return None
        shape_spec = {
            "p1": p1,
            "p2": p2,
            "line": _line_style_from_root(root),
        }
        out = {
            "slide": slide_index,
            "shape_name": getattr(shp, "name", None),
            "kind": kind,
            "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
            **shape_flip_flags(shp),
            "box": _ratio_box(*box_emu, canvas_emu),
            "shape_xml": shape_spec,
        }
        rendered = _save_ppt_graph_line_svg_image_layer(
            layer=out,
            shp=shp,
            canvas_emu=canvas_emu,
            assets_dir=assets_dir,
        )
        if rendered is None:
            raise RuntimeError(f"ppt_graph_line image export failed for shape={getattr(shp, 'name', None)!r}")
        return _validate_dual_export_layer(rendered)

    shape_spec = _extract_sppr_custgeom_spec_from_shape_xml(shape_xml_norm, canvas_emu)
    # Fallback: preset geometry (prstGeom) and non-custGeom shapes.
    if shape_spec is None:
        fallback = _graph_shape_spec_from_shape_xml(shape_xml_norm, kind, box_emu, canvas_emu)
        if fallback is None:
            return None
        if kind == "ppt_graph_geo" and _is_fully_transparent_geo(fallback.get("fill"), fallback.get("line")):
            return None
        out = {
            "slide": slide_index,
            "shape_name": getattr(shp, "name", None),
            "kind": kind,
            "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
            **shape_flip_flags(shp),
            "box": _ratio_box(*box_emu, canvas_emu),
            "shape_xml": fallback,
        }
        if _is_outer_page_border_graph_layer(out):
            return None
        rendered = _save_ppt_graph_geo_svg_image_layer(
            layer=out,
            shp=shp,
            base_box_emu=box_emu,
            final_box_emu=_shape_box_emu(shp, transform) or box_emu,
            canvas_emu=canvas_emu,
            assets_dir=assets_dir,
            transform=transform,
        )
        if rendered is None:
            raise RuntimeError(f"ppt_graph_geo image export failed for shape={getattr(shp, 'name', None)!r}")
        return _validate_dual_export_layer(rendered)

    # Extract fill/line styles. Many templates use theme colors (schemeClr) via style refs;
    # keep a compact color spec that combine.py can re-apply deterministically.
    fill = None
    line = None
    try:
        root = ET.fromstring(shape_xml_norm)
        sp_pr = root.find(".//p:spPr", NS)
        fill = (_solid_fill_spec_from_node(sp_pr) if sp_pr is not None else None) or _style_ref_color_spec_from_root(
            root, "fillRef"
        )
        line = _line_style_from_root(root)
    except Exception:
        fill = None
        line = None

    if kind == "ppt_graph_geo" and _is_fully_transparent_geo(fill, line):
        return None

    out = {
        "slide": slide_index,
        "shape_name": getattr(shp, "name", None),
        "kind": kind,
        "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
        **shape_flip_flags(shp),
        "box": _ratio_box(*box_emu, canvas_emu),
        "shape_xml": shape_spec,
        "fill": fill,
        "line": line,
    }
    if _is_outer_page_border_graph_layer(out):
        return None
    rendered = _save_ppt_graph_geo_svg_image_layer(
        layer=out,
        shp=shp,
        base_box_emu=box_emu,
        final_box_emu=_shape_box_emu(shp, transform) or box_emu,
        canvas_emu=canvas_emu,
        assets_dir=assets_dir,
        transform=transform,
    )
    if rendered is None:
        raise RuntimeError(f"ppt_graph_geo image export failed for shape={getattr(shp, 'name', None)!r}")
    return _validate_dual_export_layer(rendered)


def _line_style_from_ln_el(ln_el) -> Optional[Dict[str, Any]]:
    if ln_el is None:
        return None
    try:
        root = ET.fromstring(ET.tostring(ln_el, encoding="utf-8"))
    except Exception:
        return None
    width_emu = None
    try:
        width_emu = int(root.get("w") or 0)
    except Exception:
        width_emu = None
    color = _solid_fill_spec_from_node(root)
    out: Dict[str, Any] = {}
    if width_emu is not None and width_emu > 0:
        out["width_pt"] = round(float(width_emu) / EMU_PER_PT, 4)
    if color is not None:
        out["color"] = color
    dash = root.find("a:prstDash", NS)
    if dash is not None and dash.get("val"):
        out["dash"] = str(dash.get("val"))
    return out or None


def extract_table_layer(
    slide,
    shp,
    slide_index: int,
    canvas_emu: Optional[tuple[int, int]],
    transform: Optional[Dict[str, float]] = None,
) -> Optional[Dict[str, Any]]:
    st = getattr(shp, "shape_type", None)
    if st != MSO_SHAPE_TYPE.TABLE and not hasattr(shp, "table"):
        return None
    try:
        tbl = shp.table
    except Exception:
        return None

    table_style_id = None
    try:
        root = ET.fromstring(ET.tostring(getattr(shp, "_element", None), encoding="utf-8"))
        tbl_el = root.find(".//a:tbl", NS)
        if tbl_el is not None:
            tbl_pr = tbl_el.find("a:tblPr", NS)
            if tbl_pr is not None:
                sid = tbl_pr.find("a:tableStyleId", NS)
                if sid is not None and sid.text and sid.text.strip():
                    table_style_id = sid.text.strip()
    except Exception:
        table_style_id = None

    try:
        rows = len(tbl.rows)
        cols = len(tbl.columns)
    except Exception:
        return None
    if rows <= 0 or cols <= 0:
        return None

    col_widths_pt: List[float] = []
    for col in tbl.columns:
        try:
            col_widths_pt.append(float(col.width.pt))
        except Exception:
            col_widths_pt.append(0.0)

    row_heights_pt: List[float] = []
    for row in tbl.rows:
        try:
            row_heights_pt.append(float(row.height.pt))
        except Exception:
            row_heights_pt.append(0.0)

    cells: List[List[Dict[str, Any]]] = []
    for ri in range(rows):
        row_cells: List[Dict[str, Any]] = []
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            paragraphs: List[Dict[str, Any]] = []
            margin_left_pt = None
            margin_top_pt = None
            margin_right_pt = None
            margin_bottom_pt = None
            vertical_anchor = None
            body_pr = None
            try:
                tf = cell.text_frame
                for pi, para in enumerate(tf.paragraphs):
                    paragraphs.append(_paragraph_info(para, canvas_emu))
                try:
                    margin_left_pt = float(tf.margin_left.pt)
                    margin_top_pt = float(tf.margin_top.pt)
                    margin_right_pt = float(tf.margin_right.pt)
                    margin_bottom_pt = float(tf.margin_bottom.pt)
                except Exception:
                    pass
                try:
                    va = getattr(tf, "vertical_anchor", None)
                    vertical_anchor = str(va) if va is not None else None
                except Exception:
                    vertical_anchor = None
                try:
                    tc = getattr(cell, "_tc", None)
                    tx_body = tc.find(qn("a:txBody")) if tc is not None and hasattr(tc, "find") else None
                    bp = tx_body.find(qn("a:bodyPr")) if tx_body is not None and hasattr(tx_body, "find") else None
                    if bp is not None:
                        root = ET.fromstring(ET.tostring(bp, encoding="utf-8"))
                        attrs: Dict[str, Any] = dict(root.attrib)
                        autofit = None
                        if root.find("a:spAutoFit", NS) is not None:
                            autofit = "spAutoFit"
                        elif root.find("a:normAutoFit", NS) is not None:
                            autofit = "normAutoFit"
                        elif root.find("a:noAutoFit", NS) is not None:
                            autofit = "noAutoFit"
                        body_pr = {"attrs": attrs, "autofit": autofit}
                except Exception:
                    body_pr = None
                try:
                    tc = getattr(cell, "_tc", None)
                    tcpr = tc.tcPr if tc is not None and hasattr(tc, "tcPr") else None
                    if tcpr is not None:
                        tcpr_root = ET.fromstring(ET.tostring(tcpr, encoding="utf-8"))
                        anchor_attr = tcpr_root.get("anchor")
                        if anchor_attr:
                            vertical_anchor = str(anchor_attr)
                        for attr_name, target_name in (
                            ("marL", "margin_left_pt"),
                            ("marT", "margin_top_pt"),
                            ("marR", "margin_right_pt"),
                            ("marB", "margin_bottom_pt"),
                        ):
                            raw_value = tcpr_root.get(attr_name)
                            if raw_value not in (None, ""):
                                pt_value = float(raw_value) / EMU_PER_PT
                                if target_name == "margin_left_pt":
                                    margin_left_pt = pt_value
                                elif target_name == "margin_top_pt":
                                    margin_top_pt = pt_value
                                elif target_name == "margin_right_pt":
                                    margin_right_pt = pt_value
                                elif target_name == "margin_bottom_pt":
                                    margin_bottom_pt = pt_value
                except Exception:
                    pass
            except Exception:
                paragraphs = []

            fill = None
            borders: Dict[str, Any] = {}
            try:
                tc = getattr(cell, "_tc", None)
                tcpr = tc.tcPr if tc is not None and hasattr(tc, "tcPr") else None
                if tcpr is not None:
                    tcpr_xml = ET.fromstring(ET.tostring(tcpr, encoding="utf-8"))
                    grad_fill = _gradient_fill_spec_from_node(tcpr_xml)
                    fill = grad_fill or _solid_fill_spec_from_node(tcpr_xml)
                    # #region debug-point H:table-cell-fill-extract
                    try:
                        if fill is not None or grad_fill is not None:
                            _dbg_report_table_gradient_fill(
                                "H",
                                "test.py:extract_table_layer",
                                "[DEBUG] extracted table cell fill",
                                {
                                    "slide": slide_index,
                                    "shape_name": getattr(shp, "name", None),
                                    "row": ri,
                                    "col": ci,
                                    "solid_fill": fill,
                                    "gradient_fill": grad_fill,
                                },
                            )
                    except Exception:
                        pass
                    # #endregion
                    for tag, key in (("lnL", "l"), ("lnR", "r"), ("lnT", "t"), ("lnB", "b")):
                        ln = tcpr.find(qn(f"a:{tag}")) if hasattr(tcpr, "find") else None
                        spec = _line_style_from_ln_el(ln)
                        if spec is not None:
                            borders[key] = spec
            except Exception:
                fill = None
                borders = {}

            row_cells.append(
                {
                    "row": ri,
                    "col": ci,
                    "text": _normalize_text_breaks(getattr(cell, "text", "")),
                    "paragraphs": paragraphs,
                    "margin_left_pt": margin_left_pt,
                    "margin_top_pt": margin_top_pt,
                    "margin_right_pt": margin_right_pt,
                    "margin_bottom_pt": margin_bottom_pt,
                    "vertical_anchor": vertical_anchor,
                    "body_pr": body_pr,
                    "fill": fill,
                    "borders": borders or None,
                }
            )
        cells.append(row_cells)

    return {
        "slide": slide_index,
        "shape_name": getattr(shp, "name", None),
        "kind": "ppt_graph_table",
        "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
        **shape_flip_flags(shp),
        "box": shape_box(shp, transform, canvas_emu),
        "rows": rows,
        "cols": cols,
        "table_style_id": table_style_id,
        "col_widths_pt": col_widths_pt,
        "row_heights_pt": row_heights_pt,
        "cells": cells,
    }


def _rgba_from_color_spec(spec: Optional[Dict[str, Any]], default: Optional[tuple[int, int, int, int]] = None) -> Optional[tuple[int, int, int, int]]:
    if not isinstance(spec, dict):
        return default
    rgb = str(spec.get("rgb") or "").strip().lstrip("#")
    if len(rgb) != 6 and str(spec.get("type") or "").lower() == "scheme":
        rgb = _CURRENT_PPT_THEME_COLORS.get(_theme_scheme_alias(spec.get("scheme")), "")
    if len(rgb) != 6:
        return default
    try:
        r = int(rgb[0:2], 16)
        g = int(rgb[2:4], 16)
        b = int(rgb[4:6], 16)
    except Exception:
        return default
    alpha = 255
    if spec.get("alpha_val") is not None:
        try:
            alpha = int(round(max(0.0, min(1.0, float(spec.get("alpha_val")) / 100000.0)) * 255.0))
        except Exception:
            alpha = 255
    h, l, s = colorsys.rgb_to_hls(r / 255.0, g / 255.0, b / 255.0)
    for mod in spec.get("mods") or []:
        if not isinstance(mod, dict):
            continue
        op = str(mod.get("op") or "").lower()
        try:
            val = max(0.0, min(1.0, float(mod.get("val")) / 100000.0))
        except Exception:
            continue
        if op == "alpha":
            alpha = int(round(val * 255.0))
        elif op == "lummod":
            l *= val
        elif op == "lumoff":
            l += val
        elif op == "tint":
            l += (1.0 - l) * val
        elif op == "shade":
            l *= val
    l = max(0.0, min(1.0, l))
    r_f, g_f, b_f = colorsys.hls_to_rgb(h, l, s)
    r = int(round(max(0.0, min(1.0, r_f)) * 255.0))
    g = int(round(max(0.0, min(1.0, g_f)) * 255.0))
    b = int(round(max(0.0, min(1.0, b_f)) * 255.0))
    return r, g, b, alpha


def _positive_parts(values: Any, count: int) -> List[float]:
    out: List[float] = []
    if isinstance(values, list):
        for value in values[:count]:
            try:
                out.append(max(0.0, float(value or 0.0)))
            except Exception:
                out.append(0.0)
    if len(out) < count:
        out.extend([0.0] * (count - len(out)))
    total = sum(out)
    if total <= 0.0 and count > 0:
        return [1.0] * count
    return out


def _table_cell_rects_px(
    layer: Dict[str, Any],
    out_w: int,
    out_h: int,
) -> List[List[tuple[int, int, int, int]]]:
    rows = int(layer.get("rows") or 0)
    cols = int(layer.get("cols") or 0)
    col_parts = _positive_parts(layer.get("col_widths_pt"), cols)
    row_parts = _positive_parts(layer.get("row_heights_pt"), rows)
    col_total = sum(col_parts) or 1.0
    row_total = sum(row_parts) or 1.0

    x_edges = [0]
    acc = 0.0
    for value in col_parts:
        acc += value
        x_edges.append(int(round(acc / col_total * out_w)))
    y_edges = [0]
    acc = 0.0
    for value in row_parts:
        acc += value
        y_edges.append(int(round(acc / row_total * out_h)))
    if x_edges:
        x_edges[-1] = out_w
    if y_edges:
        y_edges[-1] = out_h

    rects: List[List[tuple[int, int, int, int]]] = []
    for ri in range(rows):
        row: List[tuple[int, int, int, int]] = []
        for ci in range(cols):
            row.append((x_edges[ci], y_edges[ri], x_edges[ci + 1], y_edges[ri + 1]))
        rects.append(row)
    return rects


def _table_border_width_px(width_pt: Any, out_w: int, box_emu: Optional[tuple[int, int, int, int]]) -> int:
    try:
        table_w_pt = float(box_emu[2]) / EMU_PER_PT if box_emu else 0.0
        width = float(width_pt or 1.0)
        if table_w_pt > 0:
            return max(1, int(round(width * float(out_w) / table_w_pt)))
    except Exception:
        pass
    return 1


def _table_layer_to_rgba_image(layer: Dict[str, Any], box_emu: Optional[tuple[int, int, int, int]]) -> Optional[Image.Image]:
    rows = int(layer.get("rows") or 0)
    cols = int(layer.get("cols") or 0)
    if rows <= 0 or cols <= 0:
        return None
    out_w, out_h = _raster_size_2048_for_box(box_emu)
    im = Image.new("RGBA", (out_w, out_h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(im, "RGBA")
    rects = _table_cell_rects_px(layer, out_w, out_h)
    cells = layer.get("cells") or []

    for ri, row in enumerate(cells):
        if ri >= rows or not isinstance(row, list):
            continue
        for ci, cell in enumerate(row):
            if ci >= cols or not isinstance(cell, dict):
                continue
            x0, y0, x1, y1 = rects[ri][ci]
            fill_spec = cell.get("fill")
            fill = _rgba_from_color_spec(fill_spec)
            # #region debug-point I:table-cell-fill-raster
            try:
                if fill_spec is not None:
                    _dbg_report_table_gradient_fill(
                        "I",
                        "test.py:_table_layer_to_rgba_image",
                        "[DEBUG] raster table cell fill",
                        {
                            "shape_name": layer.get("shape_name"),
                            "row": ri,
                            "col": ci,
                            "cell_fill": fill_spec,
                            "fill_rgba": fill,
                        },
                    )
            except Exception:
                pass
            # #endregion
            if isinstance(fill_spec, dict) and str(fill_spec.get("type") or "").lower() == "gradient":
                _paint_gradient_fill_rect(im, (x0, y0, x1, y1), fill_spec)
            elif fill is not None:
                draw.rectangle((x0, y0, x1, y1), fill=fill)

    for ri, row in enumerate(cells):
        if ri >= rows or not isinstance(row, list):
            continue
        for ci, cell in enumerate(row):
            if ci >= cols or not isinstance(cell, dict):
                continue
            x0, y0, x1, y1 = rects[ri][ci]
            borders = cell.get("borders") or {}
            for side, pts in (
                ("l", ((x0, y0), (x0, y1))),
                ("r", ((x1, y0), (x1, y1))),
                ("t", ((x0, y0), (x1, y0))),
                ("b", ((x0, y1), (x1, y1))),
            ):
                spec = borders.get(side) if isinstance(borders, dict) else None
                if not isinstance(spec, dict):
                    continue
                color = _rgba_from_color_spec(spec.get("color"), (0, 0, 0, 255))
                width = _table_border_width_px(spec.get("width_pt"), out_w, box_emu)
                draw.line(pts, fill=color, width=width)
    return im


def _save_ppt_graph_table_image_layer(
    *,
    layer: Dict[str, Any],
    shp,
    canvas_emu: Optional[tuple[int, int]],
    assets_dir: Path,
    transform: Optional[Dict[str, float]] = None,
) -> Optional[Dict[str, Any]]:
    base_box_emu = _shape_bounds_box_emu(shp, transform)
    final_box_emu = _shape_box_emu(shp, transform) or base_box_emu
    if base_box_emu is None or final_box_emu is None:
        return None
    im = _table_layer_to_rgba_image(layer, base_box_emu)
    if im is None:
        return None
    shape_id = getattr(shp, "shape_id", None) or 0
    out_path = (assets_dir / f"slide{int(layer.get('slide') or 0):03d}_shape{shape_id}_ppt_graph_table_rendered.png").resolve()
    flips = shape_flip_flags(shp)
    rendered = _save_final_image_layer(
        im=im,
        out_path=out_path,
        final_box_emu=final_box_emu,
        canvas_emu=canvas_emu,
        rotation_deg=shape_rotation_deg(shp, transform),
        flip_h=bool(flips.get("flip_h")),
        flip_v=bool(flips.get("flip_v")),
    )
    if rendered is None:
        return None
    rendered_path, rendered_box_emu = rendered
    placed_box_emu = rendered_box_emu or final_box_emu
    return {
        "slide": layer.get("slide"),
        "shape_name": layer.get("shape_name"),
        "kind": "ppt_graph_table",
        "box": _ratio_box(*placed_box_emu, canvas_emu),
        "saved_path": rendered_path,
        "structure_info": {
            "box": layer.get("box"),
            "rotation_deg": layer.get("rotation_deg"),
            "flip_h": layer.get("flip_h"),
            "flip_v": layer.get("flip_v"),
            "rows": layer.get("rows"),
            "cols": layer.get("cols"),
            "table_style_id": layer.get("table_style_id"),
            "col_widths_pt": layer.get("col_widths_pt"),
            "row_heights_pt": layer.get("row_heights_pt"),
            "cells": layer.get("cells"),
        },
    }


def _table_cell_text_body_pr(cell: Dict[str, Any]) -> Dict[str, Any]:
    attrs: Dict[str, Any] = {"rtlCol": "0", "wrap": "square"}
    body_pr = cell.get("body_pr")
    if isinstance(body_pr, dict) and isinstance(body_pr.get("attrs"), dict):
        for key, value in body_pr.get("attrs", {}).items():
            if value is not None:
                attrs[str(key)] = value
    for key, source_key in (
        ("lIns", "margin_left_pt"),
        ("tIns", "margin_top_pt"),
        ("rIns", "margin_right_pt"),
        ("bIns", "margin_bottom_pt"),
    ):
        try:
            attrs[key] = int(round(float(cell.get(source_key) or 0.0) * EMU_PER_PT))
        except Exception:
            attrs[key] = 0
    anchor = str(cell.get("vertical_anchor") or attrs.get("anchor") or "t").lower()
    if "middle" in anchor or "center" in anchor or "ctr" in anchor:
        attrs["anchor"] = "ctr"
    elif "bottom" in anchor or anchor in {"b", "bot"}:
        attrs["anchor"] = "b"
    else:
        attrs["anchor"] = "t"
    autofit = body_pr.get("autofit") if isinstance(body_pr, dict) else None
    return {"attrs": attrs, "autofit": autofit or "spAutoFit"}


def _cell_has_text(cell: Dict[str, Any]) -> bool:
    if str(cell.get("text") or "").strip():
        return True
    for para in cell.get("paragraphs") or []:
        if not isinstance(para, dict):
            continue
        for run in para.get("runs") or []:
            if isinstance(run, dict) and str(run.get("text") or "").strip():
                return True
    return False


def _table_cell_text_layers(
    layer: Dict[str, Any],
    shp,
    canvas_emu: Optional[tuple[int, int]],
    transform: Optional[Dict[str, float]] = None,
) -> List[Dict[str, Any]]:
    base_box_emu = _shape_bounds_box_emu(shp, transform)
    if base_box_emu is None or canvas_emu is None:
        return []
    left, top, width, height = base_box_emu
    rows = int(layer.get("rows") or 0)
    cols = int(layer.get("cols") or 0)
    if rows <= 0 or cols <= 0:
        return []
    col_parts = _positive_parts(layer.get("col_widths_pt"), cols)
    row_parts = _positive_parts(layer.get("row_heights_pt"), rows)
    col_total = sum(col_parts) or 1.0
    row_total = sum(row_parts) or 1.0
    col_edges = [left]
    acc = 0.0
    for value in col_parts:
        acc += value
        col_edges.append(int(round(left + acc / col_total * width)))
    row_edges = [top]
    acc = 0.0
    for value in row_parts:
        acc += value
        row_edges.append(int(round(top + acc / row_total * height)))
    col_edges[-1] = left + width
    row_edges[-1] = top + height

    out: List[Dict[str, Any]] = []
    cells = layer.get("cells") or []
    for ri, row in enumerate(cells):
        if ri >= rows or not isinstance(row, list):
            continue
        for ci, cell in enumerate(row):
            if ci >= cols or not isinstance(cell, dict) or not _cell_has_text(cell):
                continue
            cell_left = col_edges[ci]
            cell_top = row_edges[ri]
            cell_w = col_edges[ci + 1] - cell_left
            cell_h = row_edges[ri + 1] - cell_top
            body_pr = _table_cell_text_body_pr(cell)
            out.append(
                {
                    "slide": layer.get("slide"),
                    "shape_name": f"{layer.get('shape_name')}_cell_{ri}_{ci}",
                    "kind": "text",
                    "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
                    **shape_flip_flags(shp),
                    "box": _ratio_box(cell_left, cell_top, cell_w, cell_h, canvas_emu),
                    "paragraphs": cell.get("paragraphs") or [
                        {
                            "runs": [
                                {
                                    "run_index": 0,
                                    "text": _normalize_text_breaks(cell.get("text")),
                                }
                            ]
                        }
                    ],
                    "body_pr_anchor": body_pr["attrs"].get("anchor"),
                    "body_pr": body_pr,
                }
            )
    return out


def extract_table_as_image_and_text_layers(
    slide,
    shp,
    slide_index: int,
    assets_dir: Path,
    canvas_emu: Optional[tuple[int, int]],
    transform: Optional[Dict[str, float]] = None,
) -> List[Dict[str, Any]]:
    table_layer = extract_table_layer(slide, shp, slide_index, canvas_emu, transform)
    if table_layer is None:
        return []
    image_layer = _save_ppt_graph_table_image_layer(
        layer=table_layer,
        shp=shp,
        canvas_emu=canvas_emu,
        assets_dir=assets_dir,
        transform=transform,
    )
    if image_layer is None:
        raise RuntimeError(f"ppt_graph_table image export failed for shape={getattr(shp, 'name', None)!r}")
    text_layers = _table_cell_text_layers(table_layer, shp, canvas_emu, transform)
    return [*text_layers, _validate_dual_export_layer(image_layer)]


def extract_slide_background_layers(
    slide, slide_index: int, assets_dir: Path, slide_width_emu: int, slide_height_emu: int
) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    bg_pr, _bg_ref, blips, rel_owner = _resolve_slide_background_source(slide)
    if bg_pr is None and not blips:
        return out
    for idx, blip in enumerate(blips, start=1):
        r_id = blip.get(qn("r:embed"))
        if not r_id:
            continue
        try:
            rel = rel_owner.rels[r_id]
        except KeyError:
            continue
        if getattr(rel, "is_external", False):
            continue
        part = getattr(rel, "target_part", None)
        if part is None or not hasattr(part, "blob"):
            continue
        ctype = getattr(part, "content_type", None)
        if not (ctype and str(ctype).startswith("image/")):
            continue

        base = f"slide{slide_index:03d}_background{idx}"
        ext = ""
        try:
            ext = (getattr(part, "ext", "") or "").lower()
        except Exception:
            ext = ""
        exported = export_asset_image(blob=part.blob, ext=ext, content_type=ctype, out_base=base, assets_dir=assets_dir)
        out.append(
            {
                "slide": slide_index,
                "shape_name": "slide_background",
                "kind": exported["kind"],
                "box": None,
                "saved_path": exported["saved_path"],
                "source": "slide_background",
            }
        )
    gradient_spec = _gradient_fill_spec_from_node(bg_pr) if bg_pr is not None else None
    if gradient_spec is not None:
        gradient_layer = _export_gradient_background_layer(
            gradient_spec,
            slide_index,
            slide_width_emu,
            slide_height_emu,
            assets_dir,
        )
        if gradient_layer is not None:
            out.append(gradient_layer)
    return out


def _background_source_candidates(slide) -> List[tuple[Any, Any]]:
    out: List[tuple[Any, Any]] = []
    seen: set[int] = set()
    for owner in (
        slide,
        getattr(slide, "slide_layout", None),
        getattr(getattr(slide, "slide_layout", None), "slide_master", None),
    ):
        if owner is None:
            continue
        el = getattr(owner, "_element", None)
        part = getattr(owner, "part", None)
        if el is None or part is None:
            continue
        marker = id(el)
        if marker in seen:
            continue
        seen.add(marker)
        out.append((el, part))
    return out


def _resolve_slide_background_source(slide) -> tuple[Any, Any, List[Any], Any]:
    for el, part in _background_source_candidates(slide):
        try:
            bg_pr = el.find(".//p:bg/p:bgPr", NS)
            bg_ref = el.find(".//p:bg/p:bgRef", NS)
            blips = list(el.xpath(".//p:bg//a:blip[@r:embed]"))
        except Exception:
            bg_pr = None
            bg_ref = None
            blips = []
        if bg_pr is not None or bg_ref is not None or blips:
            return bg_pr, bg_ref, blips, part
    return None, None, [], getattr(slide, "part", None)


def extract_slide_background_fill(slide) -> Optional[Dict[str, Any]]:
    bg_pr, bg_ref, _blips, _rel_owner = _resolve_slide_background_source(slide)
    if bg_pr is None and bg_ref is None:
        return None

    result: Dict[str, Any] = {}

    if bg_pr is not None:
        solid_spec = _background_fill_from_color_spec(_solid_fill_spec_from_node(bg_pr))
        if solid_spec is not None:
            return solid_spec

        gradient_spec = _gradient_fill_spec_from_node(bg_pr)
        if gradient_spec is not None:
            return None

        return result or None

    # Background may be a theme reference rather than a concrete bgPr.
    if bg_ref is not None:
        try:
            bg_ref_root = ET.fromstring(ET.tostring(bg_ref, encoding="utf-8"))
            color_spec = _color_spec_from_color_el(bg_ref_root)
        except Exception:
            color_spec = None
        solid_spec = _background_fill_from_color_spec(color_spec)
        if solid_spec is not None:
            return solid_spec
        result.update({"type": "ref"})
    return result or None


def _background_rgb_from_fill(fill: Optional[Dict[str, Any]]) -> str:
    if not isinstance(fill, dict):
        return "FFFFFF"
    if str(fill.get("type") or "").lower() != "solid":
        return "FFFFFF"
    color_type = str(fill.get("color_type") or "").lower()
    spec: Dict[str, Any] = {}
    if color_type in {"rgb", "srgb"} and fill.get("rgb"):
        spec = {"type": "srgb", "rgb": str(fill.get("rgb")).upper()}
    elif color_type == "scheme" and fill.get("scheme"):
        spec = {"type": "scheme", "scheme": fill.get("scheme")}
    else:
        return "FFFFFF"
    if fill.get("mods"):
        spec["mods"] = list(fill.get("mods") or [])
    rgba = _rgba_from_color_spec(spec)
    if rgba is None:
        return "FFFFFF"
    return f"{int(rgba[0]):02X}{int(rgba[1]):02X}{int(rgba[2]):02X}"


def _slide_canvas_asset_path(slide_index: int, assets_dir: Path) -> Path:
    return assets_dir / f"slide{slide_index:03d}_canvas_image_png_rgba_rendered.png"


def _materialize_slide_canvas_asset(slide_index: int, assets_dir: Path, source_path: Any) -> Optional[str]:
    src_text = str(source_path or "").strip()
    if not src_text:
        return None
    src = Path(src_text)
    dst = _slide_canvas_asset_path(slide_index, assets_dir)
    try:
        ensure_dir(dst.parent)
        if src.exists() and src.resolve() != dst.resolve():
            shutil.copyfile(src, dst)
            try:
                src.unlink()
            except Exception:
                pass
        elif src.exists() and src.resolve() == dst.resolve():
            pass
        else:
            return src_text
    except Exception:
        return src_text
    return str(dst)


def _remove_asset_file(path_value: Any) -> None:
    path_text = str(path_value or "").strip()
    if not path_text:
        return
    try:
        path = Path(path_text)
        if path.exists() and path.is_file():
            path.unlink()
    except Exception:
        pass


def _asset_is_fully_opaque(path_value: Any) -> bool:
    path_text = str(path_value or "").strip()
    if not path_text:
        return False
    try:
        with Image.open(path_text) as im:
            bands = im.getbands()
            if "A" not in bands:
                return True
            alpha = im.getchannel("A")
            extrema = alpha.getextrema()
            if not isinstance(extrema, tuple) or len(extrema) != 2:
                return False
            return int(extrema[0]) >= 255 and int(extrema[1]) >= 255
    except Exception:
        return False


def _box_covers_slide(box: Any, *, tol: float = 1e-6) -> bool:
    if not isinstance(box, dict):
        return False
    try:
        left = float(box.get("left") or 0.0)
        top = float(box.get("top") or 0.0)
        width = float(box.get("width") or 0.0)
        height = float(box.get("height") or 0.0)
    except Exception:
        return False
    right = left + width
    bottom = top + height
    return left <= tol and top <= tol and right >= (1.0 - tol) and bottom >= (1.0 - tol)


def extract_slide_canvas_layer(
    slide_index: int,
    slide_width_emu: int,
    slide_height_emu: int,
    background_fill: Optional[Dict[str, Any]] = None,
    assets_dir: Optional[Path] = None,
) -> Dict[str, Any]:
    """
    Add a synthetic bottom-most layer representing the full slide canvas.
    This preserves the original page size even if no visible layer reaches the slide edges.
    """
    rgb = _background_rgb_from_fill(background_fill)
    out_w, out_h = _raster_size_2048_for_box((0, 0, int(slide_width_emu), int(slide_height_emu)))
    out_path = _slide_canvas_asset_path(slide_index, assets_dir or Path("."))
    Image.new(
        "RGBA",
        (out_w, out_h),
        (
            int(rgb[0:2], 16),
            int(rgb[2:4], 16),
            int(rgb[4:6], 16),
            255,
        ),
    ).save(out_path, format="PNG")
    return {
        "slide": slide_index,
        "shape_name": "slide_canvas",
        "kind": "slide_canvas",
        "box": {
            "left": 0.0,
            "top": 0.0,
            "width": 1.0,
            "height": 1.0,
        },
        "canvas_width_emu": int(slide_width_emu),
        "canvas_height_emu": int(slide_height_emu),
        "saved_path": str(out_path),
        "background_fill": dict(background_fill) if background_fill else None,
    }


def _promote_background_to_canvas(
    slide_layers: List[Dict[str, Any]],
    slide_index: int,
    slide_width_emu: int,
    slide_height_emu: int,
    assets_dir: Path,
    background_fill: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    for idx in range(len(slide_layers) - 1, -1, -1):
        layer = slide_layers[idx]
        if not isinstance(layer, dict):
            continue
        kind = str(layer.get("kind") or "")
        if kind not in IMAGE_LIKE_EXPORT_KINDS:
            continue
        if str(layer.get("source") or "") in {"slide_background", "slide_background_gradient"}:
            continue
        if not layer.get("saved_path"):
            continue
        if not _box_covers_slide(layer.get("box")):
            continue
        if not _asset_is_fully_opaque(layer.get("saved_path")):
            continue
        promoted_saved_path = str(layer.get("saved_path"))
        stale_canvas_path = _slide_canvas_asset_path(slide_index, assets_dir)
        try:
            if stale_canvas_path.resolve() != Path(promoted_saved_path).resolve():
                _remove_asset_file(stale_canvas_path)
        except Exception:
            _remove_asset_file(stale_canvas_path)
        promoted = {
            "slide": slide_index,
            "shape_name": "slide_canvas",
            "kind": "slide_canvas",
            "box": {"left": 0.0, "top": 0.0, "width": 1.0, "height": 1.0},
            "canvas_width_emu": int(slide_width_emu),
            "canvas_height_emu": int(slide_height_emu),
            "saved_path": promoted_saved_path,
            "background_fill": dict(background_fill) if background_fill else None,
        }
        slide_layers.pop(idx)
        for bg_idx in range(len(slide_layers) - 1, idx - 1, -1):
            bg_layer = slide_layers[bg_idx]
            if not isinstance(bg_layer, dict):
                continue
            if str(bg_layer.get("source") or "") not in {"slide_background", "slide_background_gradient"}:
                continue
            if not bg_layer.get("saved_path"):
                continue
            if not _box_covers_slide(bg_layer.get("box")):
                continue
            _remove_asset_file(bg_layer.get("saved_path"))
            slide_layers.pop(bg_idx)
        return promoted

    for idx in range(len(slide_layers) - 1, -1, -1):
        layer = slide_layers[idx]
        if not isinstance(layer, dict):
            continue
        if str(layer.get("source") or "") not in {"slide_background", "slide_background_gradient"}:
            continue
        if not layer.get("saved_path"):
            continue
        if layer.get("box") is not None and not _box_covers_slide(layer.get("box")):
            continue
        promoted = {
            "slide": slide_index,
            "shape_name": "slide_canvas",
            "kind": "slide_canvas",
            "box": {"left": 0.0, "top": 0.0, "width": 1.0, "height": 1.0},
            "canvas_width_emu": int(slide_width_emu),
            "canvas_height_emu": int(slide_height_emu),
            "saved_path": _materialize_slide_canvas_asset(slide_index, assets_dir, layer.get("saved_path")) or layer.get("saved_path"),
            "background_fill": dict(background_fill) if background_fill else None,
        }
        slide_layers.pop(idx)
        return promoted
    return extract_slide_canvas_layer(slide_index, slide_width_emu, slide_height_emu, background_fill, assets_dir)


def _renumber_slide_layer_shape_names(slide_layers: List[Dict[str, Any]]) -> None:
    for idx, layer in enumerate(slide_layers, start=1):
        if isinstance(layer, dict):
            layer["shape_name"] = idx


def extract_embedded_fonts_from_pptx(pptx_path: Path, fonts_dir: Path) -> List[Dict[str, Any]]:
    """
    Extract embedded font binaries from a .pptx (OOXML zip package).

    Notes:
    - PowerPoint commonly stores embedded fonts under 'ppt/fonts/' as .fntdata/.odttf.
    - These are saved as-is (no de-obfuscation / conversion).
    """
    ensure_dir(fonts_dir)
    prefixes = ("ppt/fonts/", "ppt/embeddings/", "word/fonts/", "xl/fonts/")

    # Regenerate fonts output each run to avoid mixing old "font7.fntdata" style outputs with
    # new "Typeface__style.*" outputs.
    try:
        for p in fonts_dir.iterdir():
            if p.is_file():
                p.unlink()
    except Exception:
        pass

    def _safe_filename(name: str, fallback: str) -> str:
        # Keep filenames reasonably safe across platforms.
        # We allow Unicode (font names can be non-ASCII) but remove path separators and other invalid chars.
        s = (name or "").strip()
        if not s:
            s = fallback
        bad = '<>:"/\\|?*\0'
        for ch in bad:
            s = s.replace(ch, "_")
        s = "".join(c if (c >= " " and c != "\x7f") else "_" for c in s)
        s = " ".join(s.split())
        # Avoid super long filenames.
        if len(s) > 120:
            s = s[:120].rstrip()
        return s or fallback

    def _write_unique(base_name: str, data: bytes) -> Path:
        out_path = fonts_dir / base_name
        if out_path.exists():
            stem, ext = out_path.stem, out_path.suffix
            i = 2
            while (fonts_dir / f"{stem}_{i}{ext}").exists():
                i += 1
            out_path = fonts_dir / f"{stem}_{i}{ext}"
        out_path.write_bytes(data)
        return out_path

    extracted: List[Dict[str, Any]] = []
    extracted_zip_paths: set[str] = set()

    # Prefer the authoritative embedded-font list in presentation.xml (typeface + style + rId).
    # This avoids missing fonts stored as .fntdata/.odttf with unexpected naming/extensions.
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    r_id_attr = f"{{{NS_R}}}id"

    with zipfile.ZipFile(pptx_path, "r") as zf:
        try:
            pres_xml = zf.read("ppt/presentation.xml")
            pres_rels_xml = zf.read("ppt/_rels/presentation.xml.rels")
            pres_root = ET.fromstring(pres_xml)
            rels_root = ET.fromstring(pres_rels_xml)

            rid_to_rel: Dict[str, Dict[str, str]] = {}
            for rel in list(rels_root):
                rid = rel.get("Id")
                if not rid:
                    continue
                rid_to_rel[rid] = {"type": rel.get("Type") or "", "target": rel.get("Target") or ""}

            embedded_font_lst = pres_root.find(f".//{{{NS_P}}}embeddedFontLst")
            if embedded_font_lst is not None:
                for ef in embedded_font_lst.findall(f"./{{{NS_P}}}embeddedFont"):
                    font_el = ef.find(f"./{{{NS_P}}}font")
                    typeface = font_el.get("typeface") if font_el is not None else None

                    for style_tag in ("regular", "bold", "italic", "boldItalic"):
                        style_el = ef.find(f"./{{{NS_P}}}{style_tag}")
                        if style_el is None:
                            continue
                        rid = style_el.get(r_id_attr)
                        if not rid:
                            continue
                        rel = rid_to_rel.get(rid) or {}
                        target = rel.get("target") or ""
                        if not target:
                            continue
                        # Relationships are relative to `ppt/presentation.xml`, so targets are under `ppt/`.
                        zip_path = f"ppt/{target.lstrip('/')}"
                        try:
                            data = zf.read(zip_path)
                        except KeyError:
                            continue

                        # Name by font face instead of "font7.fntdata" etc.
                        suffix = Path(zip_path).suffix or ".bin"
                        face = _safe_filename(typeface or "", "unknown_font")
                        style = _safe_filename(style_tag, "regular")
                        base_name = f"{face}__{style}{suffix}"
                        out_path = _write_unique(base_name, data)
                        extracted_zip_paths.add(zip_path)
                        extracted.append(
                            {
                                "zip_path": zip_path,
                                "saved_path": str(out_path),
                                "size_bytes": len(data),
                                "typeface": typeface,
                                "style": style_tag,
                                "rid": rid,
                                "rel_type": rel.get("type"),
                                "rel_target": target,
                                "source": "presentation_embeddedFontLst",
                            }
                        )
        except Exception:
            # Keep going with best-effort extraction by folder scan.
            pass

        # Fallback: export any files that live under the typical font directories.
        # This catches unreferenced font parts (rare) and non-standard extensions.
        for zi in zf.infolist():
            zip_name = zi.filename or ""
            zip_name_l = zip_name.lower()
            if zip_name_l.endswith("/"):
                continue
            if not any(zip_name_l.startswith(p) for p in prefixes):
                continue
            if zip_name in extracted_zip_paths:
                continue

            data = zf.read(zip_name)
            out_path = _write_unique(Path(zip_name).name, data)
            extracted_zip_paths.add(zip_name)
            extracted.append(
                {
                    "zip_path": zip_name,
                    "saved_path": str(out_path),
                    "size_bytes": len(data),
                    "source": "package_scan",
                }
            )

    return extracted


def iter_export_layers_top_to_bottom(
    slide,
    shapes,
    slide_index: int,
    assets_dir: Path,
    transform: Optional[Dict[str, float]] = None,
) -> Iterable[Dict[str, Any]]:
    """
    Iterate shapes from top -> bottom (front -> back).
    Group shapes are flattened into slide-level layers.
    """
    transform = transform or identity_transform()
    # python-pptx exposes shapes typically in back -> front order; reverse to get top -> bottom.
    for shp in reversed(list(shapes)):
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            next_transform = child_transform_for_group(shp, transform)
            for item in iter_export_layers_top_to_bottom(
                slide, shp.shapes, slide_index, assets_dir, next_transform
            ):
                yield item
            continue

        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield extract_picture_asset(slide, shp, slide_index, assets_dir, transform)
            continue

        table_layers = extract_table_as_image_and_text_layers(slide, shp, slide_index, assets_dir, _slide_canvas_emu(slide), transform)
        if table_layers:
            for table_layer in table_layers:
                yield table_layer
            continue

        text_layer = None
        if getattr(shp, "has_text_frame", False):
            text_layer = extract_text_layer(shp, slide_index, _slide_canvas_emu(slide), transform)

        graph_layer = extract_graph_layer(slide, shp, slide_index, assets_dir, transform)
        if graph_layer is not None and text_layer is not None and graph_layer.get("kind") == "ppt_graph_geo":
            shape_xml = graph_layer.get("shape_xml") or {}
            geom = shape_xml.get("geom") if isinstance(shape_xml, dict) else None
            if isinstance(geom, dict) and geom.get("type") == "prstGeom":
                graph_layer["text_content"] = {
                    "text": text_layer.get("text"),
                    "paragraphs": text_layer.get("paragraphs"),
                    "text_box": text_layer.get("text_box"),
                }
                yield graph_layer
                graph_layer = None
                text_layer = None

        if text_layer is not None:
            yield text_layer
        if graph_layer is not None:
            yield graph_layer

        # Also extract images used in shape fills (common in templates).
        # Text is drawn over fill, so yield fill after text to keep top->bottom semantics.
        for img_layer in extract_fill_image_assets(slide, shp, slide_index, assets_dir, transform):
            yield img_layer

        # Skip other layer types by default.
        continue


def main() -> None:
    global _CURRENT_PPT_THEME_COLORS
    ap = argparse.ArgumentParser(description="Extract PPTX layers: images(SVG/PNG RGBA) + editable text styles to JSON.")
    # Positional args are required by default. Use nargs='?' to allow running with defaults.
    ap.add_argument(
        "pptx",
        nargs="?",
        default="/Users/bytedance/Documents/test/code/PPTX格式/海报/~$Beige Vintage Newspaper Travel to Italy Poster.pptx",
        help="Input .pptx file path",
    )
    ap.add_argument("-o", "--out", help="Output folder", default="/Users/bytedance/Documents/test/code/PPTX格式/res")
    ap.add_argument("--json-name", default="layers.json", help="Output JSON file name (default: layers.json)")
    args = ap.parse_args()
    pptx_path = Path(args.pptx).expanduser().resolve()

    # Office often creates a temp/lock file like "~$Foo.pptx". It's not a valid pptx package to read.
    # If user passes such a path, try to locate the real file by stripping the "~$" prefix.
    if pptx_path.name.startswith("~$"):
        candidate = pptx_path.with_name(pptx_path.name[2:])
        if candidate.exists():
            pptx_path = candidate.resolve()
    fixs = "/".join(pptx_path.parts[-2:]) if len(pptx_path.parts) >= 2 else pptx_path.name
    out_dir = Path(args.out).expanduser().resolve() / fixs
    if not pptx_path.exists():
        # "~$*.pptx" is usually an Office temp/lock file; remind user to pick the real pptx.
        if pptx_path.name.startswith("~$"):
            raise SystemExit(
                f"Input not found (looks like Office temp file '~$'). "
                f"Please pass the real .pptx (usually same name without '~$'): {pptx_path}"
            )
        raise SystemExit(f"Input not found: {pptx_path}")

    ensure_dir(out_dir)
    assets_dir = out_dir / "assets"
    ensure_dir(assets_dir)
    fonts_dir = out_dir / "fonts"

    prs = Presentation(str(pptx_path))
    _CURRENT_PPT_THEME_COLORS = _extract_theme_colors_from_pptx(pptx_path)
    layers: List[Dict[str, Any]] = []

    for si, slide in enumerate(prs.slides, start=1):
        per_slide_index = 0
        slide_layers: List[Dict[str, Any]] = []
        # Shapes: top -> bottom
        for layer in iter_export_layers_top_to_bottom(slide, slide.shapes, si, assets_dir):
            per_slide_index += 1
            layer["shape_name"] = per_slide_index
            slide_layers.append(layer)

        # Background image is the bottom-most layer, append at the end (bottom).
        for layer in extract_slide_background_layers(slide, si, assets_dir, int(prs.slide_width), int(prs.slide_height)):
            per_slide_index += 1
            layer["shape_name"] = per_slide_index
            slide_layers.append(layer)

        # Synthetic slide canvas is the absolute bottom-most layer.
        background_fill = extract_slide_background_fill(slide)
        canvas_layer = _promote_background_to_canvas(
            slide_layers,
            si,
            int(prs.slide_width),
            int(prs.slide_height),
            assets_dir,
            background_fill,
        )
        per_slide_index += 1
        canvas_layer["shape_name"] = per_slide_index
        slide_layers.append(canvas_layer)
        _renumber_slide_layer_shape_names(slide_layers)
        layers.extend(slide_layers)

    out_json = out_dir / args.json_name
    # Protocol: group by slide (outer list), and omit per-layer `slide` field.
    # Example:
    #   [[{layer},{layer}], [{layer},...], ...]
    by_slide: Dict[int, List[Dict[str, Any]]] = {}
    for layer in layers:
        try:
            slide_idx = int(layer.get("slide") or 1)
        except Exception:
            slide_idx = 1
        normalized = normalize_export_layer_to_template(dict(layer))
        normalized.pop("slide", None)
        by_slide.setdefault(slide_idx, []).append(normalized)
    grouped_layers: List[List[Dict[str, Any]]] = [by_slide[k] for k in sorted(by_slide.keys())]

    out_json.write_text(
        json.dumps(_encode_scaled_numbers_for_json(grouped_layers, 1000), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    fonts = extract_embedded_fonts_from_pptx(pptx_path, fonts_dir)
    if fonts:
        (out_dir / "fonts.json").write_text(json.dumps(fonts, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"Saved {len(layers)} layers ({len(grouped_layers)} slide group(s)) to: {out_json}")
    print(f"Assets saved to: {assets_dir}")
    if fonts:
        print(f"Fonts extracted: {len(fonts)} -> {fonts_dir}")
    else:
        print("Fonts extracted: 0 (no embedded fonts found in package)")


if __name__ == "__main__":
    main()

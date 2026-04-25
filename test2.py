#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Export a PPTX into the current layer JSON protocol.

- Text layers keep paragraph/run style data.
- Image layers are baked to final RGBA PNG files plus final ratio boxes.
- Shape layers keep enough geometry/XML to rebuild the slide.

Usage:
  pip install python-pptx pillow
  python test.py /path/to/input.pptx -o /path/to/out_dir
"""

import argparse
import io
import importlib
import json
import math
import os
import re
import sys
import time
import zipfile
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from PIL import Image, ImageDraw, ImageOps


EMU_PER_PT = 12700  # PowerPoint uses EMU; 1 pt = 12700 EMU
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "asvg": "http://schemas.microsoft.com/office/drawing/2016/SVG/main",
}


def _canvas_base_emu(canvas_emu: Optional[tuple[int, int]]) -> Optional[float]:
    if not canvas_emu:
        return None
    try:
        w = float(canvas_emu[0] or 0.0)
        h = float(canvas_emu[1] or 0.0)
        if w <= 0 or h <= 0:
            return None
        return math.sqrt(w * h)
    except Exception:
        return None


def _emu_rel(value_emu: Optional[float], canvas_emu: Optional[tuple[int, int]]) -> Optional[float]:
    base = _canvas_base_emu(canvas_emu)
    if value_emu is None or not base:
        return None
    try:
        return float(value_emu) / float(base)
    except Exception:
        return None


def _pt_rel(value_pt: Optional[float], canvas_emu: Optional[tuple[int, int]]) -> Optional[float]:
    if value_pt is None:
        return None
    try:
        return _emu_rel(float(value_pt) * EMU_PER_PT, canvas_emu)
    except Exception:
        return None


def _raw_100pt_rel(raw_100pt: Optional[int], canvas_emu: Optional[tuple[int, int]]) -> Optional[float]:
    if raw_100pt is None:
        return None
    try:
        return _emu_rel((float(raw_100pt) / 100.0) * EMU_PER_PT, canvas_emu)
    except Exception:
        return None


# #region debug-point A:report
def _dbg_report_method_ppt_style(hypothesis_id: str, location: str, msg: str, data: Optional[Dict[str, Any]] = None) -> None:
    try:
        import json as _json
        import urllib.request as _urlreq

        _u, _s, _r = "http://127.0.0.1:7777/event", "method-ppt-style", os.environ.get("DEBUG_RUN_ID", "pre-fix")
        try:
            with open(".dbg/method-ppt-style.env", "r", encoding="utf-8") as _f:
                for _line in _f.read().splitlines():
                    if _line.startswith("DEBUG_SERVER_URL="):
                        _u = _line.split("=", 1)[1] or _u
                    elif _line.startswith("DEBUG_SESSION_ID="):
                        _s = _line.split("=", 1)[1] or _s
                    elif _line.startswith("DEBUG_RUN_ID="):
                        _r = _line.split("=", 1)[1] or _r
        except Exception:
            pass

        _payload = {
            "sessionId": _s,
            "runId": _r,
            "hypothesisId": hypothesis_id,
            "location": location,
            "msg": msg,
            "data": data or {},
            "ts": int(time.time() * 1000),
        }
        _urlreq.urlopen(
            _urlreq.Request(
                _u,
                data=_json.dumps(_payload).encode("utf-8"),
                headers={"Content-Type": "application/json"},
            ),
            timeout=1.0,
        ).read()
    except Exception:
        pass
# #endregion


def identity_transform() -> Dict[str, float]:
    return {"a": 1.0, "b": 0.0, "c": 0.0, "d": 1.0, "tx": 0.0, "ty": 0.0}


def compose_transform(parent: Dict[str, float], child: Dict[str, float]) -> Dict[str, float]:
    return {
        "a": parent["a"] * child["a"] + parent["b"] * child["c"],
        "b": parent["a"] * child["b"] + parent["b"] * child["d"],
        "c": parent["c"] * child["a"] + parent["d"] * child["c"],
        "d": parent["c"] * child["b"] + parent["d"] * child["d"],
        "tx": parent["a"] * child["tx"] + parent["b"] * child["ty"] + parent["tx"],
        "ty": parent["c"] * child["tx"] + parent["d"] * child["ty"] + parent["ty"],
    }


def transform_point(transform: Dict[str, float], x: float, y: float) -> tuple[float, float]:
    return (
        transform["a"] * x + transform["b"] * y + transform["tx"],
        transform["c"] * x + transform["d"] * y + transform["ty"],
    )


def transform_vector(transform: Dict[str, float], x: float, y: float) -> tuple[float, float]:
    return (
        transform["a"] * x + transform["b"] * y,
        transform["c"] * x + transform["d"] * y,
    )


def transform_rotation_deg(transform: Dict[str, float]) -> float:
    return math.degrees(math.atan2(transform["c"], transform["a"])) % 360.0


def ensure_dir(p: Path) -> None:
    p.mkdir(parents=True, exist_ok=True)


def emu_to_pt(v: Optional[int]) -> Optional[float]:
    if v is None:
        return None
    try:
        return float(v) / EMU_PER_PT
    except Exception:
        return None


def _ratio_box(
    left: Optional[int],
    top: Optional[int],
    width: Optional[int],
    height: Optional[int],
    canvas_emu: Optional[tuple[int, int]],
) -> Dict[str, Any]:
    def _round_coord(v: float) -> float:
        r = round(float(v), 5)
        return 0.0 if r == -0.0 else r

    if None in (left, top, width, height) or not canvas_emu:
        return {"left": None, "top": None, "width": None, "height": None}
    canvas_width, canvas_height = canvas_emu
    if canvas_width <= 0 or canvas_height <= 0:
        return {"left": None, "top": None, "width": None, "height": None}
    return {
        "left": _round_coord(float(left) / float(canvas_width)),
        "top": _round_coord(float(top) / float(canvas_height)),
        "width": _round_coord(float(width) / float(canvas_width)),
        "height": _round_coord(float(height) / float(canvas_height)),
    }


def _round_rotation_deg(v: float) -> float:
    """Round exported rotation degrees to 2 decimals (protocol requirement)."""
    try:
        r = round(float(v), 2)
    except Exception:
        return 0.0
    return 0.0 if r == -0.0 else r


def _encode_scaled_numbers_for_json(x: Any, scale: int = 1000) -> Any:
    """Encode protocol floats as scaled integers for JSON output."""
    if isinstance(x, float):
        try:
            v = int(round(float(x) * float(scale)))
        except Exception:
            v = 0
        return 0 if v == -0 else v
    if isinstance(x, list):
        return [_encode_scaled_numbers_for_json(i, scale) for i in x]
    if isinstance(x, dict):
        return {k: _encode_scaled_numbers_for_json(v, scale) for k, v in x.items()}
    return x


def _shape_bounds_box_emu(shp, transform: Optional[Dict[str, float]] = None) -> Optional[tuple[int, int, int, int]]:
    # left/top/width/height might not exist for some shape types; keep it defensive.
    # For shapes inside groups, `python-pptx` exposes local coordinates; apply the accumulated
    # group transform so every exported box is in slide absolute coordinates.
    transform = transform or identity_transform()
    left = getattr(shp, "left", None)
    top = getattr(shp, "top", None)
    width = getattr(shp, "width", None)
    height = getattr(shp, "height", None)
    if None not in (left, top, width, height):
        cx, cy = transform_point(transform, float(left) + float(width) / 2.0, float(top) + float(height) / 2.0)
        vx = transform_vector(transform, float(width), 0.0)
        vy = transform_vector(transform, 0.0, float(height))
        abs_width = math.hypot(vx[0], vx[1])
        abs_height = math.hypot(vy[0], vy[1])
        left = int(round(cx - abs_width / 2.0))
        top = int(round(cy - abs_height / 2.0))
        width = int(round(abs_width))
        height = int(round(abs_height))
    if None in (left, top, width, height):
        return None
    return int(left), int(top), int(width), int(height)


def shape_bounds_box(
    shp, transform: Optional[Dict[str, float]] = None, canvas_emu: Optional[tuple[int, int]] = None
) -> Dict[str, Any]:
    box_emu = _shape_bounds_box_emu(shp, transform)
    if box_emu is None:
        return _ratio_box(None, None, None, None, canvas_emu)
    left, top, width, height = box_emu
    return _ratio_box(left, top, width, height, canvas_emu)


def _shape_box_emu(shp, transform: Optional[Dict[str, float]] = None) -> Optional[tuple[int, int, int, int]]:
    bounds_box = _shape_bounds_box_emu(shp, transform)
    if bounds_box is None:
        return None
    left, top, width, height = bounds_box

    # Export the visual axis-aligned bounding box after rotation.
    rotation_deg = shape_rotation_deg(shp, transform)
    if rotation_deg:
        theta = math.radians(rotation_deg)
        bbox_w = abs(width * math.cos(theta)) + abs(height * math.sin(theta))
        bbox_h = abs(width * math.sin(theta)) + abs(height * math.cos(theta))
        cx = left + width / 2.0
        cy = top + height / 2.0
        left = int(round(cx - bbox_w / 2.0))
        top = int(round(cy - bbox_h / 2.0))
        width = int(round(bbox_w))
        height = int(round(bbox_h))
    return int(left), int(top), int(width), int(height)


def shape_box(shp, transform: Optional[Dict[str, float]] = None, canvas_emu: Optional[tuple[int, int]] = None) -> Dict[str, Any]:
    box_emu = _shape_box_emu(shp, transform)
    if box_emu is None:
        return _ratio_box(None, None, None, None, canvas_emu)
    left, top, width, height = box_emu
    return _ratio_box(left, top, width, height, canvas_emu)


def child_transform_for_group(group_shp, parent_transform: Optional[Dict[str, float]] = None) -> Dict[str, float]:
    """
    Build a transform that maps coordinates in this group's child space into slide coordinates.

    OOXML group transform:
      parent_space = off + (child_space - chOff) * (ext / chExt)
    We compose that with the parent transform so nested groups are handled correctly.
    """
    parent_transform = parent_transform or identity_transform()
    xfrm = getattr(getattr(getattr(group_shp, "_element", None), "grpSpPr", None), "xfrm", None)
    if xfrm is None:
        # Fallback: if internals differ, at least offset by the group's own box.
        left = getattr(group_shp, "left", 0) or 0
        top = getattr(group_shp, "top", 0) or 0
        return compose_transform(
            parent_transform,
            {"a": 1.0, "b": 0.0, "c": 0.0, "d": 1.0, "tx": float(left), "ty": float(top)},
        )

    off_x = float(getattr(xfrm.off, "x", 0) or 0)
    off_y = float(getattr(xfrm.off, "y", 0) or 0)
    ext_cx = float(getattr(xfrm.ext, "cx", 0) or 0)
    ext_cy = float(getattr(xfrm.ext, "cy", 0) or 0)
    ch_off_x = float(getattr(xfrm.chOff, "x", 0) or 0)
    ch_off_y = float(getattr(xfrm.chOff, "y", 0) or 0)
    ch_ext_cx = float(getattr(xfrm.chExt, "cx", 0) or 0)
    ch_ext_cy = float(getattr(xfrm.chExt, "cy", 0) or 0)

    scale_x = (ext_cx / ch_ext_cx) if ch_ext_cx else 1.0
    scale_y = (ext_cy / ch_ext_cy) if ch_ext_cy else 1.0
    rot_deg = float(xfrm.get("rot", "0") or 0) / 60000.0
    theta = math.radians(rot_deg)
    cos_t = math.cos(theta)
    sin_t = math.sin(theta)

    tx0 = off_x - ch_off_x * scale_x
    ty0 = off_y - ch_off_y * scale_y
    center_x = off_x + ext_cx / 2.0
    center_y = off_y + ext_cy / 2.0

    group_transform = {
        "a": cos_t * scale_x,
        "b": -sin_t * scale_y,
        "c": sin_t * scale_x,
        "d": cos_t * scale_y,
        "tx": cos_t * tx0 - sin_t * ty0 + center_x - cos_t * center_x + sin_t * center_y,
        "ty": sin_t * tx0 + cos_t * ty0 + center_y - sin_t * center_x - cos_t * center_y,
    }
    return compose_transform(parent_transform, group_transform)


def font_color_to_dict(font) -> Optional[Dict[str, Any]]:
    c = getattr(font, "color", None)
    if c is None:
        return None
    try:
        ctype = c.type
    except Exception:
        return None

    out: Dict[str, Any] = {"type": None, "rgb": None}
    if ctype == MSO_COLOR_TYPE.RGB and c.rgb is not None:
        out["type"] = "RGB"
        out["rgb"] = str(c.rgb)  # e.g. 'FF00AA'
    elif (str(ctype).endswith("THEME") or str(ctype) == "SCHEME") and getattr(c, "theme_color", None) is not None:
        out["type"] = "THEME"
    else:
        out["type"] = str(ctype)
    return out


def _text_color_from_run(run) -> Optional[Dict[str, Any]]:
    try:
        r_pr = run._r.find(qn("a:rPr"))
    except Exception:
        r_pr = None
    if r_pr is None:
        return None
    spec = _solid_fill_spec_from_node(r_pr)
    if not spec:
        return None
    out: Dict[str, Any] = {"type": None, "rgb": None}
    if spec.get("type") == "srgb" and spec.get("rgb"):
        out["type"] = "RGB"
        out["rgb"] = spec.get("rgb")
    elif spec.get("type") == "scheme" and spec.get("scheme"):
        out["type"] = "SCHEME"
        out["scheme"] = spec.get("scheme")
    if spec.get("mods"):
        out["mods"] = spec.get("mods")
    return out


def shape_rotation_deg(shp, transform: Optional[Dict[str, float]] = None) -> float:
    try:
        local_rot = float(getattr(shp, "rotation", 0.0) or 0.0)
    except Exception:
        local_rot = 0.0
    base_rot = transform_rotation_deg(transform or identity_transform())
    return (base_rot + local_rot) % 360.0


def shape_flip_flags(shp) -> Dict[str, bool]:
    try:
        xfrm = getattr(getattr(shp, "_element", None), "spPr", None)
        xfrm = getattr(xfrm, "xfrm", None)
        if xfrm is None:
            xfrm = getattr(getattr(getattr(shp, "_element", None), "grpSpPr", None), "xfrm", None)
        if xfrm is None:
            return {"flip_h": False, "flip_v": False}
        return {
            "flip_h": str(xfrm.get("flipH", "false")).lower() in {"1", "true"},
            "flip_v": str(xfrm.get("flipV", "false")).lower() in {"1", "true"},
        }
    except Exception:
        return {"flip_h": False, "flip_v": False}


def _resolve_image_parts(slide, blip):
    def _rel_to_part(r_id: Optional[str]):
        if not r_id:
            return None
        try:
            rel = slide.part.rels[r_id]
        except KeyError:
            return None
        if getattr(rel, "is_external", False):
            return None
        part = getattr(rel, "target_part", None)
        if part is None or not hasattr(part, "blob"):
            return None
        ctype = getattr(part, "content_type", None)
        if not (ctype and str(ctype).startswith("image/")):
            return None
        return part

    svg_rids: List[str] = []
    try:
        blip_et = ET.fromstring(ET.tostring(blip, encoding="utf-8"))
        svg_blip = blip_et.find(".//asvg:svgBlip", NS)
        if svg_blip is not None:
            rid = svg_blip.get(qn("r:embed"))
            if rid:
                svg_rids.append(rid)
    except Exception:
        svg_rids = []

    primary_part = None
    for svg_rid in svg_rids:
        part = _rel_to_part(svg_rid)
        if part is not None:
            primary_part = part
            break

    preview_part = _rel_to_part(blip.get(qn("r:embed")))
    if primary_part is not None:
        return primary_part, preview_part
    return preview_part, None


def _xml_spacing_to_dict(parent, tag_local: str) -> Optional[Dict[str, Any]]:
    if parent is None:
        return None
    node = parent.find(qn(f"a:{tag_local}"))
    if node is None:
        return None
    pts = node.find(qn("a:spcPts"))
    if pts is not None:
        try:
            raw = int(pts.get("val"))
            return {"mode": "points", "raw": raw}
        except Exception:
            return None
    pct = node.find(qn("a:spcPct"))
    if pct is not None:
        try:
            raw = int(pct.get("val"))
            return {"mode": "percent", "raw": raw}
        except Exception:
            return None
    return None


def _xml_bullet_to_dict(p_pr) -> Optional[Dict[str, Any]]:
    if p_pr is None:
        return None
    bu_none = p_pr.find(qn("a:buNone"))
    if bu_none is not None:
        return {"type": "none"}
    bu_char = p_pr.find(qn("a:buChar"))
    if bu_char is not None:
        out: Dict[str, Any] = {"type": "char", "char": bu_char.get("char")}
        bu_font = p_pr.find(qn("a:buFont"))
        if bu_font is not None and bu_font.get("typeface"):
            out["font"] = {"typeface": bu_font.get("typeface")}
        return out
    bu_auto = p_pr.find(qn("a:buAutoNum"))
    if bu_auto is not None:
        out = {"type": "autoNum"}
        if bu_auto.get("type"):
            out["scheme"] = bu_auto.get("type")
        if bu_auto.get("startAt"):
            try:
                out["startAt"] = int(bu_auto.get("startAt"))
            except Exception:
                pass
        return out
    return None


def _xml_paragraph_attrs_to_dict(p_pr, canvas_emu: Optional[tuple[int, int]]) -> Optional[Dict[str, Any]]:
    if p_pr is None:
        return None
    out: Dict[str, Any] = {}
    for key in ("marL", "marR", "indent", "lvl", "algn"):
        val = p_pr.get(key)
        if val is None:
            continue
        if key in {"marL", "marR", "indent"}:
            try:
                rel = _emu_rel(int(val), canvas_emu)
                if rel is not None:
                    out[f"{key}_rel"] = rel
            except Exception:
                continue
        elif key == "lvl":
            try:
                out[key] = int(val)
            except Exception:
                continue
        else:
            out[key] = val
    return out or None


def _normalize_spacing_dict(spec: Optional[Dict[str, Any]], canvas_emu: Optional[tuple[int, int]]) -> Optional[Dict[str, Any]]:
    if not isinstance(spec, dict):
        return spec
    mode = str(spec.get("mode") or "").lower()
    if mode != "points":
        return spec
    out = {"mode": "points"}
    rel = _raw_100pt_rel(spec.get("raw"), canvas_emu)
    if rel is not None:
        out["rel"] = rel
    return out


def _xml_int_attr(node, attr_name: str) -> Optional[int]:
    if node is None:
        return None
    val = node.get(attr_name)
    if val is None:
        return None
    try:
        return int(val)
    except Exception:
        return None


def _run_spacing_info(run, canvas_emu: Optional[tuple[int, int]]) -> Dict[str, Any]:
    r_pr = run._r.find(qn("a:rPr"))
    char_spacing_raw = _xml_int_attr(r_pr, "spc")
    kern_raw = _xml_int_attr(r_pr, "kern")
    return {
        "char_spacing_rel": _raw_100pt_rel(char_spacing_raw, canvas_emu),
        "kern_raw": kern_raw,
    }


def _run_merge_key(run_info: Dict[str, Any]) -> Tuple[Any, ...]:
    color = run_info.get("color")
    if isinstance(color, dict):
        color_key = json.dumps(color, ensure_ascii=False, sort_keys=True)
    else:
        color_key = color
    return (
        run_info.get("font_name"),
        run_info.get("font_size_rel"),
        run_info.get("bold"),
        run_info.get("italic"),
        run_info.get("underline"),
        color_key,
        run_info.get("char_spacing_rel"),
        run_info.get("kern_raw"),
    )


def _merge_adjacent_runs(runs: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    merged: List[Dict[str, Any]] = []
    for run in runs:
        if not merged:
            merged.append(dict(run))
            continue
        prev = merged[-1]
        if _run_merge_key(prev) == _run_merge_key(run):
            prev["text"] = f'{prev.get("text") or ""}{run.get("text") or ""}'
            continue
        merged.append(dict(run))
    for idx, run in enumerate(merged):
        run["run_index"] = idx
    return merged


def _normalize_text_breaks(text: Optional[str]) -> str:
    if not text:
        return ""
    # PowerPoint/python-pptx may use vertical-tab for soft line breaks inside a paragraph.
    return str(text).replace("\r\n", "\n").replace("\r", "\n").replace("\v", "\n")


def _xml_el_to_spec(el: Optional[ET.Element]) -> Optional[Dict[str, Any]]:
    if el is None:
        return None
    tag = el.tag
    ns_uri = None
    local = tag
    if tag.startswith("{") and "}" in tag:
        ns_uri, local = tag[1:].split("}", 1)
    ns_name = None
    if ns_uri:
        for k, v in NS.items():
            if v == ns_uri:
                ns_name = k
                break
    spec: Dict[str, Any] = {
        "tag": local,
        "attrs": dict(el.attrib),
        "children": [],
    }
    if ns_name:
        spec["ns"] = ns_name
    if el.text:
        spec["text"] = el.text
    for child in list(el):
        child_spec = _xml_el_to_spec(child)
        if child_spec is not None:
            spec["children"].append(child_spec)
    return spec


def _fallback_text_from_shape_xml(shp, canvas_emu: Optional[tuple[int, int]] = None) -> Optional[List[Dict[str, Any]]]:
    """Best-effort text extraction from raw shape XML.

    Some shapes contain text runs that python-pptx doesn't surface via TextFrame,
    e.g. a:fld or certain drawing constructs. This fallback keeps data minimal:
    paragraphs with runs + basic rPr-derived styling.
    """
    try:
        raw = ET.tostring(getattr(shp, "_element", None), encoding="utf-8")
        if not raw:
            return None
        root = ET.fromstring(raw)
    except Exception:
        return None

    tx = root.find(".//a:txBody", NS)
    if tx is None:
        return None

    paragraphs: List[Dict[str, Any]] = []
    for p in tx.findall("a:p", NS):
        runs: List[Dict[str, Any]] = []
        text_parts: List[str] = []
        for child in list(p):
            tag = child.tag.rsplit("}", 1)[-1]
            if tag in {"r", "fld"}:
                r_pr = child.find("a:rPr", NS)
                t_nodes = child.findall("a:t", NS)
                txt = "".join((t.text or "") for t in t_nodes)
                if not txt:
                    continue
                txt = _normalize_text_breaks(txt)
                run_info: Dict[str, Any] = {
                    "run_index": len(runs),
                    "text": txt,
                    "font_name": None,
                    "font_size_rel": None,
                    "bold": None,
                    "italic": None,
                    "underline": None,
                    "color": None,
                    # Store a compact spec for rPr. combine.py will derive concrete
                    # values from run fields (font_name/font_size_rel/bold/italic/underline/color).
                    "raw_rpr": {
                        "lang": (r_pr.get("lang") if r_pr is not None and r_pr.get("lang") else "en-US"),
                        "strike": (r_pr.get("strike") if r_pr is not None and r_pr.get("strike") else "noStrike"),
                        "from_fields": ["sz", "u", "solidFill", "latin", "ea", "cs", "sym"],
                    },
                }
                if r_pr is not None:
                    # Font name from latin/rFonts.
                    latin = r_pr.find("a:latin", NS)
                    if latin is not None and latin.get("typeface"):
                        run_info["font_name"] = latin.get("typeface")
                    r_fonts = r_pr.find("a:rFonts", NS)
                    if run_info["font_name"] is None and r_fonts is not None:
                        run_info["font_name"] = r_fonts.get("ascii") or r_fonts.get("hAnsi")
                    # Size/bold/italic/underline.
                    if r_pr.get("sz"):
                        try:
                            run_info["font_size_rel"] = _pt_rel(int(r_pr.get("sz")) / 100.0, canvas_emu)
                        except Exception:
                            pass
                    if r_pr.get("b") is not None:
                        run_info["bold"] = str(r_pr.get("b")).lower() in {"1", "true"}
                    if r_pr.get("i") is not None:
                        run_info["italic"] = str(r_pr.get("i")).lower() in {"1", "true"}
                    if r_pr.get("u") is not None:
                        run_info["underline"] = r_pr.get("u")
                    # Color
                    color_spec = _solid_fill_spec_from_node(r_pr)
                    if color_spec:
                        if color_spec.get("type") == "srgb":
                            run_info["color"] = {"type": "RGB", "rgb": color_spec.get("rgb"), "mods": color_spec.get("mods")}
                        elif color_spec.get("type") == "scheme":
                            run_info["color"] = {"type": "SCHEME", "scheme": color_spec.get("scheme"), "mods": color_spec.get("mods")}
                runs.append(run_info)
                text_parts.append(txt)
            elif tag == "br":
                text_parts.append("\n")
            else:
                continue
        joined = "".join(text_parts).strip("\n")
        if not joined and not runs:
            continue
        runs = _merge_adjacent_runs(runs)
        paragraphs.append(
            {
                "runs": runs,
                "alignment": None,
                "level": None,
                "line_spacing": None,
                "space_before": None,
                "space_after": None,
            }
        )
    return paragraphs or None


def _paragraph_info(para, canvas_emu: Optional[tuple[int, int]] = None) -> Dict[str, Any]:
    p_pr = para._p.find(qn("a:pPr"))
    end_rpr = para._p.find(qn("a:endParaRPr"))
    runs_info: List[Dict[str, Any]] = []
    paragraph_text_parts: List[str] = []

    for ri, run in enumerate(para.runs):
        if not run.text:
            continue
        f = run.font
        size_pt = None
        try:
            if f.size is not None:
                size_pt = f.size.pt
        except Exception:
            size_pt = None

        run_info = {
            "run_index": ri,
            "text": _normalize_text_breaks(run.text),
            "font_name": f.name,
            "font_size_rel": _pt_rel(size_pt, canvas_emu),
            "bold": f.bold,
            "italic": f.italic,
            "underline": f.underline,
            "color": _text_color_from_run(run) or font_color_to_dict(f),
            "raw_rpr": {
                "lang": (
                    run._r.find(qn("a:rPr")).get("lang")
                    if run._r.find(qn("a:rPr")) is not None and run._r.find(qn("a:rPr")).get("lang")
                    else "en-US"
                ),
                "strike": (
                    run._r.find(qn("a:rPr")).get("strike")
                    if run._r.find(qn("a:rPr")) is not None and run._r.find(qn("a:rPr")).get("strike")
                    else "noStrike"
                ),
                "from_fields": ["sz", "u", "solidFill", "latin", "ea", "cs", "sym"],
            },
        }
        run_info.update(_run_spacing_info(run, canvas_emu))
        # #region debug-point C:text-token-export
        if "token" in (run_info.get("text") or "").lower():
            try:
                _dbg_report_method_ppt_style(
                    "C",
                    "test.py:_paragraph_info",
                    "[DEBUG] export token run",
                    {
                        "text": run_info.get("text"),
                        "font_name": run_info.get("font_name"),
                        "font_size_rel": run_info.get("font_size_rel"),
                        "bold": run_info.get("bold"),
                        "italic": run_info.get("italic"),
                        "underline": run_info.get("underline"),
                        "color": run_info.get("color"),
                        "rpr_xml": ET.tostring(run._r.find(qn("a:rPr")), encoding="unicode") if run._r.find(qn("a:rPr")) is not None else None,
                    },
                )
            except Exception:
                pass
        # #endregion
        runs_info.append(run_info)
        paragraph_text_parts.append(_normalize_text_breaks(run.text))

    runs_info = _merge_adjacent_runs(runs_info)

    line_spacing = _normalize_spacing_dict(_xml_spacing_to_dict(p_pr, "lnSpc"), canvas_emu)
    space_before = _xml_spacing_to_dict(p_pr, "spcBef")
    space_after = _xml_spacing_to_dict(p_pr, "spcAft")

    alignment = None
    try:
        alignment = str(para.alignment) if para.alignment is not None else None
    except Exception:
        alignment = None

    end_para_rpr = None
    if end_rpr is not None:
        end_para_rpr = {"color": None}
        color_spec = _solid_fill_spec_from_node(end_rpr)
        if color_spec:
            if color_spec.get("type") == "srgb":
                end_para_rpr["color"] = {"type": "RGB", "rgb": color_spec.get("rgb"), "mods": color_spec.get("mods")}
            elif color_spec.get("type") == "scheme":
                end_para_rpr["color"] = {"type": "SCHEME", "scheme": color_spec.get("scheme"), "mods": color_spec.get("mods")}

    return {
        "runs": runs_info,
        "alignment": alignment,
        "level": getattr(para, "level", None),
        "bullet": _xml_bullet_to_dict(p_pr),
        "ppr_attrs": _xml_paragraph_attrs_to_dict(p_pr, canvas_emu),
        "line_spacing": line_spacing,
        "space_before": space_before,
        "space_after": space_after,
        "end_para_rpr": end_para_rpr,
    }


def extract_text_layer(
    shp,
    slide_index: int,
    canvas_emu: Optional[tuple[int, int]],
    transform: Optional[Dict[str, float]] = None,
) -> Optional[Dict[str, Any]]:
    tf = shp.text_frame
    paragraphs: List[Dict[str, Any]] = []
    for pi, para in enumerate(tf.paragraphs):
        para_info = _paragraph_info(para, canvas_emu)
        paragraphs.append(para_info)

    # Many shapes technically "have_text_frame" but contain no editable text.
    has_runs = any(para_info.get("runs") for para_info in paragraphs)
    if not has_runs and not (tf.text or "").strip():
        # Fallback to raw XML extraction for cases python-pptx doesn't surface.
        fallback_paras = _fallback_text_from_shape_xml(shp, canvas_emu)
        if not fallback_paras:
            return None
        paragraphs = fallback_paras

    # Capture txBody/a:bodyPr. It affects line breaking and vertical alignment.
    # Keep it as a compact spec (attrs + autofit mode) so rebuild can reproduce it faithfully.
    body_pr = None
    xfrm_rot_raw = None
    try:
        root = ET.fromstring(ET.tostring(getattr(shp, "_element", None), encoding="utf-8"))
        try:
            xfrm = root.find(".//a:xfrm", NS)
            if xfrm is not None and xfrm.get("rot") is not None:
                # Preserve the original raw rot attribute (some templates use negative values).
                xfrm_rot_raw = str(xfrm.get("rot"))
        except Exception:
            xfrm_rot_raw = None
        bp = root.find(".//p:txBody/a:bodyPr", NS)
        if bp is not None:
            attrs: Dict[str, Any] = dict(bp.attrib)
            # normalize common numeric attrs to int
            for k in ("lIns", "rIns", "tIns", "bIns"):
                if k in attrs and attrs.get(k) is not None:
                    try:
                        attrs[k] = int(attrs[k])
                    except Exception:
                        pass
            autofit = None
            if bp.find("a:spAutoFit", NS) is not None:
                autofit = "spAutoFit"
            elif bp.find("a:normAutoFit", NS) is not None:
                autofit = "normAutoFit"
            elif bp.find("a:noAutoFit", NS) is not None:
                autofit = "noAutoFit"
            body_pr = {"attrs": attrs, "autofit": autofit}
    except Exception:
        body_pr = None

    # Preserve textbox shape background (solid fill / outline) when present.
    # Many templates implement "boxed text" as a TextBox shape with a solid fill.
    shape_fill = None
    shape_line = None
    try:
        shape_xml = _shape_xml_without_text(shp, transform)
        if shape_xml:
            root = ET.fromstring(shape_xml)
            sp_pr = root.find(".//p:spPr", NS)
            shape_fill = (_solid_fill_spec_from_node(sp_pr) if sp_pr is not None else None) or _style_ref_color_spec_from_root(
                root, "fillRef"
            )
            shape_line = _line_style_from_root(root)
    except Exception:
        shape_fill = None
        shape_line = None

    out = {
        "slide": slide_index,
        "shape_name": getattr(shp, "name", None),
        "kind": "text",
        "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
        **shape_flip_flags(shp),
        # Use the un-rotated bbox (xfrm off/ext) as the placement source of truth.
        # The visual bbox after rotation would "double count" rotation on rebuild.
        "box": shape_bounds_box(shp, transform, canvas_emu),
        "paragraphs": paragraphs,
    }
    if body_pr is not None:
        out["body_pr"] = body_pr
    if xfrm_rot_raw is not None:
        out["xfrm_rot_raw"] = xfrm_rot_raw
    if shape_fill is not None:
        out["shape_fill"] = shape_fill
    if shape_line is not None:
        out["shape_line"] = shape_line
    return out


def save_image_rgba_png(blob: bytes, out_path: Path) -> None:
    im = Image.open(io.BytesIO(blob))
    im = im.convert("RGBA")
    im.save(out_path, format="PNG")


def export_asset_image(
    *,
    blob: bytes,
    ext: str,
    content_type: Optional[str],
    out_base: str,
    assets_dir: Path,
) -> Dict[str, Any]:
    ext = (ext or "").lower()
    is_svg = (ext == "svg") or (content_type == "image/svg+xml")

    if is_svg:
        out_path = assets_dir / f"{out_base}.svg"
        out_path.write_bytes(blob)
        return {
            "kind": "image_svg",
            "saved_path": str(out_path),
        }

    out_path = assets_dir / f"{out_base}.png"
    try:
        save_image_rgba_png(blob, out_path)
        return {
            "kind": "image_png_rgba",
            "saved_path": str(out_path),
        }
    except Exception:
        # Some PPT images can be EMF/WMF and not readable by Pillow.
        fallback_ext = ext if ext else "bin"
        fallback = assets_dir / f"{out_base}.{fallback_ext}"
        fallback.write_bytes(blob)
        return {
            "kind": "image_raw",
            "saved_path": str(fallback),
        }


def _load_blob_as_rgba(
    *,
    blob: bytes,
    ext: str,
    content_type: Optional[str],
    box_emu: Optional[tuple[int, int, int, int]],
) -> Optional[Image.Image]:
    ext = (ext or "").lower()
    is_svg = (ext == "svg") or (content_type == "image/svg+xml")
    out_w, out_h = _raster_size_2048_for_box(box_emu)

    if is_svg:
        try:
            import cairosvg  # type: ignore
        except Exception:
            try:
                libdir = "/opt/homebrew/lib"
                prev = os.environ.get("DYLD_FALLBACK_LIBRARY_PATH") or ""
                parts = [p for p in prev.split(":") if p]
                if libdir not in parts:
                    os.environ["DYLD_FALLBACK_LIBRARY_PATH"] = ":".join([libdir] + parts)
                prev2 = os.environ.get("DYLD_LIBRARY_PATH") or ""
                parts2 = [p for p in prev2.split(":") if p]
                if libdir not in parts2:
                    os.environ["DYLD_LIBRARY_PATH"] = ":".join([libdir] + parts2)
                for k in list(sys.modules.keys()):
                    if k == "cairosvg" or k.startswith("cairosvg.") or k == "cairocffi" or k.startswith("cairocffi."):
                        del sys.modules[k]
                cairosvg = importlib.import_module("cairosvg")
            except Exception:
                return None
        try:
            png_bytes = cairosvg.svg2png(bytestring=blob, output_width=out_w, output_height=out_h)
            return Image.open(io.BytesIO(png_bytes)).convert("RGBA")
        except Exception:
            return None

    try:
        im = Image.open(io.BytesIO(blob)).convert("RGBA")
    except Exception:
        return None
    if im.size != (out_w, out_h):
        im = im.resize((out_w, out_h), resample=Image.LANCZOS)
    return im


def _load_fill_blob_as_rgba(blob: bytes, ext: str, content_type: Optional[str]) -> Optional[Image.Image]:
    """Load fill images without pre-warping them to the destination box aspect ratio."""
    ext = (ext or "").lower()
    is_svg = (ext == "svg") or (content_type == "image/svg+xml")

    if is_svg:
        try:
            import cairosvg  # type: ignore
        except Exception:
            try:
                libdir = "/opt/homebrew/lib"
                prev = os.environ.get("DYLD_FALLBACK_LIBRARY_PATH") or ""
                parts = [p for p in prev.split(":") if p]
                if libdir not in parts:
                    os.environ["DYLD_FALLBACK_LIBRARY_PATH"] = ":".join([libdir] + parts)
                prev2 = os.environ.get("DYLD_LIBRARY_PATH") or ""
                parts2 = [p for p in prev2.split(":") if p]
                if libdir not in parts2:
                    os.environ["DYLD_LIBRARY_PATH"] = ":".join([libdir] + parts2)
                for k in list(sys.modules.keys()):
                    if k == "cairosvg" or k.startswith("cairosvg.") or k == "cairocffi" or k.startswith("cairocffi."):
                        del sys.modules[k]
                cairosvg = importlib.import_module("cairosvg")
            except Exception:
                return None

        text = blob.decode("utf-8", errors="replace")
        aspect = None
        m = re.search(r'viewBox="([^"]+)"', text)
        if m:
            try:
                _, _, w, h = [float(x) for x in m.group(1).replace(",", " ").split()]
                if w > 0 and h > 0:
                    aspect = w / h
            except Exception:
                aspect = None
        if aspect is None:
            mw = re.search(r'width="([0-9.]+)', text)
            mh = re.search(r'height="([0-9.]+)', text)
            try:
                if mw and mh:
                    w = float(mw.group(1))
                    h = float(mh.group(1))
                    if w > 0 and h > 0:
                        aspect = w / h
            except Exception:
                aspect = None
        if aspect is None or aspect <= 0:
            aspect = 1.0

        if aspect >= 1.0:
            out_w = 2048
            out_h = max(1, int(round(out_w / aspect)))
        else:
            out_h = 2048
            out_w = max(1, int(round(out_h * aspect)))
        try:
            png_bytes = cairosvg.svg2png(bytestring=blob, output_width=out_w, output_height=out_h)
            return Image.open(io.BytesIO(png_bytes)).convert("RGBA")
        except Exception:
            return None

    try:
        return Image.open(io.BytesIO(blob)).convert("RGBA")
    except Exception:
        return None


def _save_final_image_layer(
    *,
    im: Image.Image,
    out_path: Path,
    final_box_emu: Optional[tuple[int, int, int, int]],
    canvas_emu: Optional[tuple[int, int]],
    rotation_deg: float,
    flip_h: bool,
    flip_v: bool,
) -> Optional[tuple[str, Optional[tuple[int, int, int, int]]]]:
    if flip_h:
        im = ImageOps.mirror(im)
    if flip_v:
        im = ImageOps.flip(im)
    if rotation_deg % 360.0:
        im = im.rotate(-rotation_deg, expand=True, resample=Image.BICUBIC)

    placed_box_emu = final_box_emu
    if final_box_emu and canvas_emu:
        im, placed_box_emu = _clip_raster_to_canvas(im, final_box_emu, canvas_emu)

    # Match the saved raster aspect ratio to the final placement box so rebuild
    # can place it 1:1 without an extra stretch step.
    target_box_emu = placed_box_emu or final_box_emu
    if target_box_emu is not None:
        target_w, target_h = _raster_size_2048_for_box(target_box_emu)
        if im.size != (target_w, target_h):
            im = im.resize((target_w, target_h), resample=Image.LANCZOS)

    try:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        im.save(out_path, format="PNG")
        return str(out_path), placed_box_emu
    except Exception:
        return None


def _fill_rect_pct(rect: Optional[Dict[str, int]]) -> Dict[str, float]:
    if not rect:
        return {"l": 0.0, "t": 0.0, "r": 0.0, "b": 0.0}
    return {k: float(rect.get(k, 0) or 0) / 100000.0 for k in ("l", "t", "r", "b")}


def _bake_fill_image_ops(
    im: Image.Image,
    src_rect: Optional[Dict[str, int]],
    fill_rect: Optional[Dict[str, int]],
    mask_oval: bool,
    mask_paths: Optional[List[Dict[str, Any]]] = None,
) -> Image.Image:
    s = _fill_rect_pct(src_rect)
    f = _fill_rect_pct(fill_rect)
    has_crop = not (all(abs(s[k]) < 1e-9 for k in s) and all(abs(f[k]) < 1e-9 for k in f))

    if has_crop:
        src_l = min(max(s["l"], 0.0), 1.0)
        src_t = min(max(s["t"], 0.0), 1.0)
        src_r = min(max(s["r"], 0.0), 1.0)
        src_b = min(max(s["b"], 0.0), 1.0)
        src_w = max(1e-9, 1.0 - src_l - src_r)
        src_h = max(1e-9, 1.0 - src_t - src_b)

        fill_l = f["l"]
        fill_t = f["t"]
        fill_r = f["r"]
        fill_b = f["b"]
        fill_w = max(1e-9, 1.0 - fill_l - fill_r)
        fill_h = max(1e-9, 1.0 - fill_t - fill_b)

        vis_u0 = max(0.0, min(1.0, (0.0 - fill_l) / fill_w))
        vis_v0 = max(0.0, min(1.0, (0.0 - fill_t) / fill_h))
        vis_u1 = max(0.0, min(1.0, (1.0 - fill_l) / fill_w))
        vis_v1 = max(0.0, min(1.0, (1.0 - fill_t) / fill_h))

        orig_u0 = src_l + vis_u0 * src_w
        orig_v0 = src_t + vis_v0 * src_h
        orig_u1 = src_l + vis_u1 * src_w
        orig_v1 = src_t + vis_v1 * src_h

        x0 = max(0, min(im.width, int(round(orig_u0 * im.width))))
        y0 = max(0, min(im.height, int(round(orig_v0 * im.height))))
        x1 = max(x0 + 1, min(im.width, int(round(orig_u1 * im.width))))
        y1 = max(y0 + 1, min(im.height, int(round(orig_v1 * im.height))))
        im = im.crop((x0, y0, x1, y1))

    if mask_oval:
        mask = Image.new("L", im.size, 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse((0, 0, im.size[0] - 1, im.size[1] - 1), fill=255)
        im.putalpha(mask)
    elif mask_paths:
        def _sample_cubic(p0, p1, p2, p3, steps: int = 48) -> List[tuple[float, float]]:
            pts: List[tuple[float, float]] = []
            for i in range(1, steps + 1):
                t = i / float(steps)
                mt = 1.0 - t
                x = (mt ** 3) * p0[0] + 3 * (mt ** 2) * t * p1[0] + 3 * mt * (t ** 2) * p2[0] + (t ** 3) * p3[0]
                y = (mt ** 3) * p0[1] + 3 * (mt ** 2) * t * p1[1] + 3 * mt * (t ** 2) * p2[1] + (t ** 3) * p3[1]
                pts.append((x, y))
            return pts

        scale = 4
        mask = Image.new("L", (max(1, im.size[0] * scale), max(1, im.size[1] * scale)), 0)
        draw = ImageDraw.Draw(mask)
        for path_spec in mask_paths:
            if not isinstance(path_spec, dict):
                continue
            commands = path_spec.get("commands") or []
            if not isinstance(commands, list) or not commands:
                continue
            poly: List[tuple[float, float]] = []
            start_pt: Optional[tuple[float, float]] = None
            cur: Optional[tuple[float, float]] = None
            for cmd in commands:
                op = str(cmd.get("op") or "").lower()
                pts = cmd.get("pts") or []
                if op == "moveto" and pts:
                    if len(poly) >= 3:
                        draw.polygon(poly, fill=255)
                    x = float(pts[0].get("x") or 0.0) * (mask.size[0] - 1)
                    y = float(pts[0].get("y") or 0.0) * (mask.size[1] - 1)
                    cur = (x, y)
                    start_pt = cur
                    poly = [cur]
                elif op == "lnto" and pts and cur is not None:
                    x = float(pts[0].get("x") or 0.0) * (mask.size[0] - 1)
                    y = float(pts[0].get("y") or 0.0) * (mask.size[1] - 1)
                    cur = (x, y)
                    poly.append(cur)
                elif op == "cubicbezto" and len(pts) == 3 and cur is not None:
                    p1 = (float(pts[0].get("x") or 0.0) * (mask.size[0] - 1), float(pts[0].get("y") or 0.0) * (mask.size[1] - 1))
                    p2 = (float(pts[1].get("x") or 0.0) * (mask.size[0] - 1), float(pts[1].get("y") or 0.0) * (mask.size[1] - 1))
                    p3 = (float(pts[2].get("x") or 0.0) * (mask.size[0] - 1), float(pts[2].get("y") or 0.0) * (mask.size[1] - 1))
                    poly.extend(_sample_cubic(cur, p1, p2, p3))
                    cur = p3
                elif op == "close":
                    if start_pt is not None and poly and poly[-1] != start_pt:
                        poly.append(start_pt)
                    if len(poly) >= 3:
                        draw.polygon(poly, fill=255)
                    poly = []
                    cur = start_pt
            if len(poly) >= 3:
                draw.polygon(poly, fill=255)
        if scale > 1:
            mask = mask.resize(im.size, resample=Image.LANCZOS)
        im = im.convert("RGBA")
        im.putalpha(mask)

    return im


def _apply_blip_alpha_mod_fix(im: Image.Image, alpha_amt: Optional[int]) -> Image.Image:
    """
    Apply DrawingML alpha modifications from a:alphaModFix/a:alpha to the rendered bitmap.

    DrawingML uses a 0..100000 scale (100000 == 100%).
    """
    if alpha_amt is None:
        return im
    try:
        amt = int(alpha_amt)
    except Exception:
        return im
    amt = max(0, min(100000, amt))
    if amt == 100000:
        return im
    if amt == 0:
        return Image.new("RGBA", im.size, (0, 0, 0, 0))
    try:
        im = im.convert("RGBA")
        r, g, b, a = im.split()
        # Multiply existing alpha by amt/100000.
        a = a.point(lambda v: int(round(v * (amt / 100000.0))))
        return Image.merge("RGBA", (r, g, b, a))
    except Exception:
        return im


def _slide_canvas_emu(slide) -> Optional[tuple[int, int]]:
    """Best-effort slide size in EMU."""
    try:
        pres = slide.part.package.presentation_part.presentation
        w = int(pres.slide_width)
        h = int(pres.slide_height)
        return w, h
    except Exception:
        return None


def _clip_box_emu_to_canvas(
    box_emu: Optional[tuple[int, int, int, int]], canvas_emu: Optional[tuple[int, int]]
) -> Optional[tuple[int, int, int, int]]:
    """Clamp an axis-aligned box in EMU to the slide canvas."""
    if not box_emu or not canvas_emu:
        return box_emu
    cw, ch = canvas_emu
    l, t, w, h = box_emu
    if w <= 0 or h <= 0:
        return None
    vis_l = max(0, l)
    vis_t = max(0, t)
    vis_r = min(cw, l + w)
    vis_b = min(ch, t + h)
    if vis_r <= vis_l or vis_b <= vis_t:
        return None
    return vis_l, vis_t, vis_r - vis_l, vis_b - vis_t


def _clamp_line_box_emu_to_canvas(
    box_emu: Optional[tuple[int, int, int, int]], canvas_emu: Optional[tuple[int, int]]
) -> Optional[tuple[int, int, int, int]]:
    """Clamp a line/connector box in EMU to the slide canvas.

    Connector boxes commonly have width==0 (vertical line) or height==0
    (horizontal line). Those are valid and must not be discarded.
    """
    if not box_emu or not canvas_emu:
        return box_emu
    cw, ch = canvas_emu
    l, t, w, h = box_emu
    if w < 0 or h < 0 or (w == 0 and h == 0):
        return None

    # Vertical line
    if w == 0:
        if l < 0 or l > cw:
            return None
        vis_t = max(0, t)
        vis_b = min(ch, t + h)
        if vis_b <= vis_t:
            return None
        return l, vis_t, 0, vis_b - vis_t

    # Horizontal line
    if h == 0:
        if t < 0 or t > ch:
            return None
        vis_l = max(0, l)
        vis_r = min(cw, l + w)
        if vis_r <= vis_l:
            return None
        return vis_l, t, vis_r - vis_l, 0

    return _clip_box_emu_to_canvas(box_emu, canvas_emu)


def _raster_size_2048_for_box(box_emu: Optional[tuple[int, int, int, int]]) -> tuple[int, int]:
    # Default square when box is missing.
    if not box_emu:
        return 2048, 2048
    _, _, w, h = box_emu
    if w <= 0 or h <= 0:
        return 2048, 2048
    if w >= h:
        out_w = 2048
        out_h = max(1, int(round(2048 * (h / float(w)))))
    else:
        out_h = 2048
        out_w = max(1, int(round(2048 * (w / float(h)))))
    return out_w, out_h


def _clip_raster_to_canvas(
    im: Image.Image, box_emu: Optional[tuple[int, int, int, int]], canvas_emu: Optional[tuple[int, int]]
) -> tuple[Image.Image, Optional[tuple[int, int, int, int]]]:
    if not box_emu or not canvas_emu:
        return im, box_emu
    l, t, w, h = box_emu
    cw, ch = canvas_emu
    # Compute intersection between [l,l+w]x[t,t+h] and slide canvas [0,cw]x[0,ch]
    vis_l = max(0, l)
    vis_t = max(0, t)
    vis_r = min(cw, l + w)
    vis_b = min(ch, t + h)
    if vis_r <= vis_l or vis_b <= vis_t:
        # Fully off-canvas; return a 1x1 transparent pixel to avoid crashes downstream.
        return Image.new("RGBA", (1, 1), (0, 0, 0, 0)), None
    # Map visible region into raster pixel space.
    x0 = int(round((vis_l - l) / float(w) * im.width))
    y0 = int(round((vis_t - t) / float(h) * im.height))
    x1 = int(round((vis_r - l) / float(w) * im.width))
    y1 = int(round((vis_b - t) / float(h) * im.height))
    x0 = max(0, min(im.width - 1, x0))
    y0 = max(0, min(im.height - 1, y0))
    x1 = max(x0 + 1, min(im.width, x1))
    y1 = max(y0 + 1, min(im.height, y1))
    cropped = im.crop((x0, y0, x1, y1))
    clipped_box = (vis_l, vis_t, vis_r - vis_l, vis_b - vis_t)
    return cropped, clipped_box


def extract_picture_asset(
    slide, shp, slide_index: int, assets_dir: Path, transform: Optional[Dict[str, float]] = None
) -> Dict[str, Any]:
    shape_id = getattr(shp, "shape_id", None) or 0
    base = f"slide{slide_index:03d}_shape{shape_id}"
    blip = None
    try:
        blip = shp._element.find(".//a:blip", NS)
    except Exception:
        blip = None
    part = None
    if blip is not None:
        part, _ = _resolve_image_parts(slide, blip)
    if part is not None:
        blob = part.blob
        ext = (getattr(part, "ext", "") or "").lower()
        ctype = getattr(part, "content_type", None)
    else:
        img = shp.image
        blob = img.blob
        ext = (img.ext or "").lower()
        ctype = getattr(img, "content_type", None)
    canvas_emu = _slide_canvas_emu(slide)
    base_box_emu = _shape_bounds_box_emu(shp, transform)
    final_box_emu = _shape_box_emu(shp, transform) or base_box_emu
    flips = shape_flip_flags(shp)
    rotation_deg = shape_rotation_deg(shp, transform)
    source_is_svg = (ext == "svg") or (ctype == "image/svg+xml")
    src_rect = None
    fill_rect = None
    alpha_amt = None

    def _rect_from_el(rect_el) -> Optional[Dict[str, int]]:
        if rect_el is None:
            return None
        out_rect: Dict[str, int] = {}
        found = False
        for key in ("l", "t", "r", "b"):
            val = rect_el.get(key)
            if val is None:
                out_rect[key] = 0
                continue
            try:
                out_rect[key] = int(val)
                found = True
            except Exception:
                out_rect[key] = 0
        return out_rect if found else None

    if blip is not None:
        try:
            parent = blip.getparent()
            blip_fill = None
            while parent is not None:
                if parent.tag == qn("p:blipFill"):
                    blip_fill = parent
                    break
                parent = parent.getparent()
            if blip_fill is not None:
                src_rect = _rect_from_el(blip_fill.find(qn("a:srcRect")))
                stretch = blip_fill.find(qn("a:stretch"))
                if stretch is not None:
                    fill_rect = _rect_from_el(stretch.find(qn("a:fillRect")))
        except Exception:
            src_rect = None
            fill_rect = None
        try:
            amf = blip.find(qn("a:alphaModFix"))
            if amf is not None and amf.get("amt") is not None:
                alpha_amt = int(amf.get("amt"))
            else:
                ael = blip.find(qn("a:alpha"))
                if ael is not None and ael.get("val") is not None:
                    alpha_amt = int(ael.get("val"))
        except Exception:
            alpha_amt = None

    has_picture_crop = (src_rect is not None) or (fill_rect is not None)
    if has_picture_crop:
        im = _load_fill_blob_as_rgba(blob=blob, ext=ext, content_type=ctype)
        if im is not None:
            im = _bake_fill_image_ops(im, src_rect, fill_rect, False)
            im = _apply_blip_alpha_mod_fix(im, alpha_amt)
    else:
        im = _load_blob_as_rgba(blob=blob, ext=ext, content_type=ctype, box_emu=base_box_emu or final_box_emu)
        if im is not None:
            im = _apply_blip_alpha_mod_fix(im, alpha_amt)
    if im is not None and final_box_emu is not None:
        rendered = _save_final_image_layer(
            im=im,
            out_path=(assets_dir / f"{base}_rendered.png").resolve(),
            final_box_emu=final_box_emu,
            canvas_emu=canvas_emu,
            rotation_deg=rotation_deg,
            flip_h=bool(flips.get("flip_h")),
            flip_v=bool(flips.get("flip_v")),
        )
        if rendered is not None:
            rendered_path, rendered_box_emu = rendered
            placed_box_emu = rendered_box_emu or final_box_emu
            return {
                "slide": slide_index,
                "shape_name": getattr(shp, "name", None),
                "kind": "svg_image_png_rgba" if source_is_svg else "image_png_rgba",
                "box": _ratio_box(*placed_box_emu, canvas_emu),
                "saved_path": rendered_path,
            }

    exported = export_asset_image(blob=blob, ext=ext, content_type=ctype, out_base=base, assets_dir=assets_dir)
    return {
        "slide": slide_index,
        "shape_name": getattr(shp, "name", None),
        "kind": exported["kind"],
        "box": shape_box(shp, transform, canvas_emu),
        "saved_path": exported["saved_path"],
    }


def extract_fill_image_assets(
    slide, shp, slide_index: int, assets_dir: Path, transform: Optional[Dict[str, float]] = None
) -> List[Dict[str, Any]]:
    """
    Extract images used as shape fills (and similar) by scanning for a:blip/@r:embed in the shape XML.
    Many templates use picture-fill shapes instead of PICTURE shapes.
    """
    el = getattr(shp, "_element", None)
    if el is None or not hasattr(el, "xpath"):
        return []

    out: List[Dict[str, Any]] = []
    # python-pptx oxml elements provide .xpath() with built-in namespace mapping
    blips = el.xpath(".//a:blip[@r:embed]")
    for idx, blip in enumerate(blips, start=1):
        part, preview_part = _resolve_image_parts(slide, blip)
        if part is None:
            continue
        ctype = getattr(part, "content_type", None)

        blip_fill = None
        try:
            parent = blip.getparent()
            while parent is not None:
                if parent.tag == qn("a:blipFill"):
                    blip_fill = parent
                    break
                parent = parent.getparent()
        except Exception:
            blip_fill = None

        def _rect_from_el(rect_el) -> Optional[Dict[str, int]]:
            if rect_el is None:
                return None
            out_rect: Dict[str, int] = {}
            found = False
            for key in ("l", "t", "r", "b"):
                val = rect_el.get(key)
                if val is None:
                    out_rect[key] = 0
                    continue
                try:
                    out_rect[key] = int(val)
                    found = True
                except Exception:
                    out_rect[key] = 0
            return out_rect if found else None

        src_rect = None
        fill_rect = None
        alpha_amt = None
        if blip_fill is not None:
            src_rect = _rect_from_el(blip_fill.find(qn("a:srcRect")))
            stretch = blip_fill.find(qn("a:stretch"))
            if stretch is not None:
                fill_rect = _rect_from_el(stretch.find(qn("a:fillRect")))
        # Some templates apply transparency to picture fills via a:alphaModFix/a:alpha on the blip.
        try:
            amf = blip.find(qn("a:alphaModFix"))
            if amf is not None and amf.get("amt") is not None:
                alpha_amt = int(amf.get("amt"))
            else:
                ael = blip.find(qn("a:alpha"))
                if ael is not None and ael.get("val") is not None:
                    alpha_amt = int(ael.get("val"))
        except Exception:
            alpha_amt = None

        shape_id = getattr(shp, "shape_id", None) or 0
        base = f"slide{slide_index:03d}_shape{shape_id}_fill{idx}"
        ext = ""
        try:
            ext = (getattr(part, "ext", "") or "").lower()
        except Exception:
            ext = ""

        mask_oval = False
        mask_paths = None
        try:
            mask_oval = (
                getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE
                and getattr(shp, "auto_shape_type", None) == MSO_AUTO_SHAPE_TYPE.OVAL
            )
        except Exception:
            mask_oval = False

        # Many templates use a custom-geometry "freeform circle" instead of a preset oval.
        # Detect the common cubic-bezier circle approximation:
        # - a:custGeom -> a:path w==h
        # - moveTo at (w/2, 0)
        # - 4 cubicBezTo segments
        def _is_custgeom_circle_like() -> bool:
            try:
                shp_el = getattr(shp, "_element", None)
                if shp_el is None:
                    return False
                root = ET.fromstring(ET.tostring(shp_el, encoding="utf-8"))
                path = root.find(".//a:custGeom/a:pathLst/a:path", NS)
                if path is None:
                    return False
                w = int(path.get("w") or "0")
                h = int(path.get("h") or "0")
                tol = max(2, max(w, h) // 500)
                if w <= 0 or h <= 0 or abs(w - h) > tol:
                    return False
                move = path.find("a:moveTo/a:pt", NS)
                if move is None:
                    return False
                mx = int(move.get("x") or "0")
                my = int(move.get("y") or "0")
                # Circle paths may start at any cardinal midpoint:
                # top-center, right-center, bottom-center, or left-center.
                cardinal_midpoints = (
                    (w // 2, 0),
                    (w, h // 2),
                    (w // 2, h),
                    (0, h // 2),
                )
                if not any(abs(mx - ex) <= tol and abs(my - ey) <= tol for ex, ey in cardinal_midpoints):
                    return False
                cubics = path.findall("a:cubicBezTo", NS)
                if len(cubics) != 4:
                    return False
                close = path.find("a:close", NS)
                if close is None:
                    return False
                return True
            except Exception:
                return False

        def _is_custgeom_rect_like() -> bool:
            try:
                shp_el = getattr(shp, "_element", None)
                if shp_el is None:
                    return False
                root = ET.fromstring(ET.tostring(shp_el, encoding="utf-8"))
                path = root.find(".//a:custGeom/a:pathLst/a:path", NS)
                if path is None:
                    return False
                w = int(path.get("w") or "0")
                h = int(path.get("h") or "0")
                if w <= 0 or h <= 0:
                    return False
                tol = max(2, max(w, h) // 500)
                cmds = list(path)
                ops = [c.tag.rsplit("}", 1)[-1] for c in cmds]
                if ops not in (["moveTo", "lnTo", "lnTo", "lnTo", "lnTo", "close"], ["moveTo", "lnTo", "lnTo", "lnTo", "close"]):
                    return False
                pts = []
                for child in cmds:
                    tag = child.tag.rsplit("}", 1)[-1]
                    if tag not in {"moveTo", "lnTo"}:
                        continue
                    pt = child.find("a:pt", NS)
                    if pt is None:
                        return False
                    pts.append((int(pt.get("x") or "0"), int(pt.get("y") or "0")))
                allowed = {(0, 0), (w, 0), (w, h), (0, h)}
                for x, y in pts:
                    if not any(abs(x - ax) <= tol and abs(y - ay) <= tol for ax, ay in allowed):
                        return False
                return True
            except Exception:
                return False

        if not mask_oval and _is_custgeom_circle_like():
            mask_oval = True

        if not mask_oval:
            try:
                shp_el = getattr(shp, "_element", None)
                if shp_el is not None:
                    root = ET.fromstring(ET.tostring(shp_el, encoding="utf-8"))
                    cust = _cust_geom_paths_spec_from_root(root)
                    if (not _is_custgeom_rect_like()) and isinstance(cust, dict) and isinstance(cust.get("paths"), list) and cust.get("paths"):
                        mask_paths = cust.get("paths")
            except Exception:
                mask_paths = None

        canvas_emu = _slide_canvas_emu(slide)
        base_box_emu = _shape_bounds_box_emu(shp, transform)
        final_box_emu = _shape_box_emu(shp, transform) or base_box_emu
        flips = shape_flip_flags(shp)
        rotation_deg = shape_rotation_deg(shp, transform)
        source_is_svg = (ext == "svg") or (ctype == "image/svg+xml")

        im = _load_fill_blob_as_rgba(blob=part.blob, ext=ext, content_type=ctype)
        if im is None and preview_part is not None:
            preview_ext = (getattr(preview_part, "ext", "") or "").lower()
            preview_ctype = getattr(preview_part, "content_type", None)
            im = _load_fill_blob_as_rgba(blob=preview_part.blob, ext=preview_ext, content_type=preview_ctype)

        if im is not None and final_box_emu is not None:
            im = _bake_fill_image_ops(im, src_rect, fill_rect, mask_oval, mask_paths)
            im = _apply_blip_alpha_mod_fix(im, alpha_amt)
            rendered = _save_final_image_layer(
                im=im,
                out_path=(assets_dir / f"{base}_rendered.png").resolve(),
                final_box_emu=final_box_emu,
                canvas_emu=canvas_emu,
                rotation_deg=rotation_deg,
                flip_h=bool(flips.get("flip_h")),
                flip_v=bool(flips.get("flip_v")),
            )
            if rendered is not None:
                rendered_path, rendered_box_emu = rendered
                placed_box_emu = rendered_box_emu or final_box_emu
                out.append(
                    {
                        "slide": slide_index,
                        "shape_name": getattr(shp, "name", None),
                        "kind": "svg_image_png_rgba" if source_is_svg else "image_png_rgba",
                        "box": _ratio_box(*placed_box_emu, canvas_emu),
                        "saved_path": rendered_path,
                    }
                )
                continue

        exported = export_asset_image(blob=part.blob, ext=ext, content_type=ctype, out_base=base, assets_dir=assets_dir)
        out.append(
            {
                "slide": slide_index,
                "shape_name": getattr(shp, "name", None),
                "kind": exported["kind"],
                "box": shape_box(shp, transform, canvas_emu),
                "saved_path": exported["saved_path"],
            }
        )
    return out


def _shape_has_blip_fill(shp) -> bool:
    try:
        return getattr(shp, "_element", None).find(".//a:blip", NS) is not None
    except Exception:
        return False


def _remove_direct_child(root: ET.Element, tag: str) -> None:
    for child in list(root):
        if child.tag == tag:
            root.remove(child)


def _shape_xml_without_text(shp, transform: Optional[Dict[str, float]] = None) -> Optional[str]:
    try:
        root = ET.fromstring(ET.tostring(shp._element, encoding="utf-8"))
    except Exception:
        return None

    # Text is exported separately; remove it from the geometry XML to avoid duplication on rebuild.
    _remove_direct_child(root, f"{{{NS['p']}}}txBody")

    xfrm = root.find(".//a:xfrm", NS)
    if xfrm is not None:
        off = xfrm.find("a:off", NS)
        ext = xfrm.find("a:ext", NS)
        bounds_box = _shape_bounds_box_emu(shp, transform)
        if bounds_box is not None:
            left, top, width, height = bounds_box
        else:
            left = top = width = height = None
        if off is not None:
            if left is not None:
                off.set("x", str(int(left)))
            if top is not None:
                off.set("y", str(int(top)))
        if ext is not None:
            if width is not None:
                ext.set("cx", str(int(width)))
            if height is not None:
                ext.set("cy", str(int(height)))
        try:
            rot = int(round(shape_rotation_deg(shp, transform) * 60000))
            if rot:
                xfrm.set("rot", str(rot))
            elif "rot" in xfrm.attrib:
                del xfrm.attrib["rot"]
        except Exception:
            pass

    return ET.tostring(root, encoding="unicode")


def _color_spec_from_color_el(color_el: ET.Element) -> Optional[Dict[str, Any]]:
    tag = color_el.tag.rsplit("}", 1)[-1]
    if tag == "srgbClr" and color_el.get("val"):
        spec: Dict[str, Any] = {"type": "srgb", "rgb": str(color_el.get("val")).upper()}
    elif tag == "schemeClr" and color_el.get("val"):
        spec = {"type": "scheme", "scheme": str(color_el.get("val"))}
    else:
        return None
    mods: List[Dict[str, Any]] = []
    for child in list(color_el):
        tag = child.tag.rsplit("}", 1)[-1]
        if tag in {"tint", "shade", "lumMod", "lumOff", "alpha"} and child.get("val"):
            try:
                mods.append({"op": tag, "val": int(child.get("val"))})
            except Exception:
                mods.append({"op": tag, "val": child.get("val")})
    if mods:
        spec["mods"] = mods
    return spec


def _solid_fill_spec_from_node(node: ET.Element) -> Optional[Dict[str, Any]]:
    solid = node.find("a:solidFill", NS)
    if solid is None and node.tag.rsplit("}", 1)[-1] == "solidFill":
        solid = node
    if solid is None:
        return None
    srgb = solid.find("a:srgbClr", NS)
    if srgb is not None and srgb.get("val"):
        return _color_spec_from_color_el(srgb)
    scheme = solid.find("a:schemeClr", NS)
    if scheme is not None and scheme.get("val"):
        return _color_spec_from_color_el(scheme)
    return None


def _background_fill_from_color_spec(color_spec: Optional[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    if not isinstance(color_spec, dict):
        return None
    out: Dict[str, Any] = {"type": "solid"}
    color_type = str(color_spec.get("type") or "").lower()
    if color_type in {"rgb", "srgb"} and color_spec.get("rgb"):
        out.update({"color_type": "RGB", "rgb": str(color_spec.get("rgb")).upper()})
    elif color_type == "scheme" and color_spec.get("scheme"):
        out.update({"color_type": "SCHEME", "scheme": color_spec.get("scheme")})
    if color_spec.get("mods"):
        out["mods"] = list(color_spec.get("mods") or [])
    return out


def _gradient_fill_spec_from_node(node: ET.Element) -> Optional[Dict[str, Any]]:
    grad = node.find("a:gradFill", NS)
    if grad is None and node.tag.rsplit("}", 1)[-1] == "gradFill":
        grad = node
    if grad is None:
        return None

    out: Dict[str, Any] = {"type": "gradient"}
    stops: List[Dict[str, Any]] = []
    gs_lst = grad.find("a:gsLst", NS)
    if gs_lst is not None:
        for gs in gs_lst.findall("a:gs", NS):
            stop: Dict[str, Any] = {}
            if gs.get("pos") is not None:
                try:
                    stop["pos"] = int(gs.get("pos"))
                except Exception:
                    stop["pos"] = gs.get("pos")

            color_spec: Optional[Dict[str, Any]] = None
            for child in list(gs):
                color_spec = _color_spec_from_color_el(child)
                if color_spec is not None:
                    break

            if color_spec is not None:
                color_type = str(color_spec.get("type") or "").lower()
                if color_type in {"rgb", "srgb"} and color_spec.get("rgb"):
                    stop.update({"color_type": "RGB", "rgb": str(color_spec.get("rgb")).upper()})
                elif color_type == "scheme" and color_spec.get("scheme"):
                    stop.update({"color_type": "SCHEME", "scheme": color_spec.get("scheme")})
                if color_spec.get("mods"):
                    stop["mods"] = list(color_spec.get("mods") or [])
            if stop:
                stops.append(stop)

    if stops:
        out["stops"] = stops
        first_stop = next((stop for stop in stops if isinstance(stop, dict)), None)
        if isinstance(first_stop, dict):
            if first_stop.get("color_type") == "RGB" and first_stop.get("rgb"):
                out.update({"color_type": "RGB", "rgb": first_stop.get("rgb")})
            elif first_stop.get("color_type") == "SCHEME" and first_stop.get("scheme"):
                out.update({"color_type": "SCHEME", "scheme": first_stop.get("scheme")})
            if first_stop.get("mods"):
                out["mods"] = list(first_stop.get("mods") or [])

    lin = grad.find("a:lin", NS)
    if lin is not None and lin.get("ang") is not None:
        try:
            out["angle"] = int(lin.get("ang"))
        except Exception:
            out["angle"] = lin.get("ang")

    path = grad.find("a:path", NS)
    if path is not None and path.get("path"):
        out["path"] = path.get("path")

    return out if len(out) > 1 else None


def _gradient_stop_rgba(stop: Dict[str, Any]) -> Optional[tuple[int, int, int, float]]:
    if not isinstance(stop, dict):
        return None
    if str(stop.get("color_type") or "").upper() != "RGB":
        return None
    rgb = str(stop.get("rgb") or "").strip().lstrip("#")
    if len(rgb) != 6:
        return None
    try:
        r = int(rgb[0:2], 16)
        g = int(rgb[2:4], 16)
        b = int(rgb[4:6], 16)
    except Exception:
        return None

    alpha = 1.0
    for mod in stop.get("mods") or []:
        if not isinstance(mod, dict):
            continue
        if str(mod.get("op") or "").lower() != "alpha":
            continue
        try:
            alpha = max(0.0, min(1.0, float(mod.get("val")) / 100000.0))
        except Exception:
            alpha = 1.0
        break
    return r, g, b, alpha


def _gradient_fill_to_svg(
    gradient_spec: Dict[str, Any], out_w: int, out_h: int
) -> Optional[str]:
    if out_w <= 0 or out_h <= 0:
        return None

    stops_svg: List[str] = []
    for stop in gradient_spec.get("stops") or []:
        rgba = _gradient_stop_rgba(stop)
        if rgba is None:
            continue
        r, g, b, alpha = rgba
        try:
            pos = max(0.0, min(1.0, float(stop.get("pos", 0)) / 100000.0))
        except Exception:
            pos = 0.0
        stops_svg.append(
            f'<stop offset="{pos * 100:.4f}%" stop-color="rgb({r},{g},{b})" stop-opacity="{alpha:.6f}"/>'
        )
    if not stops_svg:
        return None

    gradient_markup = ""
    if str(gradient_spec.get("path") or "").lower() == "circle":
        gradient_markup = (
            '<radialGradient id="bggrad" cx="50%" cy="50%" r="70.710678%" '
            'fx="50%" fy="50%">'
            + "".join(stops_svg)
            + "</radialGradient>"
        )
    else:
        try:
            angle_deg = float(gradient_spec.get("angle") or 0.0) / 60000.0
        except Exception:
            angle_deg = 0.0
        theta = math.radians(angle_deg)
        dx = math.cos(theta)
        dy = -math.sin(theta)
        x1 = 50.0 - (dx * 50.0)
        y1 = 50.0 - (dy * 50.0)
        x2 = 50.0 + (dx * 50.0)
        y2 = 50.0 + (dy * 50.0)
        gradient_markup = (
            f'<linearGradient id="bggrad" x1="{x1:.4f}%" y1="{y1:.4f}%" '
            f'x2="{x2:.4f}%" y2="{y2:.4f}%">'
            + "".join(stops_svg)
            + "</linearGradient>"
        )

    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{out_w}" height="{out_h}" '
        f'viewBox="0 0 {out_w} {out_h}">'
        "<defs>"
        f"{gradient_markup}"
        "</defs>"
        f'<rect x="0" y="0" width="{out_w}" height="{out_h}" fill="url(#bggrad)"/>'
        "</svg>"
    )


def _svg_bytes_to_png_bytes(svg_bytes: bytes, out_w: int, out_h: int) -> Optional[bytes]:
    try:
        import cairosvg  # type: ignore
    except Exception:
        try:
            libdir = "/opt/homebrew/lib"
            prev = os.environ.get("DYLD_FALLBACK_LIBRARY_PATH") or ""
            parts = [p for p in prev.split(":") if p]
            if libdir not in parts:
                os.environ["DYLD_FALLBACK_LIBRARY_PATH"] = ":".join([libdir] + parts)
            prev2 = os.environ.get("DYLD_LIBRARY_PATH") or ""
            parts2 = [p for p in prev2.split(":") if p]
            if libdir not in parts2:
                os.environ["DYLD_LIBRARY_PATH"] = ":".join([libdir] + parts2)
            for k in list(sys.modules.keys()):
                if k == "cairosvg" or k.startswith("cairosvg.") or k == "cairocffi" or k.startswith("cairocffi."):
                    del sys.modules[k]
            cairosvg = importlib.import_module("cairosvg")
        except Exception:
            return None
    try:
        return cairosvg.svg2png(bytestring=svg_bytes, output_width=out_w, output_height=out_h)
    except Exception:
        return None


def _export_gradient_background_layer(
    gradient_spec: Dict[str, Any],
    slide_index: int,
    slide_width_emu: int,
    slide_height_emu: int,
    assets_dir: Path,
    *,
    suffix: str = "gradient",
) -> Optional[Dict[str, Any]]:
    full_box_emu = (0, 0, int(slide_width_emu), int(slide_height_emu))
    out_w, out_h = _raster_size_2048_for_box(full_box_emu)
    svg_text = _gradient_fill_to_svg(gradient_spec, out_w, out_h)
    if not svg_text:
        return None
    png_bytes = _svg_bytes_to_png_bytes(svg_text.encode("utf-8"), out_w, out_h)
    if not png_bytes:
        return None

    out_path = assets_dir / f"slide{slide_index:03d}_background_{suffix}.png"
    try:
        Image.open(io.BytesIO(png_bytes)).convert("RGBA").save(out_path, format="PNG")
    except Exception:
        return None

    return {
        "slide": slide_index,
        "shape_name": "slide_background_gradient",
        "kind": "svg_image_png_rgba",
        "box": {"left": 0.0, "top": 0.0, "width": 1.0, "height": 1.0},
        "saved_path": str(out_path),
    }


def _style_ref_color_spec_from_root(root: ET.Element, ref_name: str) -> Optional[Dict[str, Any]]:
    ref = root.find(f".//p:style/a:{ref_name}", NS)
    if ref is None:
        return None
    for child in list(ref):
        spec = _color_spec_from_color_el(child)
        if spec:
            return spec
    return None


def _line_style_from_root(root: ET.Element) -> Optional[Dict[str, Any]]:
    ln = root.find(".//a:ln", NS)
    if ln is None:
        return None
    width_emu = None
    try:
        width_emu = int(ln.get("w") or 0)
    except Exception:
        width_emu = None
    color = _solid_fill_spec_from_node(ln) or _style_ref_color_spec_from_root(root, "lnRef")
    out: Dict[str, Any] = {}
    if width_emu is not None and width_emu > 0:
        out["width_pt"] = round(float(width_emu) / EMU_PER_PT, 4)
    if color is not None:
        out["color"] = color
    dash = ln.find("a:prstDash", NS)
    if dash is not None and dash.get("val"):
        out["dash"] = str(dash.get("val"))
    head = ln.find("a:headEnd", NS)
    if head is not None and head.get("type"):
        out["head_end"] = {
            "type": str(head.get("type")),
            "w": head.get("w"),
            "len": head.get("len"),
        }
    tail = ln.find("a:tailEnd", NS)
    if tail is not None and tail.get("type"):
        out["tail_end"] = {
            "type": str(tail.get("type")),
            "w": tail.get("w"),
            "len": tail.get("len"),
        }
    return out or None


def _has_visible_fill(fill: Optional[Dict[str, Any]]) -> bool:
    return isinstance(fill, dict) and bool(fill)


def _has_visible_line(line: Optional[Dict[str, Any]]) -> bool:
    if not isinstance(line, dict):
        return False
    try:
        width_pt = float(line.get("width_pt") or 0.0)
    except Exception:
        width_pt = 0.0
    return width_pt > 0 and line.get("color") is not None


def _is_fully_transparent_geo(fill: Optional[Dict[str, Any]], line: Optional[Dict[str, Any]]) -> bool:
    return (not _has_visible_fill(fill)) and (not _has_visible_line(line))


def _cust_geom_paths_spec_from_root(root: ET.Element) -> Optional[Dict[str, Any]]:
    cust = root.find(".//a:custGeom", NS)
    if cust is None:
        return None
    path_lst = cust.find("a:pathLst", NS)
    if path_lst is None:
        return None

    paths_out: List[Dict[str, Any]] = []
    for path in path_lst.findall("a:path", NS):
        try:
            w = int(path.get("w") or 0)
            h = int(path.get("h") or 0)
        except Exception:
            w, h = 0, 0
        if w <= 0 or h <= 0:
            return None

        commands: List[Dict[str, Any]] = []
        def _round_coord(v: float) -> float:
            r = round(float(v), 5)
            return 0.0 if r == -0.0 else r
        for child in list(path):
            tag = child.tag.rsplit("}", 1)[-1]
            if tag == "moveTo":
                pt = child.find("a:pt", NS)
                if pt is None:
                    continue
                x = int(pt.get("x") or 0)
                y = int(pt.get("y") or 0)
                commands.append({"op": "moveTo", "pts": [{"x": _round_coord(x / w), "y": _round_coord(y / h)}]})
            elif tag == "lnTo":
                pt = child.find("a:pt", NS)
                if pt is None:
                    continue
                x = int(pt.get("x") or 0)
                y = int(pt.get("y") or 0)
                commands.append({"op": "lnTo", "pts": [{"x": _round_coord(x / w), "y": _round_coord(y / h)}]})
            elif tag == "cubicBezTo":
                pts = child.findall("a:pt", NS)
                if len(pts) != 3:
                    continue
                out_pts = []
                for pt in pts:
                    x = int(pt.get("x") or 0)
                    y = int(pt.get("y") or 0)
                    out_pts.append({"x": _round_coord(x / w), "y": _round_coord(y / h)})
                commands.append({"op": "cubicBezTo", "pts": out_pts})
            elif tag == "close":
                commands.append({"op": "close", "pts": []})
            else:
                # Unknown ops are rare in current templates; keep it strict.
                return None

        paths_out.append({"w": w, "h": h, "commands": commands})

    return {"type": "custGeom", "paths": paths_out}


def _extract_sppr_custgeom_spec_from_shape_xml(shape_xml: str, canvas_emu: tuple[int, int]) -> Optional[Dict[str, Any]]:
    """Export a compact spPr spec with fully relative geometry.

    - xfrm.off/ext are omitted (set to None); rebuild uses `box` for placement.
    - path w/h are stored relative to canvas width/height.
    - path points are stored relative to the path's own w/h (0..1).
    """
    try:
        root = ET.fromstring(shape_xml)
    except Exception:
        return None

    cw, ch = canvas_emu
    if cw <= 0 or ch <= 0:
        return None

    def _round_coord(v: float) -> float:
        r = round(float(v), 5)
        return 0.0 if r == -0.0 else r

    xfrm = root.find(".//a:xfrm", NS)
    flip_h = False
    flip_v = False
    try:
        if xfrm is not None:
            flip_h = str(xfrm.get("flipH") or "").lower() in {"1", "true"}
            flip_v = str(xfrm.get("flipV") or "").lower() in {"1", "true"}
    except Exception:
        flip_h = False
        flip_v = False

    cust = root.find(".//a:custGeom", NS)
    if cust is None:
        return None
    path_lst = cust.find("a:pathLst", NS)
    if path_lst is None:
        return None

    path_items: List[Dict[str, Any]] = []
    for path in path_lst.findall("a:path", NS):
        try:
            w = int(path.get("w") or 0)
            h = int(path.get("h") or 0)
        except Exception:
            w, h = 0, 0
        if w <= 0 or h <= 0:
            continue

        commands: List[Dict[str, Any]] = []
        for child in list(path):
            tag = child.tag.rsplit("}", 1)[-1]
            if tag in {"moveTo", "lnTo"}:
                pt = child.find("a:pt", NS)
                if pt is None:
                    continue
                x = int(pt.get("x") or 0)
                y = int(pt.get("y") or 0)
                commands.append({"type": tag, "pt": {"x": _round_coord(x / w), "y": _round_coord(y / h)}})
            elif tag == "cubicBezTo":
                pts = child.findall("a:pt", NS)
                if len(pts) != 3:
                    continue
                out_pts = []
                for pt in pts:
                    x = int(pt.get("x") or 0)
                    y = int(pt.get("y") or 0)
                    out_pts.append({"x": _round_coord(x / w), "y": _round_coord(y / h)})
                commands.append({"type": "cubicBezTo", "pts": out_pts})
            elif tag == "close":
                commands.append({"type": "close"})
            else:
                continue

        # w/h relative to canvas
        path_items.append({"w": _round_coord(w / cw), "h": _round_coord(h / ch), "commands": commands})

    if not path_items:
        return None

    return {
        "spPr": {
            "xfrm": {
                "flipH": bool(flip_h),
                "flipV": bool(flip_v),
                "off": None,
                "ext": None,
            },
            "custGeom": {
                "avLst": [],
                "gdLst": [],
                "ahLst": [],
                "cxnLst": [],
                "rect": {"l": "l", "t": "t", "r": "r", "b": "b"},
                "pathLst": path_items,
            },
        }
    }


def _prst_geom_spec_from_root(root: ET.Element) -> Optional[Dict[str, Any]]:
    prst = root.find(".//a:prstGeom", NS)
    if prst is None:
        return None
    name = prst.get("prst")
    if not name:
        return None
    return {"type": "prstGeom", "prst": str(name)}

def _graph_shape_spec_from_shape_xml(shape_xml: str, kind: str, box_emu: tuple[int, int, int, int], canvas_emu: tuple[int, int]) -> Optional[Dict[str, Any]]:
    try:
        root = ET.fromstring(shape_xml)
    except Exception:
        return None

    if kind == "ppt_graph_geo":
        geom = _cust_geom_paths_spec_from_root(root) or _prst_geom_spec_from_root(root)
        if geom is None:
            return None
        sp_pr = root.find(".//p:spPr", NS)
        fill = (_solid_fill_spec_from_node(sp_pr) if sp_pr is not None else None) or _style_ref_color_spec_from_root(root, "fillRef")
        if isinstance(geom, dict) and geom.get("type") == "prstGeom" and geom.get("prst") in {
            "leftBrace",
            "rightBrace",
            "leftBracket",
            "rightBracket",
            "bracePair",
            "bracketPair",
        }:
            fill = None
        return {
            "format": "shape_spec_v1",
            "type": "geo",
            "geom": geom,
            "fill": fill,
            "line": _line_style_from_root(root),
        }

    return None


def _is_outer_page_border_graph_layer(layer: Optional[Dict[str, Any]]) -> bool:
    """
    Filter decorative page-edge frames at export time only.
    We target a nearly full-page no-fill rectangle with a single normalized rect path.
    """
    if not isinstance(layer, dict):
        return False
    if layer.get("kind") != "ppt_graph_geo":
        return False
    if layer.get("fill") is not None:
        return False
    if not isinstance(layer.get("line"), dict):
        return False

    box = layer.get("box") or {}
    try:
        left = float(box.get("left") or 0.0)
        top = float(box.get("top") or 0.0)
        width = float(box.get("width") or 0.0)
        height = float(box.get("height") or 0.0)
    except Exception:
        return False
    if not (width >= 0.90 and height >= 0.90 and left <= 0.06 and top <= 0.06):
        return False

    shape_xml = layer.get("shape_xml") or {}
    sppr = (shape_xml.get("spPr") or {}) if isinstance(shape_xml, dict) else {}
    cust = sppr.get("custGeom") if isinstance(sppr, dict) else None
    if not isinstance(cust, dict):
        return False
    path_lst = cust.get("pathLst") or []
    if not (isinstance(path_lst, list) and len(path_lst) == 1 and isinstance(path_lst[0], dict)):
        return False
    cmds = path_lst[0].get("commands") or []
    if not (isinstance(cmds, list) and len(cmds) >= 5):
        return False
    try:
        pts = []
        for cmd in cmds[:4]:
            if cmd.get("type") not in {"moveTo", "lnTo"}:
                return False
            pt = cmd.get("pt") or {}
            pts.append((float(pt.get("x") or 0.0), float(pt.get("y") or 0.0)))
        if str(cmds[4].get("type") or "").lower() != "close":
            return False
        if pts != [(0.0, 0.0), (1.0, 0.0), (1.0, 1.0), (0.0, 1.0)]:
            return False
    except Exception:
        return False
    return True


def extract_graph_layer(
    slide, shp, slide_index: int, transform: Optional[Dict[str, float]] = None
) -> Optional[Dict[str, Any]]:
    st = getattr(shp, "shape_type", None)
    # TextBoxes are usually exported as editable text layers, but some templates use empty
    # TextBox shapes as colored rectangles/panels (no text). Preserve those as geo layers.
    allow_textbox_geo = False
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        has_text = False
        try:
            if getattr(shp, "has_text_frame", False) and getattr(shp, "text_frame", None) is not None:
                tf = shp.text_frame
                if (tf.text or "").strip():
                    has_text = True
                else:
                    for p in tf.paragraphs:
                        for r in getattr(p, "runs", []):
                            if (getattr(r, "text", "") or "").strip():
                                has_text = True
                                break
                        if has_text:
                            break
        except Exception:
            has_text = False
        allow_textbox_geo = not has_text
        if not allow_textbox_geo:
            return None
    # Connector/arrow lines are not consistently exposed as a named enum across
    # python-pptx versions. Detect them from the underlying OOXML tag as well.
    el = getattr(shp, "_element", None)
    is_connector = False
    try:
        is_connector = el is not None and str(getattr(el, "tag", "")).endswith("}cxnSp")
    except Exception:
        is_connector = False
    if st not in {MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.TEXT_BOX} and not is_connector:
        return None
    if _shape_has_blip_fill(shp):
        # Picture-fill shapes are already exported via fill-image extraction.
        return None

    shape_xml = _shape_xml_without_text(shp, transform)
    if not shape_xml:
        return None

    # Some templates use preset geometry `prstGeom prst="line"` which python-pptx
    # exposes as AUTO_SHAPE instead of LINE, so detect it from OOXML too.
    is_prst_line = 'prst="line"' in shape_xml
    kind = "ppt_graph_line" if (st == MSO_SHAPE_TYPE.LINE or is_connector or is_prst_line) else "ppt_graph_geo"
    canvas_emu = _slide_canvas_emu(slide)

    # Export the unrotated xfrm-like box. Rotation is exported separately and applied on rebuild.
    # Using a post-rotation AABB here would double-apply rotation in combine.py.
    box_emu = _shape_bounds_box_emu(shp, transform) or _shape_box_emu(shp, transform)
    if box_emu is None:
        return None

    # Normalize ET-generated prefixes (ns0/ns1) to stable p/a for readability and diffing.
    shape_xml_norm = (
        shape_xml.replace('xmlns:ns0="http://schemas.openxmlformats.org/presentationml/2006/main"', 'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"')
        .replace('xmlns:ns1="http://schemas.openxmlformats.org/drawingml/2006/main"', 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"')
        .replace("<ns0:", "<p:")
        .replace("</ns0:", "</p:")
        .replace("<ns1:", "<a:")
        .replace("</ns1:", "</a:")
    )

    # Lines/connectors: export as endpoints + line style (incl. arrowheads).
    if kind == "ppt_graph_line":
        cw, ch = canvas_emu
        if cw <= 0 or ch <= 0:
            return None
        p1 = None
        p2 = None
        try:
            p1 = {"x": float(shp.begin_x) / cw, "y": float(shp.begin_y) / ch}
            p2 = {"x": float(shp.end_x) / cw, "y": float(shp.end_y) / ch}
        except Exception:
            l, t, w, h = (float(v) for v in box_emu)
            if w == 0 and h == 0:
                p1 = {"x": l / cw, "y": t / ch}
                p2 = {"x": l / cw, "y": t / ch}
            else:
                if w == 0:
                    lp1 = (l, t)
                    lp2 = (l, t + h)
                elif h == 0:
                    lp1 = (l, t)
                    lp2 = (l + w, t)
                else:
                    lp1 = (l, t)
                    lp2 = (l + w, t + h)

                cx = l + (w / 2.0)
                cy = t + (h / 2.0)
                theta = math.radians(float(shape_rotation_deg(shp, transform) or 0.0))

                def _rot_pt(px: float, py: float) -> tuple[float, float]:
                    dx = px - cx
                    dy = py - cy
                    rx = cx + (dx * math.cos(theta) - dy * math.sin(theta))
                    ry = cy + (dx * math.sin(theta) + dy * math.cos(theta))
                    return rx, ry

                x1, y1 = _rot_pt(*lp1)
                x2, y2 = _rot_pt(*lp2)
                p1 = {"x": x1 / cw, "y": y1 / ch}
                p2 = {"x": x2 / cw, "y": y2 / ch}
        try:
            root = ET.fromstring(shape_xml_norm)
        except Exception:
            return None
        shape_spec = {
            "format": "shape_spec_v1",
            "type": "line",
            "p1": p1,
            "p2": p2,
            "line": _line_style_from_root(root),
        }
        return {
            "slide": slide_index,
            "shape_name": getattr(shp, "name", None),
            "kind": kind,
            "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
            **shape_flip_flags(shp),
            "box": _ratio_box(*box_emu, canvas_emu),
            "shape_xml": shape_spec,
        }

    shape_spec = _extract_sppr_custgeom_spec_from_shape_xml(shape_xml_norm, canvas_emu)
    # Fallback: preset geometry (prstGeom) and non-custGeom shapes.
    if shape_spec is None:
        fallback = _graph_shape_spec_from_shape_xml(shape_xml_norm, kind, box_emu, canvas_emu)
        if fallback is None:
            return None
        if kind == "ppt_graph_geo" and _is_fully_transparent_geo(fallback.get("fill"), fallback.get("line")):
            return None
        out = {
            "slide": slide_index,
            "shape_name": getattr(shp, "name", None),
            "kind": kind,
            "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
            **shape_flip_flags(shp),
            "box": _ratio_box(*box_emu, canvas_emu),
            "shape_xml": fallback,
        }
        if _is_outer_page_border_graph_layer(out):
            return None
        return out

    # Extract fill/line styles. Many templates use theme colors (schemeClr) via style refs;
    # keep a compact color spec that combine.py can re-apply deterministically.
    fill = None
    line = None
    try:
        root = ET.fromstring(shape_xml_norm)
        sp_pr = root.find(".//p:spPr", NS)
        fill = (_solid_fill_spec_from_node(sp_pr) if sp_pr is not None else None) or _style_ref_color_spec_from_root(
            root, "fillRef"
        )
        line = _line_style_from_root(root)
    except Exception:
        fill = None
        line = None

    if kind == "ppt_graph_geo" and _is_fully_transparent_geo(fill, line):
        return None

    out = {
        "slide": slide_index,
        "shape_name": getattr(shp, "name", None),
        "kind": kind,
        "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
        **shape_flip_flags(shp),
        "box": _ratio_box(*box_emu, canvas_emu),
        "shape_xml": shape_spec,
        "fill": fill,
        "line": line,
    }
    if _is_outer_page_border_graph_layer(out):
        return None
    return out


def _line_style_from_ln_el(ln_el) -> Optional[Dict[str, Any]]:
    if ln_el is None:
        return None
    try:
        root = ET.fromstring(ET.tostring(ln_el, encoding="utf-8"))
    except Exception:
        return None
    width_emu = None
    try:
        width_emu = int(root.get("w") or 0)
    except Exception:
        width_emu = None
    color = _solid_fill_spec_from_node(root)
    out: Dict[str, Any] = {}
    if width_emu is not None and width_emu > 0:
        out["width_pt"] = round(float(width_emu) / EMU_PER_PT, 4)
    if color is not None:
        out["color"] = color
    dash = root.find("a:prstDash", NS)
    if dash is not None and dash.get("val"):
        out["dash"] = str(dash.get("val"))
    return out or None


def extract_table_layer(
    slide,
    shp,
    slide_index: int,
    canvas_emu: Optional[tuple[int, int]],
    transform: Optional[Dict[str, float]] = None,
) -> Optional[Dict[str, Any]]:
    st = getattr(shp, "shape_type", None)
    if st != MSO_SHAPE_TYPE.TABLE and not hasattr(shp, "table"):
        return None
    try:
        tbl = shp.table
    except Exception:
        return None

    table_style_id = None
    try:
        root = ET.fromstring(ET.tostring(getattr(shp, "_element", None), encoding="utf-8"))
        tbl_el = root.find(".//a:tbl", NS)
        if tbl_el is not None:
            tbl_pr = tbl_el.find("a:tblPr", NS)
            if tbl_pr is not None:
                sid = tbl_pr.find("a:tableStyleId", NS)
                if sid is not None and sid.text and sid.text.strip():
                    table_style_id = sid.text.strip()
    except Exception:
        table_style_id = None

    try:
        rows = len(tbl.rows)
        cols = len(tbl.columns)
    except Exception:
        return None
    if rows <= 0 or cols <= 0:
        return None

    col_widths_pt: List[float] = []
    for col in tbl.columns:
        try:
            col_widths_pt.append(float(col.width.pt))
        except Exception:
            col_widths_pt.append(0.0)

    row_heights_pt: List[float] = []
    for row in tbl.rows:
        try:
            row_heights_pt.append(float(row.height.pt))
        except Exception:
            row_heights_pt.append(0.0)

    cells: List[List[Dict[str, Any]]] = []
    for ri in range(rows):
        row_cells: List[Dict[str, Any]] = []
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            paragraphs: List[Dict[str, Any]] = []
            margin_left_pt = None
            margin_top_pt = None
            margin_right_pt = None
            margin_bottom_pt = None
            vertical_anchor = None
            body_pr = None
            try:
                tf = cell.text_frame
                for pi, para in enumerate(tf.paragraphs):
                    paragraphs.append(_paragraph_info(para, canvas_emu))
                try:
                    margin_left_pt = float(tf.margin_left.pt)
                    margin_top_pt = float(tf.margin_top.pt)
                    margin_right_pt = float(tf.margin_right.pt)
                    margin_bottom_pt = float(tf.margin_bottom.pt)
                except Exception:
                    pass
                try:
                    va = getattr(tf, "vertical_anchor", None)
                    vertical_anchor = str(va) if va is not None else None
                except Exception:
                    vertical_anchor = None
                try:
                    tc = getattr(cell, "_tc", None)
                    tx_body = tc.find(qn("a:txBody")) if tc is not None and hasattr(tc, "find") else None
                    bp = tx_body.find(qn("a:bodyPr")) if tx_body is not None and hasattr(tx_body, "find") else None
                    if bp is not None:
                        root = ET.fromstring(ET.tostring(bp, encoding="utf-8"))
                        attrs: Dict[str, Any] = dict(root.attrib)
                        autofit = None
                        if root.find("a:spAutoFit", NS) is not None:
                            autofit = "spAutoFit"
                        elif root.find("a:normAutoFit", NS) is not None:
                            autofit = "normAutoFit"
                        elif root.find("a:noAutoFit", NS) is not None:
                            autofit = "noAutoFit"
                        body_pr = {"attrs": attrs, "autofit": autofit}
                except Exception:
                    body_pr = None
            except Exception:
                paragraphs = []

            fill = None
            borders: Dict[str, Any] = {}
            try:
                tc = getattr(cell, "_tc", None)
                tcpr = tc.tcPr if tc is not None and hasattr(tc, "tcPr") else None
                if tcpr is not None:
                    tcpr_xml = ET.fromstring(ET.tostring(tcpr, encoding="utf-8"))
                    fill = _solid_fill_spec_from_node(tcpr_xml)
                    for tag, key in (("lnL", "l"), ("lnR", "r"), ("lnT", "t"), ("lnB", "b")):
                        ln = tcpr.find(qn(f"a:{tag}")) if hasattr(tcpr, "find") else None
                        spec = _line_style_from_ln_el(ln)
                        if spec is not None:
                            borders[key] = spec
            except Exception:
                fill = None
                borders = {}

            row_cells.append(
                {
                    "row": ri,
                    "col": ci,
                    "text": _normalize_text_breaks(getattr(cell, "text", "")),
                    "paragraphs": paragraphs,
                    "margin_left_pt": margin_left_pt,
                    "margin_top_pt": margin_top_pt,
                    "margin_right_pt": margin_right_pt,
                    "margin_bottom_pt": margin_bottom_pt,
                    "vertical_anchor": vertical_anchor,
                    "body_pr": body_pr,
                    "fill": fill,
                    "borders": borders or None,
                }
            )
        cells.append(row_cells)

    return {
        "slide": slide_index,
        "shape_name": getattr(shp, "name", None),
        "kind": "ppt_graph_table",
        "rotation_deg": _round_rotation_deg(shape_rotation_deg(shp, transform)),
        **shape_flip_flags(shp),
        "box": shape_box(shp, transform, canvas_emu),
        "rows": rows,
        "cols": cols,
        "table_style_id": table_style_id,
        "col_widths_pt": col_widths_pt,
        "row_heights_pt": row_heights_pt,
        "cells": cells,
    }


def extract_slide_background_layers(
    slide, slide_index: int, assets_dir: Path, slide_width_emu: int, slide_height_emu: int
) -> List[Dict[str, Any]]:
    el = getattr(slide, "_element", None)
    if el is None or not hasattr(el, "xpath"):
        return []
    out: List[Dict[str, Any]] = []
    bg_pr = None
    try:
        bg_pr = el.find(".//p:bg/p:bgPr", NS)
    except Exception:
        bg_pr = None
    blips = el.xpath(".//p:bg//a:blip[@r:embed]")
    for idx, blip in enumerate(blips, start=1):
        r_id = blip.get(qn("r:embed"))
        if not r_id:
            continue
        try:
            rel = slide.part.rels[r_id]
        except KeyError:
            continue
        if getattr(rel, "is_external", False):
            continue
        part = getattr(rel, "target_part", None)
        if part is None or not hasattr(part, "blob"):
            continue
        ctype = getattr(part, "content_type", None)
        if not (ctype and str(ctype).startswith("image/")):
            continue

        base = f"slide{slide_index:03d}_background{idx}"
        ext = ""
        try:
            ext = (getattr(part, "ext", "") or "").lower()
        except Exception:
            ext = ""
        exported = export_asset_image(blob=part.blob, ext=ext, content_type=ctype, out_base=base, assets_dir=assets_dir)
        out.append(
            {
                "slide": slide_index,
                "shape_name": "slide_background",
                "kind": exported["kind"],
                "box": None,
                "saved_path": exported["saved_path"],
                "source": "slide_background",
            }
        )
    gradient_spec = _gradient_fill_spec_from_node(bg_pr) if bg_pr is not None else None
    if gradient_spec is not None:
        gradient_layer = _export_gradient_background_layer(
            gradient_spec,
            slide_index,
            slide_width_emu,
            slide_height_emu,
            assets_dir,
        )
        if gradient_layer is not None:
            out.append(gradient_layer)
    return out


def extract_slide_background_fill(slide) -> Optional[Dict[str, Any]]:
    el = getattr(slide, "_element", None)
    if el is None:
        return None
    try:
        bg = el.find(".//p:bg", NS)
        bg_pr = el.find(".//p:bg/p:bgPr", NS)
        bg_ref = el.find(".//p:bg/p:bgRef", NS)
    except Exception:
        bg = None
        bg_pr = None
        bg_ref = None
    if bg_pr is None and bg_ref is None:
        return None

    result: Dict[str, Any] = {}

    if bg_pr is not None:
        solid_spec = _background_fill_from_color_spec(_solid_fill_spec_from_node(bg_pr))
        if solid_spec is not None:
            return solid_spec

        gradient_spec = _gradient_fill_spec_from_node(bg_pr)
        if gradient_spec is not None:
            return None

        return result or None

    # Background may be a theme reference rather than a concrete bgPr.
    if bg_ref is not None:
        result.update({"type": "ref"})
    return result or None


def extract_slide_canvas_layer(
    slide_index: int, slide_width_emu: int, slide_height_emu: int, background_fill: Optional[Dict[str, Any]] = None
) -> Dict[str, Any]:
    """
    Add a synthetic bottom-most layer representing the full slide canvas.
    This preserves the original page size even if no visible layer reaches the slide edges.
    """
    return {
        "slide": slide_index,
        "shape_name": "slide_canvas",
        "kind": "slide_canvas",
        "box": {
            "left": 0.0,
            "top": 0.0,
            "width": 1.0,
            "height": 1.0,
        },
        "canvas_width_emu": int(slide_width_emu),
        "canvas_height_emu": int(slide_height_emu),
        "background_fill": dict(background_fill) if background_fill else None,
    }


def extract_embedded_fonts_from_pptx(pptx_path: Path, fonts_dir: Path) -> List[Dict[str, Any]]:
    """
    Extract embedded font binaries from a .pptx (OOXML zip package).

    Notes:
    - PowerPoint commonly stores embedded fonts under 'ppt/fonts/' as .fntdata/.odttf.
    - These are saved as-is (no de-obfuscation / conversion).
    """
    ensure_dir(fonts_dir)
    prefixes = ("ppt/fonts/", "ppt/embeddings/", "word/fonts/", "xl/fonts/")

    # Regenerate fonts output each run to avoid mixing old "font7.fntdata" style outputs with
    # new "Typeface__style.*" outputs.
    try:
        for p in fonts_dir.iterdir():
            if p.is_file():
                p.unlink()
    except Exception:
        pass

    def _safe_filename(name: str, fallback: str) -> str:
        # Keep filenames reasonably safe across platforms.
        # We allow Unicode (font names can be non-ASCII) but remove path separators and other invalid chars.
        s = (name or "").strip()
        if not s:
            s = fallback
        bad = '<>:"/\\|?*\0'
        for ch in bad:
            s = s.replace(ch, "_")
        s = "".join(c if (c >= " " and c != "\x7f") else "_" for c in s)
        s = " ".join(s.split())
        # Avoid super long filenames.
        if len(s) > 120:
            s = s[:120].rstrip()
        return s or fallback

    def _write_unique(base_name: str, data: bytes) -> Path:
        out_path = fonts_dir / base_name
        if out_path.exists():
            stem, ext = out_path.stem, out_path.suffix
            i = 2
            while (fonts_dir / f"{stem}_{i}{ext}").exists():
                i += 1
            out_path = fonts_dir / f"{stem}_{i}{ext}"
        out_path.write_bytes(data)
        return out_path

    extracted: List[Dict[str, Any]] = []
    extracted_zip_paths: set[str] = set()

    # Prefer the authoritative embedded-font list in presentation.xml (typeface + style + rId).
    # This avoids missing fonts stored as .fntdata/.odttf with unexpected naming/extensions.
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    r_id_attr = f"{{{NS_R}}}id"

    with zipfile.ZipFile(pptx_path, "r") as zf:
        try:
            pres_xml = zf.read("ppt/presentation.xml")
            pres_rels_xml = zf.read("ppt/_rels/presentation.xml.rels")
            pres_root = ET.fromstring(pres_xml)
            rels_root = ET.fromstring(pres_rels_xml)

            rid_to_rel: Dict[str, Dict[str, str]] = {}
            for rel in list(rels_root):
                rid = rel.get("Id")
                if not rid:
                    continue
                rid_to_rel[rid] = {"type": rel.get("Type") or "", "target": rel.get("Target") or ""}

            embedded_font_lst = pres_root.find(f".//{{{NS_P}}}embeddedFontLst")
            if embedded_font_lst is not None:
                for ef in embedded_font_lst.findall(f"./{{{NS_P}}}embeddedFont"):
                    font_el = ef.find(f"./{{{NS_P}}}font")
                    typeface = font_el.get("typeface") if font_el is not None else None

                    for style_tag in ("regular", "bold", "italic", "boldItalic"):
                        style_el = ef.find(f"./{{{NS_P}}}{style_tag}")
                        if style_el is None:
                            continue
                        rid = style_el.get(r_id_attr)
                        if not rid:
                            continue
                        rel = rid_to_rel.get(rid) or {}
                        target = rel.get("target") or ""
                        if not target:
                            continue
                        # Relationships are relative to `ppt/presentation.xml`, so targets are under `ppt/`.
                        zip_path = f"ppt/{target.lstrip('/')}"
                        try:
                            data = zf.read(zip_path)
                        except KeyError:
                            continue

                        # Name by font face instead of "font7.fntdata" etc.
                        suffix = Path(zip_path).suffix or ".bin"
                        face = _safe_filename(typeface or "", "unknown_font")
                        style = _safe_filename(style_tag, "regular")
                        base_name = f"{face}__{style}{suffix}"
                        out_path = _write_unique(base_name, data)
                        extracted_zip_paths.add(zip_path)
                        extracted.append(
                            {
                                "zip_path": zip_path,
                                "saved_path": str(out_path),
                                "size_bytes": len(data),
                                "typeface": typeface,
                                "style": style_tag,
                                "rid": rid,
                                "rel_type": rel.get("type"),
                                "rel_target": target,
                                "source": "presentation_embeddedFontLst",
                            }
                        )
        except Exception:
            # Keep going with best-effort extraction by folder scan.
            pass

        # Fallback: export any files that live under the typical font directories.
        # This catches unreferenced font parts (rare) and non-standard extensions.
        for zi in zf.infolist():
            zip_name = zi.filename or ""
            zip_name_l = zip_name.lower()
            if zip_name_l.endswith("/"):
                continue
            if not any(zip_name_l.startswith(p) for p in prefixes):
                continue
            if zip_name in extracted_zip_paths:
                continue

            data = zf.read(zip_name)
            out_path = _write_unique(Path(zip_name).name, data)
            extracted_zip_paths.add(zip_name)
            extracted.append(
                {
                    "zip_path": zip_name,
                    "saved_path": str(out_path),
                    "size_bytes": len(data),
                    "source": "package_scan",
                }
            )

    return extracted


def iter_export_layers_top_to_bottom(
    slide,
    shapes,
    slide_index: int,
    assets_dir: Path,
    transform: Optional[Dict[str, float]] = None,
) -> Iterable[Dict[str, Any]]:
    """
    Iterate shapes from top -> bottom (front -> back).
    Group shapes are flattened into slide-level layers.
    """
    transform = transform or identity_transform()
    # python-pptx exposes shapes typically in back -> front order; reverse to get top -> bottom.
    for shp in reversed(list(shapes)):
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            next_transform = child_transform_for_group(shp, transform)
            for item in iter_export_layers_top_to_bottom(
                slide, shp.shapes, slide_index, assets_dir, next_transform
            ):
                yield item
            continue

        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield extract_picture_asset(slide, shp, slide_index, assets_dir, transform)
            continue

        table_layer = extract_table_layer(slide, shp, slide_index, _slide_canvas_emu(slide), transform)
        if table_layer is not None:
            yield table_layer
            continue

        text_layer = None
        if getattr(shp, "has_text_frame", False):
            text_layer = extract_text_layer(shp, slide_index, _slide_canvas_emu(slide), transform)

        graph_layer = extract_graph_layer(slide, shp, slide_index, transform)
        if graph_layer is not None and text_layer is not None and graph_layer.get("kind") == "ppt_graph_geo":
            shape_xml = graph_layer.get("shape_xml") or {}
            geom = shape_xml.get("geom") if isinstance(shape_xml, dict) else None
            if isinstance(geom, dict) and geom.get("type") == "prstGeom":
                graph_layer["text_content"] = {
                    "text": text_layer.get("text"),
                    "paragraphs": text_layer.get("paragraphs"),
                    "text_box": text_layer.get("text_box"),
                }
                yield graph_layer
                graph_layer = None
                text_layer = None

        if text_layer is not None:
            yield text_layer
        if graph_layer is not None:
            yield graph_layer

        # Also extract images used in shape fills (common in templates).
        # Text is drawn over fill, so yield fill after text to keep top->bottom semantics.
        for img_layer in extract_fill_image_assets(slide, shp, slide_index, assets_dir, transform):
            yield img_layer

        # Skip other layer types by default.
        continue


def main() -> None:
    ap = argparse.ArgumentParser(description="Extract PPTX layers: images(SVG/PNG RGBA) + editable text styles to JSON.")
    # Positional args are required by default. Use nargs='?' to allow running with defaults.
    ap.add_argument(
        "pptx",
        nargs="?",
        default="/Users/bytedance/Documents/test/code/PPTX格式/海报/~$Beige Vintage Newspaper Travel to Italy Poster.pptx",
        help="Input .pptx file path",
    )
    ap.add_argument("-o", "--out", help="Output folder", default="/Users/bytedance/Documents/test/code/PPTX格式/res")
    ap.add_argument("--json-name", default="layers.json", help="Output JSON file name (default: layers.json)")
    args = ap.parse_args()
    pptx_path = Path(args.pptx).expanduser().resolve()

    # Office often creates a temp/lock file like "~$Foo.pptx". It's not a valid pptx package to read.
    # If user passes such a path, try to locate the real file by stripping the "~$" prefix.
    if pptx_path.name.startswith("~$"):
        candidate = pptx_path.with_name(pptx_path.name[2:])
        if candidate.exists():
            pptx_path = candidate.resolve()
    fixs = "/".join(pptx_path.parts[-2:]) if len(pptx_path.parts) >= 2 else pptx_path.name
    out_dir = Path(args.out).expanduser().resolve() / fixs
    if not pptx_path.exists():
        # "~$*.pptx" is usually an Office temp/lock file; remind user to pick the real pptx.
        if pptx_path.name.startswith("~$"):
            raise SystemExit(
                f"Input not found (looks like Office temp file '~$'). "
                f"Please pass the real .pptx (usually same name without '~$'): {pptx_path}"
            )
        raise SystemExit(f"Input not found: {pptx_path}")

    ensure_dir(out_dir)
    assets_dir = out_dir / "assets"
    ensure_dir(assets_dir)
    fonts_dir = out_dir / "fonts"

    prs = Presentation(str(pptx_path))
    layers: List[Dict[str, Any]] = []

    for si, slide in enumerate(prs.slides, start=1):
        per_slide_index = 0
        # Shapes: top -> bottom
        for layer in iter_export_layers_top_to_bottom(slide, slide.shapes, si, assets_dir):
            per_slide_index += 1
            layer["shape_name"] = per_slide_index
            layers.append(layer)

        # Background image is the bottom-most layer, append at the end (bottom).
        for layer in extract_slide_background_layers(slide, si, assets_dir, int(prs.slide_width), int(prs.slide_height)):
            per_slide_index += 1
            layer["shape_name"] = per_slide_index
            layers.append(layer)

        # Synthetic slide canvas is the absolute bottom-most layer.
        background_fill = extract_slide_background_fill(slide)
        canvas_layer = extract_slide_canvas_layer(si, int(prs.slide_width), int(prs.slide_height), background_fill)
        per_slide_index += 1
        canvas_layer["shape_name"] = per_slide_index
        layers.append(canvas_layer)

    out_json = out_dir / args.json_name
    # Protocol: group by slide (outer list), and omit per-layer `slide` field.
    # Example:
    #   [[{layer},{layer}], [{layer},...], ...]
    by_slide: Dict[int, List[Dict[str, Any]]] = {}
    for layer in layers:
        try:
            slide_idx = int(layer.get("slide") or 1)
        except Exception:
            slide_idx = 1
        normalized = dict(layer)
        normalized.pop("slide", None)
        by_slide.setdefault(slide_idx, []).append(normalized)
    grouped_layers: List[List[Dict[str, Any]]] = [by_slide[k] for k in sorted(by_slide.keys())]

    out_json.write_text(
        json.dumps(_encode_scaled_numbers_for_json(grouped_layers, 1000), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    fonts = extract_embedded_fonts_from_pptx(pptx_path, fonts_dir)
    if fonts:
        (out_dir / "fonts.json").write_text(json.dumps(fonts, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"Saved {len(layers)} layers ({len(grouped_layers)} slide group(s)) to: {out_json}")
    print(f"Assets saved to: {assets_dir}")
    if fonts:
        print(f"Fonts extracted: {len(fonts)} -> {fonts_dir}")
    else:
        print("Fonts extracted: 0 (no embedded fonts found in package)")


if __name__ == "__main__":
    main()
